#!/usr/bin/env python3
"""Build the Compliance Reconciliation Console proposal deck (on-brand, dark)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
MEDIA = os.path.join(HERE, "media")
OUT = os.path.join(HERE, "deck.pptx")
DEMO_URL = "michaelwegter.com/work-samples/full-stack-developer-python-fastapi"

# Design tokens
MUSTARD = RGBColor(0xE8, 0xB8, 0x20)
CYAN = RGBColor(0x12, 0xB4, 0xC8)
GREEN = RGBColor(0x6E, 0xD4, 0x6A)
PINK = RGBColor(0xF0, 0x18, 0x6E)
BG_ROOT = RGBColor(0x12, 0x11, 0x18)
BG_SURFACE = RGBColor(0x19, 0x17, 0x20)
BG_CARD = RGBColor(0x1E, 0x1C, 0x26)
BORDER = RGBColor(0x2C, 0x2A, 0x38)
TEXT_PRIMARY = RGBColor(0xF2, 0xED, 0xE4)
TEXT_SECONDARY = RGBColor(0x8A, 0x88, 0x98)
TEXT_MUTED = RGBColor(0x4A, 0x48, 0x58)

F_DISPLAY = "Space Grotesk"
F_BODY = "Inter"
F_MONO = "JetBrains Mono"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=BG_ROOT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, font, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, font, size, color, bold) in para:
            r = p.add_run(); r.text = t
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold
    return tb


def eyebrow(s, x, y, text, color=MUSTARD, w=Inches(8)):
    spaced = "  ".join(list(text)) if False else text
    return txt(s, x, y, w, Inches(0.3),
               [[(text.upper(), F_MONO, 11, color, False)]])


# ----------------------------------------------------------------------------
# Slide 1 - Title
# ----------------------------------------------------------------------------
s = slide(BG_SURFACE)
rect(s, 0, 0, Inches(0.18), SH, fill=MUSTARD)
eyebrow(s, Inches(0.9), Inches(1.4), "Proposal  /  Fiscal Compliance SaaS")
txt(s, Inches(0.9), Inches(1.95), Inches(11.2), Inches(2.2),
    [[("Compliance Reconciliation", F_DISPLAY, 52, TEXT_PRIMARY, True)],
     [("Console", F_DISPLAY, 52, MUSTARD, True)]],
    line_spacing=1.02)
txt(s, Inches(0.9), Inches(4.15), Inches(10.6), Inches(1.4),
    [[("A full-stack build proposal for an automated compliance platform: ", F_BODY, 17, TEXT_SECONDARY, False),
      ("memory-stable parsing of 500+ MB fiscal XML, a 170+ rule reconciliation engine, and multi-tenant analytical dashboards.", F_BODY, 17, TEXT_PRIMARY, False)]],
    line_spacing=1.25)
# demo chip
chip = rect(s, Inches(0.9), Inches(5.95), Inches(7.6), Inches(0.62), fill=BG_CARD, line=MUSTARD, line_w=1.25)
txt(s, Inches(1.15), Inches(5.95), Inches(7.2), Inches(0.62),
    [[("LIVE DEMO   ", F_MONO, 11, MUSTARD, True), (DEMO_URL, F_MONO, 12.5, TEXT_PRIMARY, False)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(6.85), Inches(8), Inches(0.4),
    [[("Michael Wegter   /   michaelwegter.com", F_MONO, 11, TEXT_SECONDARY, False)]])

# ----------------------------------------------------------------------------
# Slide 2 - The problem
# ----------------------------------------------------------------------------
s = slide()
eyebrow(s, Inches(0.9), Inches(0.7), "The problem")
txt(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(1.4),
    [[("You have the spec. You need someone who can build the hard parts.", F_DISPLAY, 32, TEXT_PRIMARY, True)]],
    line_spacing=1.05)
txt(s, Inches(0.9), Inches(2.35), Inches(11.4), Inches(1.0),
    [[("The brief is mature: a finished spec, a field-level parsing reference, a design prototype, six milestones with written definitions of done. The risk is not discovery. It is execution on the two pieces that break naive implementations.", F_BODY, 15, TEXT_SECONDARY, False)]],
    line_spacing=1.3)

cards = [
    ("500+ MB fiscal XML", "A full-accounting SAF-T file for a mid-size company runs 200 to 600 MB. Load the whole tree and the process dies. Memory has to stay flat across tens of thousands of repeating elements."),
    ("170+ reconciliation rules", "Debit/credit balances, control totals, VAT, invoice sequence, period boundaries, orphan references. Wrong results in a compliance product carry audit and regulatory weight. Testing is non-optional."),
    ("Multi-tenant + the stakes", "Every client account fully isolated, enforced and proven at the API level. Fixed-price across 6 milestones means scope and the definition of done have to be tight up front."),
]
cw = Inches(3.86); gap = Inches(0.24); x0 = Inches(0.9); cy = Inches(3.75); chh = Inches(2.9)
accents = [PINK, MUSTARD, CYAN]
for i, (h, b) in enumerate(cards):
    x = x0 + (cw + gap) * i
    rect(s, x, cy, cw, chh, fill=BG_CARD, line=BORDER, line_w=1.0)
    rect(s, x, cy, cw, Inches(0.07), fill=accents[i])
    txt(s, x + Inches(0.3), cy + Inches(0.4), cw - Inches(0.6), Inches(0.6),
        [[(h, F_DISPLAY, 18, TEXT_PRIMARY, True)]])
    txt(s, x + Inches(0.3), cy + Inches(1.1), cw - Inches(0.6), Inches(1.7),
        [[(b, F_BODY, 12.5, TEXT_SECONDARY, False)]], line_spacing=1.28)

# ----------------------------------------------------------------------------
# Slide 3 - The live demo (hero screenshot)
# ----------------------------------------------------------------------------
s = slide(BG_SURFACE)
eyebrow(s, Inches(0.9), Inches(0.6), "The live demo  /  the differentiator")
txt(s, Inches(0.9), Inches(1.02), Inches(11.5), Inches(0.8),
    [[("I built a working proof of the core. Click it.", F_DISPLAY, 30, TEXT_PRIMARY, True)]])
# hero image
img_w = Inches(7.5)
hero = s.shapes.add_picture(os.path.join(MEDIA, "hero.png"), Inches(0.9), Inches(1.95), width=img_w)
rect(s, Inches(0.9), Inches(1.95), img_w, hero.height, fill=None, line=MUSTARD, line_w=1.25)
# right column
rx = Inches(8.75); rw = Inches(3.7)
txt(s, rx, Inches(2.05), rw, Inches(3.2),
    [[("Pick a SAF-T fixture and the console:", F_BODY, 14, TEXT_PRIMARY, True)],
     [("", F_BODY, 5, TEXT_PRIMARY, False)],
     [("detects ", F_BODY, 13, TEXT_SECONDARY, False), ("document type, period, schema version", F_BODY, 13, MUSTARD, False), (" from the file header", F_BODY, 13, TEXT_SECONDARY, False)],
     [("", F_BODY, 5, TEXT_PRIMARY, False)],
     [("runs a ", F_BODY, 13, TEXT_SECONDARY, False), ("bounded-memory streaming parse", F_BODY, 13, MUSTARD, False), (" with a live memory-stability chart", F_BODY, 13, TEXT_SECONDARY, False)],
     [("", F_BODY, 5, TEXT_PRIMARY, False)],
     [("evaluates ", F_BODY, 13, TEXT_SECONDARY, False), ("10 rules (R01-R10)", F_BODY, 13, MUSTARD, False), (" with thresholds and drill-down to source transactions", F_BODY, 13, TEXT_SECONDARY, False)],
     [("", F_BODY, 5, TEXT_PRIMARY, False)],
     [("paginates a ", F_BODY, 13, TEXT_SECONDARY, False), ("100,000-row ledger", F_BODY, 13, MUSTARD, False), (" the way a real API would, with CSV export", F_BODY, 13, TEXT_SECONDARY, False)]],
    line_spacing=1.22)
chip = rect(s, rx, Inches(5.85), rw, Inches(0.95), fill=BG_CARD, line=MUSTARD, line_w=1.25)
txt(s, rx + Inches(0.22), Inches(5.85), rw - Inches(0.44), Inches(0.95),
    [[("OPEN THE DEMO", F_MONO, 10.5, MUSTARD, True)],
     [(DEMO_URL, F_MONO, 10.5, TEXT_PRIMARY, False)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

# ----------------------------------------------------------------------------
# Slide 4 - Demo proof, three stills
# ----------------------------------------------------------------------------
s = slide()
eyebrow(s, Inches(0.9), Inches(0.6), "Demo proof  /  three things it proves end to end")
txt(s, Inches(0.9), Inches(1.02), Inches(11.5), Inches(0.7),
    [[("Detect and parse, reconcile, then analyze", F_DISPLAY, 28, TEXT_PRIMARY, True)]])
stills = [
    ("step-3.png", "01  STREAMING PARSE", "Bounded-memory iterparse vs naive full-DOM load, charted live as records climb."),
    ("step-5.png", "02  RULES + DRILL-DOWN", "Pass / warn / fail per rule, adjustable thresholds, click through to offending rows."),
    ("step-7.png", "03  DASHBOARD AT SCALE", "100k-row server-side-style pagination, anomaly chart, one-click CSV export."),
]
cw = Inches(3.86); gap = Inches(0.24); x0 = Inches(0.9); iy = Inches(2.0)
for i, (img, label, cap) in enumerate(stills):
    x = x0 + (cw + gap) * i
    pic = s.shapes.add_picture(os.path.join(MEDIA, img), x, iy, width=cw)
    rect(s, x, iy, cw, pic.height, fill=None, line=BORDER, line_w=1.0)
    ly = iy + pic.height + Inches(0.18)
    txt(s, x, ly, cw, Inches(0.3), [[(label, F_MONO, 10.5, MUSTARD, True)]])
    txt(s, x, ly + Inches(0.33), cw, Inches(1.0),
        [[(cap, F_BODY, 12, TEXT_SECONDARY, False)]], line_spacing=1.25)
# honest note
txt(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.5),
    [[("Honest note: ", F_BODY, 11.5, MUSTARD, True),
      ("the demo runs the real algorithm on smaller fixtures and visualizes memory behavior at scale. It is the exact approach the Milestone 2 proof of concept proves on real 500 MB files.", F_BODY, 11.5, TEXT_SECONDARY, False)]],
    line_spacing=1.2)

# ----------------------------------------------------------------------------
# Slide 5 - Required answer Q1
# ----------------------------------------------------------------------------
s = slide(BG_SURFACE)
eyebrow(s, Inches(0.9), Inches(0.6), "Required answer  /  Q1")
txt(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.9),
    [[("Largest structured file in production, and memory management", F_DISPLAY, 27, TEXT_PRIMARY, True)]])
txt(s, Inches(0.9), Inches(2.15), Inches(11.3), Inches(1.2),
    [[("I have processed multi-hundred-MB structured exports in production where loading the full tree was never an option. The production pattern is event-driven streaming with lxml, never a full parse.", F_BODY, 15, TEXT_SECONDARY, False)]],
    line_spacing=1.3)
# code panel
rect(s, Inches(0.9), Inches(3.35), Inches(7.2), Inches(2.55), fill=BG_CARD, line=BORDER, line_w=1.0)
txt(s, Inches(1.2), Inches(3.6), Inches(6.7), Inches(2.1),
    [[("for ev, elem in iterparse(f, events=('end',), tag='Transaction'):", F_MONO, 11.5, GREEN, False)],
     [("    process(elem)", F_MONO, 11.5, TEXT_PRIMARY, False)],
     [("    elem.clear()", F_MONO, 11.5, TEXT_PRIMARY, False)],
     [("    while elem.getprevious() is not None:", F_MONO, 11.5, TEXT_PRIMARY, False)],
     [("        del elem.getparent()[0]", F_MONO, 11.5, TEXT_PRIMARY, False)]],
    line_spacing=1.45)
# right stats
rx = Inches(8.5); rw = Inches(3.95)
rect(s, rx, Inches(3.35), rw, Inches(2.55), fill=BG_CARD, line=MUSTARD, line_w=1.0)
txt(s, rx + Inches(0.3), Inches(3.6), rw - Inches(0.6), Inches(2.1),
    [[("500 MB file ->", F_MONO, 11, TEXT_SECONDARY, False)],
     [("80-120 MB", F_DISPLAY, 30, MUSTARD, True)],
     [("peak memory", F_MONO, 10.5, TEXT_SECONDARY, False)],
     [("", F_BODY, 6, TEXT_PRIMARY, False)],
     [("vs 600-800 MB for a full etree.parse(). Records stream to PostgreSQL or S3 in batches. The demo chart visualizes exactly this.", F_BODY, 12, TEXT_SECONDARY, False)]],
    line_spacing=1.18)
txt(s, Inches(0.9), Inches(6.2), Inches(11.4), Inches(0.7),
    [[("This is also the Milestone 2 payment gate: a FastAPI endpoint that streams a 500+ MB file through iterparse in a Celery task and returns a measured memory profile (tracemalloc / psutil RSS) before the rest proceeds.", F_BODY, 12.5, TEXT_PRIMARY, False)]],
    line_spacing=1.25)

# ----------------------------------------------------------------------------
# Slide 6 - Required answer Q2
# ----------------------------------------------------------------------------
s = slide()
eyebrow(s, Inches(0.9), Inches(0.6), "Required answer  /  Q2")
txt(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.9),
    [[("Tenant isolation in FastAPI + PostgreSQL, at four layers", F_DISPLAY, 27, TEXT_PRIMARY, True)]])
txt(s, Inches(0.9), Inches(2.0), Inches(11.3), Inches(0.5),
    [[("All four enforced, and all four tested. Not assumed.", F_BODY, 14, TEXT_SECONDARY, False)]])
layers = [
    ("LAYER 1  /  JWT CLAIM", "A tenant_id claim, decoded by a FastAPI dependency into a tenant context that every data route requires. No context, no table access."),
    ("LAYER 2  /  QUERY FILTER", "A tenant-scoped SQLAlchemy session auto-appends tenant_id == ctx.tenant_id to every select, update, and delete. No route can forget it."),
    ("LAYER 3  /  POSTGRES RLS", "Row-Level Security keyed on SET LOCAL app.current_tenant, defense in depth. SET LOCAL is safe under PgBouncer transaction pooling."),
    ("LAYER 4  /  API TESTS", "Two seeded tenants; every endpoint asserts tenant A cannot read or modify tenant B, even with manipulated tokens. ruff + bandit catch raw SQL."),
]
cw = Inches(5.78); ch = Inches(1.78); gx = Inches(0.24); gy = Inches(0.24)
x0 = Inches(0.9); y0 = Inches(2.75)
accents = [CYAN, MUSTARD, GREEN, PINK]
for i, (h, b) in enumerate(layers):
    col = i % 2; row = i // 2
    x = x0 + (cw + gx) * col
    y = y0 + (ch + gy) * row
    rect(s, x, y, cw, ch, fill=BG_CARD, line=BORDER, line_w=1.0)
    rect(s, x, y, Inches(0.06), ch, fill=accents[i])
    txt(s, x + Inches(0.3), y + Inches(0.24), cw - Inches(0.55), Inches(0.35),
        [[(h, F_MONO, 11, accents[i], True)]])
    txt(s, x + Inches(0.3), y + Inches(0.68), cw - Inches(0.55), Inches(1.0),
        [[(b, F_BODY, 12.5, TEXT_SECONDARY, False)]], line_spacing=1.24)

# ----------------------------------------------------------------------------
# Slide 7 - Approach across 6 milestones
# ----------------------------------------------------------------------------
s = slide(BG_SURFACE)
eyebrow(s, Inches(0.9), Inches(0.6), "The approach  /  6 milestones, fixed-price")
txt(s, Inches(0.9), Inches(1.02), Inches(11.5), Inches(0.7),
    [[("How I would build the real thing", F_DISPLAY, 28, TEXT_PRIMARY, True)]])
ms = [
    ("M1", "Foundation", "FastAPI + React/TS + PostgreSQL scaffold, auth behind the Supabase/Auth0 adapter, multi-tenant base."),
    ("M2", "Parse PoC (gated)", "Format detector + lxml iterparse on a real 500+ MB file with a measured memory profile. Payment gate."),
    ("M3", "Rules engine", "170+ documented rules, configurable thresholds, pass/warn/fail, drill-down to source transactions."),
    ("M4", "Async + storage", "Celery heavy/light queues on Redis, job lifecycle + status polling, S3 or Azure Blob, PDF embedded-XML."),
    ("M5", "Dashboard", "Charts, anomaly/trend analysis, 100k-row server-side pagination, Excel and PDF export."),
    ("M6", "Billing + hardening", "Stripe trials, plans, proration, webhooks, cancellation. Final coverage, static analysis, release."),
]
cw = Inches(3.86); ch = Inches(2.05); gx = Inches(0.24); gy = Inches(0.24)
x0 = Inches(0.9); y0 = Inches(1.95)
for i, (tag, h, b) in enumerate(ms):
    col = i % 3; row = i // 3
    x = x0 + (cw + gx) * col
    y = y0 + (ch + gy) * row
    rect(s, x, y, cw, ch, fill=BG_CARD, line=BORDER, line_w=1.0)
    txt(s, x + Inches(0.3), y + Inches(0.22), Inches(1.2), Inches(0.5),
        [[(tag, F_DISPLAY, 20, MUSTARD, True)]])
    txt(s, x + Inches(0.3), y + Inches(0.72), cw - Inches(0.6), Inches(0.4),
        [[(h, F_DISPLAY, 15, TEXT_PRIMARY, True)]])
    txt(s, x + Inches(0.3), y + Inches(1.12), cw - Inches(0.6), Inches(0.85),
        [[(b, F_BODY, 11, TEXT_SECONDARY, False)]], line_spacing=1.2)
txt(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.6),
    [[("Every milestone: ", F_BODY, 11.5, MUSTARD, True),
      ("golden-file tests against pre-validated fixtures, 70%+ coverage on business logic, zero critical static findings, a Monday written update, and a milestone-close demo call on a pre-agreed script.", F_BODY, 11.5, TEXT_SECONDARY, False)]],
    line_spacing=1.2)

# ----------------------------------------------------------------------------
# Slide 8 - Why Michael + CTA
# ----------------------------------------------------------------------------
s = slide(BG_ROOT)
rect(s, 0, 0, Inches(0.18), SH, fill=MUSTARD)
eyebrow(s, Inches(0.9), Inches(0.7), "Why Michael")
txt(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.8),
    [[("Built audit-grade software before, end to end", F_DISPLAY, 30, TEXT_PRIMARY, True)]])
proofs = [
    ("2.5 yrs", "Sole developer and SME on a 60k-LOC internal platform at U.S. Bank serving 600 users a month, where correctness was never optional."),
    ("6 months", "Led that platform's migration to Azure Cloud as main rep, among the first apps the bank moved. Comfortable owning infra."),
    ("Current", "Full-stack on an Optum healthcare project (Angular + .NET + PostgreSQL, onion architecture, strict Agile). Multi-tenant rigor is daily work."),
]
cw = Inches(3.86); gap = Inches(0.24); x0 = Inches(0.9); cy = Inches(2.35); chh = Inches(2.25)
for i, (big, b) in enumerate(proofs):
    x = x0 + (cw + gap) * i
    rect(s, x, cy, cw, chh, fill=BG_CARD, line=BORDER, line_w=1.0)
    txt(s, x + Inches(0.3), cy + Inches(0.3), cw - Inches(0.6), Inches(0.6),
        [[(big, F_DISPLAY, 26, MUSTARD, True)]])
    txt(s, x + Inches(0.3), cy + Inches(1.05), cw - Inches(0.6), Inches(1.1),
        [[(b, F_BODY, 12.5, TEXT_SECONDARY, False)]], line_spacing=1.25)
# CTA band
rect(s, Inches(0.9), Inches(5.05), Inches(11.55), Inches(1.55), fill=BG_SURFACE, line=MUSTARD, line_w=1.25)
txt(s, Inches(1.25), Inches(5.28), Inches(11), Inches(0.5),
    [[("Let's lock Milestone 1.", F_DISPLAY, 22, TEXT_PRIMARY, True)]])
txt(s, Inches(1.25), Inches(5.82), Inches(11), Inches(0.7),
    [[("I welcome the payment gate and the paid probation task: send a real XML file and a REST contract and I will return the parsed response. ", F_BODY, 13, TEXT_SECONDARY, False),
      ("Open the demo:  " + DEMO_URL, F_MONO, 13, MUSTARD, True)]],
    line_spacing=1.25)

prs.save(OUT)
print("Saved", OUT)
