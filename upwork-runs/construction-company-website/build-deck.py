"""
Build the proposal deck for construction-company-website.
Uses the bespoke construction design system:
  bg:     #0D0F12  (Charcoal Black)
  amber:  #E8991A  (Safety Amber)
  white:  #F0EDE8  (Off-White Warm)
  stone:  #9C9790  (Warm Stone)
  card:   #1E2228  (Deep Slate)
  border: #2A2F38  (Steel Edge)
Fonts: Calibri (system fallback for Oswald/DM Sans - PPTX does not embed web fonts).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ---------- helpers ----------
def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# palette
BG      = rgb('#0D0F12')
SURFACE = rgb('#161A1F')
CARD    = rgb('#1E2228')
BORDER  = rgb('#2A2F38')
AMBER   = rgb('#E8991A')
WHITE   = rgb('#F0EDE8')
STONE   = rgb('#9C9790')
BLACK   = rgb('#000000')

W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def label(slide, text, l, t, w, h, size=Pt(11), bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txb

def heading(slide, text, l, t, w, h, size=Pt(36), color=WHITE, align=PP_ALIGN.LEFT):
    return label(slide, text, l, t, w, h, size=size, bold=True, color=color, align=align)

def amber_rule(slide, l, t, w, thickness=Pt(3)):
    """Thin horizontal amber rule."""
    rule = slide.shapes.add_shape(1, l, t, w, thickness)
    rule.fill.solid()
    rule.fill.fore_color.rgb = AMBER
    rule.line.fill.background()
    return rule

def eyebrow(slide, text, l, t, w):
    """Small uppercase amber label above a heading."""
    return label(slide, text.upper(), l, t, w, Inches(0.3),
                 size=Pt(10), bold=True, color=AMBER)

def para_run(tf, text, size=Pt(13), color=STONE, bold=False):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return p

# margin/grid constants
ML = Inches(0.7)    # left margin
MT = Inches(0.6)    # top margin
CW = W - Inches(1.4)  # content width
CH = H - Inches(1.2)  # content height


# ============================================================
# SLIDE 1 -- TITLE
# ============================================================
s1 = add_slide()
fill_bg(s1, BG)

# amber left bar
box(s1, Inches(0), Inches(0), Inches(0.22), H, fill_color=AMBER)

# surface panel right side
box(s1, Inches(5.5), Inches(0), Inches(7.83), H, fill_color=SURFACE)

# eyebrow
eyebrow(s1, 'Upwork Proposal  |  Full-Stack Developer, U.S.-based', Inches(0.55), Inches(1.5), Inches(4.7))

# main headline
heading(s1, 'CONSTRUCTION\nCOMPANY\nWEBSITE', Inches(0.55), Inches(2.0), Inches(4.7), Inches(3.0), size=Pt(52), color=WHITE)

# sub
label(s1, 'A four-page production-quality marketing site built before this proposal was submitted.', Inches(0.55), Inches(5.2), Inches(4.7), Inches(1.0), size=Pt(14), color=STONE)

# right panel content
label(s1, 'Michael Wegter', Inches(6.0), Inches(2.0), Inches(6.8), Inches(0.5), size=Pt(20), bold=True, color=WHITE)
label(s1, 'Full-Stack Developer', Inches(6.0), Inches(2.55), Inches(6.8), Inches(0.4), size=Pt(13), color=STONE)

amber_rule(s1, Inches(6.0), Inches(3.1), Inches(3.0))

for i, line in enumerate([
    'michaelwegter.com',
    'Demo: michaelwegter.com/demos/construction-company-website/',
    'June 15, 2026',
]):
    label(s1, line, Inches(6.0), Inches(3.3) + Inches(i*0.4), Inches(7.0), Inches(0.4), size=Pt(12), color=STONE)


# ============================================================
# SLIDE 2 -- THE DEMO
# ============================================================
s2 = add_slide()
fill_bg(s2, BG)

eyebrow(s2, 'Already Built', ML, MT, CW)
heading(s2, 'THE LIVE DEMO', ML, MT + Inches(0.35), CW, Inches(0.7), size=Pt(38))
amber_rule(s2, ML, MT + Inches(1.1), Inches(3))

label(s2, 'I built a fully working four-page construction company website before submitting this proposal.\nClick the link below to explore it live.', ML, MT + Inches(1.3), CW, Inches(0.7), size=Pt(14), color=STONE)

# demo URL box
url_box = box(s2, ML, Inches(2.6), CW, Inches(0.85), fill_color=CARD, line_color=AMBER, line_width=Pt(2))
label(s2, 'https://michaelwegter.com/demos/construction-company-website/', ML + Inches(0.2), Inches(2.7), CW - Inches(0.4), Inches(0.6), size=Pt(14), bold=True, color=AMBER)

# demo name label
label(s2, 'Cornerstone Construction Co. Website Demo', ML, Inches(3.65), CW, Inches(0.4), size=Pt(13), bold=True, color=WHITE)

# 4 page badges
badge_labels = [
    ('🎬', 'HOME', 'Video hero, stats,\nservices teaser, testimonials'),
    ('🏗️', 'SERVICES', 'Pricing tiers for all 3\nservice lines + FAQ'),
    ('📸', 'PROJECTS', '9-project gallery,\nfilter tabs, lightbox'),
    ('📋', 'CONTACT', 'Validated form,\nsuccess state'),
]

bw = (CW - Inches(0.45)) / 4
for i, (icon, title, desc) in enumerate(badge_labels):
    bx = ML + i * (bw + Inches(0.15))
    by = Inches(4.2)
    bh = Inches(2.6)
    box(s2, bx, by, bw, bh, fill_color=CARD, line_color=BORDER, line_width=Pt(1))
    # amber top border strip
    box(s2, bx, by, bw, Inches(0.07), fill_color=AMBER)
    label(s2, icon, bx, by + Inches(0.2), bw, Inches(0.4), size=Pt(22), align=PP_ALIGN.CENTER)
    label(s2, title, bx, by + Inches(0.75), bw, Inches(0.35), size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(s2, desc, bx + Inches(0.1), by + Inches(1.15), bw - Inches(0.2), Inches(1.1), size=Pt(11), color=STONE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3 -- REQUIREMENTS COVERAGE
# ============================================================
s3 = add_slide()
fill_bg(s3, BG)

eyebrow(s3, 'Every requirement, accounted for', ML, MT, CW)
heading(s3, 'WHAT YOU ASKED FOR, BUILT', ML, MT + Inches(0.35), CW, Inches(0.7), size=Pt(36))
amber_rule(s3, ML, MT + Inches(1.1), Inches(4))

reqs = [
    ('Video-Background Homepage', 'Full-viewport construction B-roll. Autoplay, muted, looped. Dark gradient overlay for legibility. JS fallback to poster image.'),
    ('Services + Pricing Page', 'Home Building ($180-$260/sq ft), Excavation ($125-$175/hr), Commercial (custom quote). Tiered pricing cards with included-items lists.'),
    ('Contact Page', 'Validated form: name, email, phone, service type, project description. Inline errors on blur. Animated success state. No page reload.'),
    ('Projects Gallery', '9 projects (Residential / Commercial / Excavation). Instant filter tabs. Full lightbox with arrow nav and ESC keyboard close.'),
    ('Construction Industry Theme', 'Charcoal black + safety amber. Oswald headlines, DM Sans body. Bespoke design system. Zero template feel.'),
    ('Professional UX/UI Polish', 'Sticky transparent-to-solid nav. Scroll fade-ins. Animated stats counters. Mobile hamburger. WCAG AA. Reduced-motion support.'),
]

rw = (CW - Inches(0.3)) / 3
rh = Inches(1.6)
for i, (title, desc) in enumerate(reqs):
    col = i % 3
    row = i // 3
    rx = ML + col * (rw + Inches(0.15))
    ry = Inches(2.0) + row * (rh + Inches(0.15))
    box(s3, rx, ry, rw, rh, fill_color=CARD, line_color=BORDER, line_width=Pt(1))
    # amber left stripe
    box(s3, rx, ry, Inches(0.06), rh, fill_color=AMBER)
    label(s3, title, rx + Inches(0.15), ry + Inches(0.12), rw - Inches(0.2), Inches(0.35), size=Pt(11), bold=True, color=WHITE)
    label(s3, desc, rx + Inches(0.15), ry + Inches(0.5), rw - Inches(0.2), Inches(1.0), size=Pt(10), color=STONE)


# ============================================================
# SLIDE 4 -- KEY DIFFERENTIATORS
# ============================================================
s4 = add_slide()
fill_bg(s4, BG)

eyebrow(s4, 'Why this demo stands out', ML, MT, CW)
heading(s4, 'KEY DIFFERENTIATORS', ML, MT + Inches(0.35), CW, Inches(0.7), size=Pt(36))
amber_rule(s4, ML, MT + Inches(1.1), Inches(3))

diffs = [
    ('1', 'Video done correctly', 'Same-origin MP4, correct autoplay attributes, overlay for legibility on every frame, and a JS error handler for fallback. Most implementations break. This one holds.'),
    ('2', 'Pricing transparency', 'Most construction sites hide prices. Firms that show ballpark numbers convert better and attract qualified leads. The Services page shows real starting figures with tier breakdowns.'),
    ('3', 'Gallery that converts', 'Filter tabs let homeowners, commercial tenants, and excavation buyers self-identify and see only relevant work. Lightbox shows project details. This is the #1 conversion driver on construction sites.'),
    ('4', 'No template feel', 'Custom design system from scratch: charcoal black, safety amber, Oswald headlines, sharp edges. Modeled on the visual language of Turner, Suffolk, and Ryan Companies.'),
    ('5', 'Mobile-first + accessible', 'Responsive across mobile, tablet, desktop. Hamburger nav below 768px. WCAG 2.1 AA: aria labels, keyboard lightbox nav, focus management, reduced-motion guards.'),
]

dh = Inches(0.98)
for i, (num, title, desc) in enumerate(diffs):
    dy = Inches(2.0) + i * (dh + Inches(0.1))
    # number badge
    box(s4, ML, dy, Inches(0.5), dh - Inches(0.05), fill_color=AMBER)
    label(s4, num, ML, dy + Inches(0.2), Inches(0.5), Inches(0.5), size=Pt(16), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    # title
    label(s4, title, ML + Inches(0.65), dy + Inches(0.05), Inches(2.5), Inches(0.4), size=Pt(13), bold=True, color=WHITE)
    # desc
    label(s4, desc, ML + Inches(0.65), dy + Inches(0.45), CW - Inches(0.65), Inches(0.45), size=Pt(11), color=STONE)


# ============================================================
# SLIDE 5 -- APPROACH / HOW I BUILD THE REAL THING
# ============================================================
s5 = add_slide()
fill_bg(s5, BG)

eyebrow(s5, 'From demo to delivery', ML, MT, CW)
heading(s5, 'HOW I\'D BUILD THE REAL THING', ML, MT + Inches(0.35), CW, Inches(0.7), size=Pt(36))
amber_rule(s5, ML, MT + Inches(1.1), Inches(3.5))

label(s5, 'The demo is the foundation. The production build adds your real assets, backend, and CMS.', ML, MT + Inches(1.3), CW, Inches(0.45), size=Pt(13), color=STONE)

phases = [
    ('WEEK 1', 'Assets + Brand Setup', 'Swap in your real logo, brand colors, and company name. Collect project photos, real testimonials, and service details. Set up domain and hosting.'),
    ('WEEK 1 TO 2', 'Content + Backend', 'Replace placeholder copy and imagery with your real content. Add backend contact form handler with email notifications. Wire up CMS for self-serve updates.'),
    ('WEEK 2 TO 3', 'Polish + Launch', 'SEO metadata for your service area, Google Analytics, cross-browser testing, performance optimization, and go-live. I stay on for feedback after launch.'),
]

pw = (CW - Inches(0.3)) / 3
ph = Inches(2.7)
py = Inches(2.2)

for i, (phase, title, desc) in enumerate(phases):
    px = ML + i * (pw + Inches(0.15))
    # card
    box(s5, px, py, pw, ph, fill_color=SURFACE, line_color=BORDER, line_width=Pt(1))
    # amber top border
    box(s5, px, py, pw, Inches(0.09), fill_color=AMBER)
    label(s5, phase, px + Inches(0.15), py + Inches(0.18), pw - Inches(0.3), Inches(0.3), size=Pt(10), bold=True, color=AMBER)
    label(s5, title, px + Inches(0.15), py + Inches(0.55), pw - Inches(0.3), Inches(0.45), size=Pt(14), bold=True, color=WHITE)
    amber_rule(s5, px + Inches(0.15), py + Inches(1.05), pw - Inches(0.3), Pt(1))
    label(s5, desc, px + Inches(0.15), py + Inches(1.2), pw - Inches(0.3), Inches(1.3), size=Pt(11), color=STONE)

# production extras note
note_items = [
    'Your real brand assets (logo, photos, copy)',
    'Backend form handler with email notifications',
    'CMS integration for self-serve content updates',
    'SEO configuration for your service area',
    'Hosting setup and domain configuration',
]
label(s5, 'Production build includes:', ML, Inches(5.45), Inches(3.5), Inches(0.35), size=Pt(11), bold=True, color=WHITE)
for i, item in enumerate(note_items):
    label(s5, '  ' + chr(8226) + '  ' + item, ML, Inches(5.85) + Inches(i * 0.3), Inches(5.5), Inches(0.28), size=Pt(11), color=STONE)


# ============================================================
# SLIDE 6 -- ABOUT MICHAEL
# ============================================================
s6 = add_slide()
fill_bg(s6, BG)

eyebrow(s6, 'Who I am', ML, MT, CW)
heading(s6, 'MICHAEL WEGTER', ML, MT + Inches(0.35), CW, Inches(0.7), size=Pt(40))
amber_rule(s6, ML, MT + Inches(1.1), Inches(3.5))

# left column: bio
bio_text = (
    "I'm a full-stack developer based in the U.S. with about 2.5 years of "
    "experience building and maintaining production applications at scale.\n\n"
    "I communicate plainly, don't disappear mid-project, and become the subject "
    "matter expert on whatever I'm working on. At U.S. Bank I took over a critical "
    "app as the sole developer two months in and owned it from there.\n\n"
    "I build quickly. If something is broken, I fix it, often within minutes."
)
txb = slide_shapes = s6.shapes.add_textbox(ML, MT + Inches(1.3), Inches(5.5), Inches(3.2))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = bio_text
run.font.size = Pt(13)
run.font.color.rgb = STONE
run.font.name = 'Calibri'

# right column: proof points
proofs = [
    ('US Bank via Turnberry', '~60k LOC React + Python + Java (TDAAS). 600 users/month. Sole developer and SME. Led Azure Cloud migration.'),
    ('Optum RHRP 4 (current)', 'Angular + .NET/C#, Agile, 150+ story points delivered. Team go-to for AI-assisted development.'),
    ('Personal Portfolio', '6 live apps on michaelwegter.com: gallery planner, SEO analyzer, Spotify + Apple Music tools, and more.'),
    ('Stack', 'React, Angular, Python, .NET/C#, SQL, PostgreSQL, Docker, Kubernetes, CI/CD.'),
]

rx = ML + Inches(6.1)
ry_start = MT + Inches(1.3)
for i, (title, desc) in enumerate(proofs):
    ry = ry_start + i * Inches(1.35)
    box(s6, rx, ry, Inches(6.5), Inches(1.2), fill_color=CARD, line_color=BORDER, line_width=Pt(1))
    box(s6, rx, ry, Inches(0.05), Inches(1.2), fill_color=AMBER)
    label(s6, title, rx + Inches(0.18), ry + Inches(0.1), Inches(6.1), Inches(0.35), size=Pt(12), bold=True, color=WHITE)
    label(s6, desc, rx + Inches(0.18), ry + Inches(0.48), Inches(6.1), Inches(0.6), size=Pt(11), color=STONE)


# ============================================================
# SLIDE 7 -- CLOSING CTA
# ============================================================
s7 = add_slide()
fill_bg(s7, BG)

# large amber strip across middle
box(s7, Inches(0), Inches(2.5), W, Inches(2.6), fill_color=AMBER)

label(s7, 'READY TO SEE IT LIVE?', ML, Inches(2.8), CW, Inches(0.7), size=Pt(44), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
label(s7, 'The demo is already built. Click it, explore all four pages, then reach out.', ML, Inches(3.55), CW, Inches(0.45), size=Pt(15), color=BLACK, align=PP_ALIGN.CENTER)

# bottom content on dark
label(s7, 'https://michaelwegter.com/demos/construction-company-website/', ML, Inches(5.4), CW, Inches(0.5), size=Pt(16), bold=True, color=AMBER, align=PP_ALIGN.CENTER)

label(s7, 'Michael Wegter   |   michaelwegter.com   |   U.S.-based Full-Stack Developer', ML, Inches(6.4), CW, Inches(0.4), size=Pt(12), color=STONE, align=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
out_path = '/Users/michaelwegter/Desktop/Projects/upwork-agentic-workflow/upwork-runs/construction-company-website/proposal-deck.pptx'
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Slides: {len(prs.slides)}')
