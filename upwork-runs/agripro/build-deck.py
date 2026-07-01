#!/usr/bin/env python3
"""Build proposal/deck.pptx for AgriPro using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Paths
BASE = os.path.dirname(os.path.abspath(__file__))
MEDIA = os.path.join(BASE, "proposal", "media")
OUT = os.path.join(BASE, "proposal", "deck.pptx")

# Colours (agri-industrial palette)
NAV      = RGBColor(0x0F, 0x2A, 0x1A)   # dark forest green
PRIMARY  = RGBColor(0x1B, 0x43, 0x32)   # agronomy deep green
ACCENT   = RGBColor(0xB5, 0x65, 0x1D)   # earth amber
ACTION   = RGBColor(0x1D, 0x4E, 0xD8)   # strong blue
PASS_CLR = RGBColor(0x15, 0x80, 0x3D)   # status pass
FLAG_CLR = RGBColor(0xB4, 0x53, 0x09)   # status flag
BG       = RGBColor(0xF5, 0xF2, 0xEC)   # warm parchment
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT     = RGBColor(0x1C, 0x19, 0x17)
MUTED    = RGBColor(0x57, 0x53, 0x4E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=16, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "IBM Plex Sans"
    return txb


def title_slide():
    slide = prs.slides.add_slide(blank_layout)
    # Full-bleed dark green background
    add_rect(slide, 0, 0, 13.33, 7.5, fill=NAV)
    # Accent bar at bottom
    add_rect(slide, 0, 7.0, 13.33, 0.5, fill=ACCENT)
    # Title
    add_text(slide, "AgriPro Operations Console",
             1.0, 1.8, 11.33, 1.4, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Subtitle
    add_text(slide, "Full-Featured AI-Powered Procurement Operations Platform",
             1.0, 3.2, 11.33, 0.8, size=20, color=RGBColor(0x9E, 0xCE, 0x9E), align=PP_ALIGN.CENTER)
    # Demo link
    add_text(slide, "michaelwegter.com/demos/agripro/",
             1.0, 4.1, 11.33, 0.6, size=15, color=RGBColor(0xC8, 0xD9, 0xC8), align=PP_ALIGN.CENTER, italic=True)
    # Author
    add_text(slide, "Proposal by Michael Wegter  |  github.com/mwegter95",
             1.0, 5.2, 11.33, 0.5, size=13, color=MUTED, align=PP_ALIGN.CENTER)


def problem_slide():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BG)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=PRIMARY)
    add_text(slide, "The Problem", 0.5, 0.15, 12, 0.7, size=22, bold=True, color=WHITE)

    pain_points = [
        ("Spreadsheets driving multi-state ops",
         "Field inspections, lab results, and procurement approvals live in disconnected spreadsheets and email threads across states."),
        ("No visibility across teams",
         "Procurement, quality, warehouse, and management each work in silos. No shared source of truth, no audit trail, no escalation path."),
        ("Manual document processing",
         "Lab certificates are read by hand, re-keyed into spreadsheets. OCR extraction, classification, and anomaly flags do not exist today."),
        ("Zero governance",
         "No approval gate enforces quality standards before a lot moves to warehouse. Anomalous values can slip through undetected."),
    ]
    for i, (title, body) in enumerate(pain_points):
        col = 0.4 + (i % 2) * 6.5
        row = 1.3 + (i // 2) * 2.6
        add_rect(slide, col, row, 6.2, 2.3, fill=WHITE, line=RGBColor(0xD6, 0xD3, 0xC8))
        add_rect(slide, col, row, 0.18, 2.3, fill=ACCENT)
        add_text(slide, title, col + 0.35, row + 0.12, 5.7, 0.5, size=13, bold=True, color=PRIMARY)
        add_text(slide, body, col + 0.35, row + 0.6, 5.7, 1.5, size=11, color=MUTED)


def approach_slide():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BG)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=PRIMARY)
    add_text(slide, "The Approach", 0.5, 0.15, 12, 0.7, size=22, bold=True, color=WHITE)

    milestones = [
        ("M1", "Core Operational Spine", "Weeks 1-8",
         "Authenticated RBAC, procurement campaigns, field inspection records, 6-stage approval pipeline, search and filtering, audit logs, mobile layout, admin portal."),
        ("M2", "AI and Quality Layer", "Weeks 9-12",
         "Production OCR on lab cert ingestion, document classification, configurable anomaly detection, warehouse allocation, document management, notifications, role dashboards with PDF and Excel export."),
        ("M3", "Hardening and Delivery", "Weeks 13-16",
         "Unit and E2E test suite (Vitest + Playwright), technical documentation, Docker + Nginx deploy, CI/CD pipeline, knowledge transfer, 30-day post-launch stabilization."),
    ]
    for i, (num, title, dur, body) in enumerate(milestones):
        y = 1.3 + i * 1.9
        add_rect(slide, 0.5, y, 12.3, 1.65, fill=WHITE, line=RGBColor(0xD6, 0xD3, 0xC8))
        add_rect(slide, 0.5, y, 0.22, 1.65, fill=PRIMARY)
        add_text(slide, num, 0.8, y + 0.12, 0.8, 0.4, size=12, bold=True, color=PRIMARY)
        add_text(slide, title, 1.5, y + 0.08, 7.0, 0.45, size=14, bold=True, color=PRIMARY)
        add_text(slide, dur, 9.5, y + 0.08, 3.0, 0.45, size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
        add_text(slide, body, 1.5, y + 0.6, 11.0, 0.95, size=11, color=MUTED)


def demo_slide():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=NAV)
    add_rect(slide, 0, 0, 13.33, 0.85, fill=PRIMARY)
    add_text(slide, "The Live Demo", 0.5, 0.1, 10, 0.65, size=22, bold=True, color=WHITE)
    add_text(slide, "michaelwegter.com/demos/agripro/", 8.5, 0.1, 4.5, 0.65,
             size=13, color=RGBColor(0xC8, 0xD9, 0xC8), align=PP_ALIGN.RIGHT, italic=True)

    hero_path = os.path.join(MEDIA, "hero.png")
    if os.path.exists(hero_path):
        slide.shapes.add_picture(hero_path, Inches(0.4), Inches(1.0), Inches(8.4), Inches(5.8))

    add_rect(slide, 9.1, 1.0, 4.0, 1.5, fill=RGBColor(0x0A, 0x1F, 0x12))
    add_text(slide, "What you can do right now:", 9.2, 1.05, 3.8, 0.4, size=11, bold=True, color=ACCENT)
    add_text(slide, "Log in as any of 4 roles\nDrop a grain cert to trigger OCR\nWatch anomaly flags fire in real time\nApprove / Flag / Reject inspections\nExport a dashboard PDF or Excel file",
             9.2, 1.45, 3.8, 2.2, size=11, color=RGBColor(0xC8, 0xD9, 0xC8))

    add_rect(slide, 9.1, 2.7, 4.0, 4.1, fill=RGBColor(0x0A, 0x1F, 0x12))
    step1_path = os.path.join(MEDIA, "step-1.png")
    step2_path = os.path.join(MEDIA, "step-2.png")
    if os.path.exists(step1_path):
        slide.shapes.add_picture(step1_path, Inches(9.15), Inches(2.75), Inches(3.9), Inches(1.85))
    add_text(slide, "OCR Certificate Scanner", 9.2, 4.62, 3.8, 0.35, size=10, bold=True, color=ACCENT)
    if os.path.exists(step2_path):
        slide.shapes.add_picture(step2_path, Inches(9.15), Inches(5.05), Inches(3.9), Inches(1.75))
    add_text(slide, "Anomaly Detection", 9.2, 6.82, 3.8, 0.35, size=10, bold=True, color=FLAG_CLR)


def requirements_slide():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BG)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=PRIMARY)
    add_text(slide, "Requirement Coverage", 0.5, 0.15, 12, 0.7, size=22, bold=True, color=WHITE)

    rows_left = [
        ("R1", "Secure auth + RBAC", "Live Demo"),
        ("R2", "Procurement campaigns", "Live Demo"),
        ("R3", "Field inspection workflows", "Live Demo"),
        ("R4", "Lab testing + quality grading", "Live Demo"),
        ("R5", "Multi-stage approval workflows", "Live Demo"),
        ("R6", "Warehouse allocation", "Rich Mock"),
        ("R7", "Document management", "Rich Mock"),
        ("R8", "Search and filtering", "Live Demo"),
    ]
    rows_right = [
        ("R9", "Role-aware dashboards", "Live Demo"),
        ("R10", "PDF and Excel export", "Live Demo"),
        ("R11", "Notifications + reminders", "Rich Mock"),
        ("R12", "Audit logs", "Rich Mock"),
        ("R13", "Mobile-responsive", "Live Demo"),
        ("R14", "OCR for lab certs", "Live Demo"),
        ("R15", "Intelligent classification", "Live Demo"),
        ("R16", "Anomaly detection", "Live Demo"),
    ]

    pill_colors = {"Live Demo": PASS_CLR, "Rich Mock": FLAG_CLR, "Proposal": ACTION}

    def render_table(rows, start_x, start_y):
        for i, (rid, req, status) in enumerate(rows):
            y = start_y + i * 0.72
            row_bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xED, 0xE7)
            add_rect(slide, start_x, y, 5.9, 0.66, fill=row_bg)
            add_text(slide, rid, start_x + 0.08, y + 0.1, 0.45, 0.45, size=10,
                     bold=True, color=ACCENT)
            add_text(slide, req, start_x + 0.55, y + 0.1, 3.4, 0.45, size=11, color=TEXT)
            pill_c = pill_colors.get(status, MUTED)
            add_rect(slide, start_x + 4.05, y + 0.1, 1.65, 0.44, fill=pill_c)
            add_text(slide, status, start_x + 4.08, y + 0.12, 1.6, 0.4, size=10,
                     bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    render_table(rows_left, 0.5, 1.15)
    render_table(rows_right, 6.93, 1.15)


def why_michael_slide():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BG)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=PRIMARY)
    add_text(slide, "Why Michael Wegter", 0.5, 0.15, 12, 0.7, size=22, bold=True, color=WHITE)

    proofs = [
        ("Sole developer, 600-user internal platform",
         "At U.S. Bank I owned TDAAS within months and kept a 60,000-line React/Python/SQL/Java platform running for 600 users per month. If something broke, I fixed it, often within 10 minutes."),
        ("AI in real production workflows",
         "I use GitHub Copilot and Azure OpenAI in daily delivery. At Optum I am the team go-to for AI-assisted development. The OCR, classification, and anomaly detection in this demo are the same class of AI I ship at work."),
        ("Full-stack, cloud-native, and deployment-ready",
         "React, Angular, Python, .NET/C#, Java, SQL, PostgreSQL, Docker, Kubernetes, CI/CD. I led the TDAAS Azure Cloud migration as the main technical rep, among the first apps the bank moved. I own the full stack."),
    ]
    for i, (title, body) in enumerate(proofs):
        x = 0.5 + i * 4.28
        add_rect(slide, x, 1.3, 4.05, 4.8, fill=WHITE, line=RGBColor(0xD6, 0xD3, 0xC8))
        add_rect(slide, x, 1.3, 4.05, 0.18, fill=ACCENT)
        add_text(slide, title, x + 0.2, 1.55, 3.65, 0.75, size=13, bold=True, color=PRIMARY)
        add_text(slide, body, x + 0.2, 2.45, 3.65, 3.5, size=11, color=MUTED)

    add_rect(slide, 0.5, 6.4, 12.33, 0.72, fill=NAV)
    add_text(slide,
             "150+ story points delivered at Optum (current)  |  "
             "6-month Azure migration lead at U.S. Bank  |  "
             "Full portfolio: michaelwegter.com  |  github.com/mwegter95",
             0.7, 6.48, 12.0, 0.55, size=11, color=RGBColor(0xC8, 0xD9, 0xC8), align=PP_ALIGN.CENTER)


def scope_slide():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BG)
    add_rect(slide, 0, 0, 13.33, 1.0, fill=PRIMARY)
    add_text(slide, "Scope, Timeline, and Next Step", 0.5, 0.15, 12, 0.7, size=22, bold=True, color=WHITE)

    # Tech stack
    add_rect(slide, 0.5, 1.2, 6.0, 3.5, fill=WHITE, line=RGBColor(0xD6, 0xD3, 0xC8))
    add_text(slide, "Production Tech Stack", 0.8, 1.35, 5.5, 0.45, size=13, bold=True, color=PRIMARY)
    stack_items = [
        "Frontend: React 18 (or Angular per preference) + Vite",
        "Backend: Python (FastAPI) or Node.js",
        "Database: PostgreSQL with Row-Level Security for RBAC",
        "OCR: Tesseract + cloud fallback (AWS Textract / Azure DI)",
        "AI: OpenAI or Azure OpenAI for classification, anomaly scoring",
        "Deploy: Docker + Nginx, CI/CD (GitHub Actions or similar)",
        "Export: jsPDF + SheetJS (or server-side Puppeteer)",
    ]
    for i, item in enumerate(stack_items):
        add_text(slide, item, 0.8, 1.9 + i * 0.37, 5.5, 0.35, size=11, color=MUTED)

    # Timeline
    add_rect(slide, 7.0, 1.2, 6.0, 3.5, fill=WHITE, line=RGBColor(0xD6, 0xD3, 0xC8))
    add_text(slide, "Timeline (Indicative)", 7.3, 1.35, 5.5, 0.45, size=13, bold=True, color=PRIMARY)
    timeline_items = [
        ("Weeks 1-8",   "M1: Core spine, RBAC, inspection + approval pipeline"),
        ("Weeks 9-12",  "M2: AI layer, OCR, anomaly detection, dashboards"),
        ("Weeks 13-16", "M3: Hardening, testing, documentation, deploy"),
        ("Day 1+",      "30-day post-launch stabilization included"),
    ]
    for i, (wk, desc) in enumerate(timeline_items):
        y = 1.9 + i * 0.62
        add_text(slide, wk, 7.3, y, 1.6, 0.5, size=11, bold=True, color=ACCENT)
        add_text(slide, desc, 9.0, y, 3.8, 0.5, size=11, color=MUTED)

    # CTA
    add_rect(slide, 0.5, 5.0, 12.33, 2.0, fill=NAV)
    add_text(slide, "Next Step", 0.9, 5.12, 11.5, 0.45, size=14, bold=True, color=ACCENT)
    add_text(slide,
             "Try the live demo: michaelwegter.com/demos/agripro/\n"
             "Happy to walk through it together and talk about your timeline, stack preferences, and budget.\n"
             "Contact: github.com/mwegter95",
             0.9, 5.62, 11.5, 1.25, size=12, color=RGBColor(0xC8, 0xD9, 0xC8))


# Build all slides
title_slide()
problem_slide()
approach_slide()
demo_slide()
requirements_slide()
why_michael_slide()
scope_slide()

prs.save(OUT)
print(f"Saved: {OUT}")
