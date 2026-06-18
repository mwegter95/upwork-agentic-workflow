#!/usr/bin/env python3
"""Build the proposal deck for the AI Chatbot for Customer Support Upwork run."""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
MEDIA = os.path.join(HERE, "media")
OUT = os.path.join(HERE, "deck.pptx")

# SaaS support-tool palette (blue/slate, no portfolio theme)
BG       = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
SURFACE  = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
CARD     = RGBColor(0x27, 0x35, 0x49)   # slate-700 approx
BORDER   = RGBColor(0x33, 0x41, 0x55)   # slate-700
BLUE     = RGBColor(0x37, 0x99, 0xF8)   # blue-400 (on dark)
BLUE_BTN = RGBColor(0x25, 0x63, 0xEB)   # blue-600
GREEN    = RGBColor(0x10, 0xB9, 0x81)   # green-500
TEXT     = RGBColor(0xF8, 0xFA, 0xFC)   # slate-50
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
TEXT_MID = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
INK      = RGBColor(0x0F, 0x17, 0x2A)

FONT_HDR  = "DM Sans"
FONT_BODY = "Inter"

DEMO_URL  = "michaelwegter.com/demos/ai-chatbot-for-customer-support/"
DEMO_FULL = "https://" + DEMO_URL

def no_dashes(s):
    if re.search(r"[–—]", s):
        raise AssertionError(f"Em/en dash in: {s!r}")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width
H = prs.slide_height

def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False

def rect(slide, x, y, w, h, fill=SURFACE, line_color=None, lw=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line_color: r.line.color.rgb = line_color; r.line.width = lw or Pt(1)
    else: r.line.fill.background()
    r.shadow.inherit = False
    return r

def txt(slide, x, y, w, h, text, *, font=FONT_BODY, size=16, color=TEXT,
        bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    no_dashes(text)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color
    return tb

def eyebrow(slide, x, y, w, text, color=GREEN):
    txt(slide, x, y, w, Inches(0.28), text.upper(),
        font=FONT_BODY, size=10, color=color, bold=True)

def footer(slide, n, total=9):
    txt(slide, Inches(0.5), Inches(7.1), Inches(9), Inches(0.28),
        "MICHAEL WEGTER  /  AI CUSTOMER SUPPORT CHATBOT  /  PROPOSAL",
        size=8.5, color=TEXT_DIM, font=FONT_BODY)
    txt(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.28),
        f"{n} / {total}", size=8.5, color=TEXT_DIM,
        font=FONT_BODY, align=PP_ALIGN.RIGHT)

def add_image(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)

# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
# accent bar top
rect(s, 0, 0, W, Inches(0.06), fill=BLUE_BTN)
# content
txt(s, Inches(1), Inches(1.5), Inches(8), Inches(0.35),
    "PROPOSAL FOR UPWORK", size=11, color=GREEN, bold=True, font=FONT_BODY)
txt(s, Inches(1), Inches(2.0), Inches(8.5), Inches(1.4),
    "AI Customer Support Chatbot", size=44, font=FONT_HDR, bold=True, color=TEXT)
txt(s, Inches(1), Inches(3.55), Inches(8), Inches(0.7),
    "In-browser RAG chatbot with FAQ knowledge base, live order lookup,\nhuman escalation, and admin dashboard. Zero recurring AI cost.",
    size=17, color=TEXT_MID, font=FONT_BODY)
# Demo badge
rect(s, Inches(9.5), Inches(1.8), Inches(3.3), Inches(1.9), fill=CARD, line_color=BORDER)
txt(s, Inches(9.7), Inches(2.0), Inches(2.9), Inches(0.28),
    "LIVE DEMO", size=9.5, color=GREEN, bold=True, font=FONT_BODY)
txt(s, Inches(9.7), Inches(2.36), Inches(2.9), Inches(1.0),
    DEMO_URL, size=12, color=BLUE, font=FONT_BODY)
# Author
txt(s, Inches(1), Inches(6.2), Inches(6), Inches(0.35),
    "Michael Wegter  |  Full-Stack Developer  |  United States",
    size=13, color=TEXT_DIM, font=FONT_BODY)
footer(s, 1)

# ── Slide 2: The Problem ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.9), Inches(0.65), Inches(6), "The Problem")
txt(s, Inches(0.9), Inches(1.1), Inches(11), Inches(1.0),
    "Your small team is manually answering the same questions every day",
    size=34, font=FONT_HDR, bold=True, color=TEXT)
# Pain points grid
pains = [
    ("Repetitive volume", "Orders, shipping, returns. Same 20 questions, every day, from every customer."),
    ("Business-hours only", "Customers waiting until 9am ET for answers they could get in seconds."),
    ("Agent capacity", "Your team's time goes to tickets a bot could close, not the ones that need a human."),
    ("Integration gaps", "Fragmented stack: website, help desk, order system. No unified first-line response."),
]
cols = [Inches(0.9), Inches(6.8)]
for i, (title, body) in enumerate(pains):
    col = cols[i % 2]
    row = Inches(2.45) if i < 2 else Inches(4.5)
    rect(s, col, row, Inches(5.6), Inches(1.65), fill=SURFACE, line_color=BORDER, lw=Pt(0.75))
    txt(s, col + Inches(0.22), row + Inches(0.18), Inches(5.0), Inches(0.3),
        title, size=14, font=FONT_HDR, bold=True, color=TEXT)
    txt(s, col + Inches(0.22), row + Inches(0.58), Inches(5.0), Inches(0.9),
        body, size=13, color=TEXT_DIM, font=FONT_BODY)
footer(s, 2)

# ── Slide 3: The Solution ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.9), Inches(0.65), Inches(6), "The Solution")
txt(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(0.85),
    "SupportAI: a fully in-browser AI chatbot, zero recurring cost",
    size=34, font=FONT_HDR, bold=True, color=TEXT)
txt(s, Inches(0.9), Inches(2.1), Inches(11), Inches(0.5),
    "The AI model runs in the customer's browser using WebAssembly. No OpenAI. No Anthropic. No monthly token bill.\nYour FAQ knowledge base stays in a JSON file your team can edit without developer help.",
    size=14, color=TEXT_MID, font=FONT_BODY)

features = [
    ("📚", "RAG Knowledge Base", "20+ FAQ entries retrieved by semantic similarity. Bot cites sources. Grounded answers, not hallucinations."),
    ("📦", "Order Lookup", "Parses order numbers from chat. Queries your API (mock in demo). Returns status, carrier, ETA."),
    ("🙋", "Human Escalation", "Business-hours routing (9am to 6pm ET). Ticket confirmation. Zendesk handoff with full chat context."),
    ("⚙️", "Admin Dashboard", "Add/edit FAQ entries. Analytics: deflection rate, top queries, escalations. Maintainable without a developer."),
]
for i, (icon, title, body) in enumerate(features):
    col = cols[i % 2]
    row = Inches(2.85) if i < 2 else Inches(4.75)
    rect(s, col, row, Inches(5.6), Inches(1.65), fill=SURFACE, line_color=BORDER, lw=Pt(0.75))
    txt(s, col + Inches(0.18), row + Inches(0.18), Inches(0.5), Inches(0.4), icon, size=18)
    txt(s, col + Inches(0.65), row + Inches(0.18), Inches(4.8), Inches(0.35),
        title, size=14, font=FONT_HDR, bold=True, color=TEXT)
    txt(s, col + Inches(0.65), row + Inches(0.60), Inches(4.8), Inches(0.85),
        body, size=12.5, color=TEXT_DIM, font=FONT_BODY)
footer(s, 3)

# ── Slide 4: Demo Screenshots ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.9), Inches(0.55), Inches(6), "Live Demo")
txt(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.7),
    "Built for this proposal. Click it yourself.",
    size=32, font=FONT_HDR, bold=True, color=TEXT)
txt(s, Inches(0.9), Inches(1.8), Inches(11), Inches(0.35),
    DEMO_FULL, size=14, color=BLUE, font=FONT_BODY)

pics = [
    ("01-hero.png", "Hero: 4 live features on first paint"),
    ("02-chat-ready.png", "Chat UI with quick actions + source citations"),
    ("03-admin-dashboard.png", "Admin: analytics + KB editor"),
]
pw = Inches(3.8); ph = Inches(4.2)
for i, (fname, cap) in enumerate(pics):
    px = Inches(0.55) + i * (pw + Inches(0.2))
    py = Inches(2.3)
    rect(s, px, py, pw, ph + Inches(0.36), fill=CARD, line_color=BORDER, lw=Pt(0.75))
    add_image(s, os.path.join(MEDIA, fname), px + Inches(0.05), py + Inches(0.05), pw - Inches(0.1), ph)
    txt(s, px + Inches(0.1), py + ph + Inches(0.08), pw - Inches(0.2), Inches(0.28),
        cap, size=10.5, color=TEXT_DIM, font=FONT_BODY)
footer(s, 4)

# ── Slide 5: Architecture ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.9), Inches(0.55), Inches(6), "Architecture")
txt(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.7),
    "In-browser RAG pipeline: no backend required",
    size=32, font=FONT_HDR, bold=True, color=TEXT)

nodes = [
    ("User Input", "Customer types a message in the chat widget"),
    ("Retrieval", "Transformers.js embeddings rank top-3 FAQ chunks by semantic similarity"),
    ("Generation", "Qwen2.5-0.5B WASM model generates a grounded answer, streamed back"),
    ("Response", "Answer shown with source citations. Order and escalation handled in parallel."),
]
nw = Inches(2.65); nh = Inches(2.2)
for i, (title, body) in enumerate(nodes):
    nx = Inches(0.55) + i * (nw + Inches(0.55))
    ny = Inches(2.1)
    rect(s, nx, ny, nw, nh, fill=SURFACE, line_color=BORDER, lw=Pt(0.75))
    txt(s, nx + Inches(0.18), ny + Inches(0.2), nw - Inches(0.36), Inches(0.35),
        str(i + 1), size=22, font=FONT_HDR, bold=True, color=BLUE)
    txt(s, nx + Inches(0.18), ny + Inches(0.72), nw - Inches(0.36), Inches(0.4),
        title, size=14, font=FONT_HDR, bold=True, color=TEXT)
    txt(s, nx + Inches(0.18), ny + Inches(1.2), nw - Inches(0.36), Inches(0.85),
        body, size=12, color=TEXT_DIM, font=FONT_BODY)
    if i < 3:
        txt(s, nx + nw + Inches(0.1), ny + Inches(0.9), Inches(0.35), Inches(0.4),
            ">", size=24, color=BORDER, font=FONT_HDR, bold=True)

# Note row
rect(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(1.1), fill=CARD, line_color=BORDER, lw=Pt(0.75))
txt(s, Inches(0.85), Inches(4.72), Inches(2.2), Inches(0.35),
    "Zero recurring cost", size=13, font=FONT_HDR, bold=True, color=GREEN)
txt(s, Inches(3.2), Inches(4.72), Inches(9.3), Inches(0.7),
    "The LLM runs in the visitor's browser via WebAssembly. No OpenAI, no Anthropic, no cloud inference bill. The FAQ is a plain JSON file. Order status calls your existing endpoint. Zendesk handoff uses the Web Widget SDK.",
    size=12, color=TEXT_MID, font=FONT_BODY)
footer(s, 5)

# ── Slide 6: Zendesk Integration ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.9), Inches(0.55), Inches(6), "Zendesk Integration Design")
txt(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.75),
    "Smooth handoff to your support team",
    size=32, font=FONT_HDR, bold=True, color=TEXT)
txt(s, Inches(0.9), Inches(1.9), Inches(11), Inches(0.45),
    "The demo simulates this flow. In production, step 3 is a live Zendesk Web Widget SDK call.",
    size=14, color=TEXT_MID, font=FONT_BODY)

steps = [
    ("Trigger detected", "Bot identifies escalation signal: customer says 'talk to a person,' repeated confusion, or unresolved question."),
    ("Business hours check", "Current time checked against 9am to 6pm ET. Live chat offered during hours; next available window shown after hours."),
    ("Zendesk handoff", "Zendesk Web Widget opens pre-filled with customer name, issue summary, and full conversation transcript."),
    ("Ticket confirmed", "Ticket ID shown to customer. Bot remains open for follow-up questions while agent reviews the case."),
]
sw = Inches(2.8); sh = Inches(2.6)
for i, (title, body) in enumerate(steps):
    sx = Inches(0.6) + i * (sw + Inches(0.35))
    sy = Inches(2.55)
    rect(s, sx, sy, sw, sh, fill=SURFACE, line_color=BORDER, lw=Pt(0.75))
    rect(s, sx, sy, sw, Inches(0.07), fill=BLUE_BTN)
    txt(s, sx + Inches(0.18), sy + Inches(0.2), sw - Inches(0.36), Inches(0.3),
        f"STEP {i + 1}", size=9.5, color=GREEN, bold=True, font=FONT_BODY)
    txt(s, sx + Inches(0.18), sy + Inches(0.6), sw - Inches(0.36), Inches(0.5),
        title, size=14, font=FONT_HDR, bold=True, color=TEXT)
    txt(s, sx + Inches(0.18), sy + Inches(1.15), sw - Inches(0.36), Inches(1.3),
        body, size=12.5, color=TEXT_DIM, font=FONT_BODY)
    if i < 3:
        txt(s, sx + sw + Inches(0.05), sy + Inches(1.1), Inches(0.26), Inches(0.4),
            ">", size=22, color=BORDER, font=FONT_HDR, bold=True)
footer(s, 6)

# ── Slide 7: Timeline + Pricing ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.9), Inches(0.55), Inches(6), "Timeline + Pricing")
txt(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.75),
    "Fixed price. Scope is clear, delivery is bounded.",
    size=32, font=FONT_HDR, bold=True, color=TEXT)

# Pricing cards
cards = [
    (True, "Base Build", "$2,500", "In-browser chatbot, FAQ KB, order lookup, escalation, admin panel, embed code + docs"),
    (False, "With Live Zendesk", "$3,500", "Everything above plus real Zendesk ticket creation, live chat handoff, and setup guide"),
    (False, "Full Production", "$4,500", "Everything above plus multi-language KB, Zendesk app config, custom branding, post-launch support"),
]
cw = Inches(3.7)
for i, (featured, label, price, note) in enumerate(cards):
    cx = Inches(0.7) + i * (cw + Inches(0.35))
    cy = Inches(2.1)
    ch = Inches(2.5)
    fill = RGBColor(0x1E, 0x3A, 0x5F) if featured else SURFACE
    lc = BLUE_BTN if featured else BORDER
    rect(s, cx, cy, cw, ch, fill=fill, line_color=lc, lw=Pt(1.5 if featured else 0.75))
    txt(s, cx + Inches(0.22), cy + Inches(0.2), cw - Inches(0.44), Inches(0.3),
        label.upper(), size=10, color=GREEN if featured else TEXT_DIM, bold=True, font=FONT_BODY)
    txt(s, cx + Inches(0.22), cy + Inches(0.65), cw - Inches(0.44), Inches(0.7),
        price, size=36, font=FONT_HDR, bold=True, color=BLUE if featured else TEXT)
    txt(s, cx + Inches(0.22), cy + Inches(1.45), cw - Inches(0.44), Inches(0.95),
        note, size=12, color=TEXT_DIM, font=FONT_BODY)

# Timeline table
ty = Inches(4.85)
phases = [
    ("Architecture + knowledge base setup", "1 day"),
    ("Core chatbot + RAG + in-browser LLM", "3 days"),
    ("Order lookup + escalation flow", "1 day"),
    ("Admin panel + analytics", "1 day"),
    ("Zendesk integration + docs", "0.5 days"),
    ("QA + polish + embed package", "1 day"),
]
rect(s, Inches(0.7), ty, Inches(11.9), Inches(2.25), fill=SURFACE, line_color=BORDER, lw=Pt(0.75))
for i, (phase, days) in enumerate(phases):
    row_y = ty + Inches(0.05) + i * Inches(0.34)
    if i % 2 == 1:
        rect(s, Inches(0.7), row_y, Inches(11.9), Inches(0.34), fill=CARD)
    txt(s, Inches(0.9), row_y + Inches(0.05), Inches(9.5), Inches(0.28),
        phase, size=12, color=TEXT_MID, font=FONT_BODY)
    txt(s, Inches(11.0), row_y + Inches(0.05), Inches(1.4), Inches(0.28),
        days, size=12, color=GREEN, font=FONT_BODY, bold=True, align=PP_ALIGN.RIGHT)
footer(s, 7)

# ── Slide 8: About Michael ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.9), Inches(0.55), Inches(6), "About Michael Wegter")
txt(s, Inches(0.9), Inches(1.0), Inches(11), Inches(0.75),
    "Full-stack developer, currently based in the United States",
    size=32, font=FONT_HDR, bold=True, color=TEXT)
txt(s, Inches(0.9), Inches(1.95), Inches(10), Inches(0.55),
    "I confirmed: I am currently based in the United States. I have built the stack this project requires and ship AI-assisted features daily.",
    size=14, color=TEXT_MID, font=FONT_BODY)

proof = [
    ("2.5 years at U.S. Bank", "Primary developer on a React + Python platform serving 600 users/month. Became sole developer and SME within months. Fixed incidents in under 10 minutes. 60,000+ LOC."),
    ("Azure Cloud Migration", "Led a 6-month migration to Azure as the project's main representative, among the first internal apps the bank migrated."),
    ("Optum (current)", "Angular + .NET + PostgreSQL on a large Agile team. 150+ story points delivered. Team's go-to for AI-assisted development. Contract won April 2026."),
    ("AI delivery, daily", "Built this demo's in-browser RAG chatbot in a single session using Transformers.js, the same stack I would use on your production build."),
]
for i, (title, body) in enumerate(proof):
    col = Inches(0.9) if i % 2 == 0 else Inches(7.1)
    row = Inches(2.75) if i < 2 else Inches(4.65)
    pw = Inches(5.7); ph = Inches(1.65)
    rect(s, col, row, pw, ph, fill=SURFACE, line_color=BORDER, lw=Pt(0.75))
    txt(s, col + Inches(0.2), row + Inches(0.18), pw - Inches(0.4), Inches(0.35),
        title, size=14, font=FONT_HDR, bold=True, color=TEXT)
    txt(s, col + Inches(0.2), row + Inches(0.6), pw - Inches(0.4), Inches(0.9),
        body, size=12, color=TEXT_DIM, font=FONT_BODY)
footer(s, 8)

# ── Slide 9: CTA ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, H, fill=RGBColor(0x0D, 0x14, 0x26))
rect(s, 0, 0, W, Inches(0.08), fill=BLUE_BTN)
rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=BLUE_BTN)
txt(s, Inches(1.5), Inches(1.6), Inches(10), Inches(0.35),
    "NEXT STEP", size=11, color=GREEN, bold=True, font=FONT_BODY,
    align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(2.1), Inches(10), Inches(1.2),
    "Try the demo. Let's talk about the real build.",
    size=40, font=FONT_HDR, bold=True, color=TEXT,
    align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(3.55), Inches(10), Inches(0.55),
    DEMO_FULL,
    size=18, color=BLUE, font=FONT_BODY, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.6),
    "Ask about orders (ORD-1001), returns, shipping. Try escalating to a human. Open the admin panel (admin123).",
    size=14, color=TEXT_DIM, font=FONT_BODY, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
    "Michael Wegter  |  michaelwegter.com  |  United States",
    size=13, color=TEXT_MID, font=FONT_BODY, align=PP_ALIGN.CENTER)
footer(s, 9)

prs.save(OUT)
print(f"Deck saved: {OUT}")
