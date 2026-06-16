"""
Generate the proposal deck for EdTech School Connect Web Portal.
Design: Bespoke EdTech palette (blue/green/warm-white), Nunito + Calibri fallback.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color palette (from plan.md / demo design system) ──────────────────────
C_PRIMARY     = RGBColor(0x1A, 0x56, 0xDB)   # #1A56DB
C_PRIMARY_DK  = RGBColor(0x1E, 0x42, 0x9F)   # #1E429F
C_ACCENT      = RGBColor(0x0E, 0x9F, 0x6E)   # #0E9F6E
C_WARNING     = RGBColor(0xE3, 0xA0, 0x08)   # #E3A008
C_BG          = RGBColor(0xF9, 0xFA, 0xFB)   # #F9FAFB
C_CARD        = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
C_BORDER      = RGBColor(0xE5, 0xE7, 0xEB)   # #E5E7EB
C_TEXT        = RGBColor(0x11, 0x18, 0x27)   # #111827
C_MUTED       = RGBColor(0x6B, 0x72, 0x80)   # #6B7280
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)   # #DBEAFE (very light blue)
C_GRAD_END    = RGBColor(0x16, 0x3A, 0x8A)   # deepened primary for gradient effect


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout

# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def add_label(slide, text, left, top, width=Inches(6)):
    """Small all-caps eyebrow label in primary blue."""
    txb = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = C_PRIMARY
    return txb


def hex_pill(slide, text, left, top, fill, text_color=None):
    """A small badge pill."""
    w = Inches(1.4)
    h = Inches(0.3)
    r = add_rect(slide, left, top, w, h, fill)
    r.line.fill.background()
    txt = slide.shapes.add_textbox(left, top, w, h)
    tf  = txt.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = text_color or C_WHITE
    return r, txt


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ═══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# Full-bleed blue background
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, C_PRIMARY)

# Accent strip bottom-left
add_rect(s1, 0, Inches(6.8), Inches(5), Inches(0.7), C_ACCENT)

# Title
add_text(s1, "SchoolConnect Portal",
         Inches(0.7), Inches(1.4), Inches(11), Inches(1.5),
         font_name="Calibri", font_size=52, bold=True,
         color=C_WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text(s1, "A live, two-role school-parent communication platform with real-time\nmessaging, student performance updates, and an activity feed.",
         Inches(0.7), Inches(2.9), Inches(9), Inches(1.2),
         font_name="Calibri", font_size=20,
         color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.LEFT)

# URL strip
add_rect(s1, Inches(0.7), Inches(4.4), Inches(9.5), Inches(0.55),
         RGBColor(0x1E, 0x40, 0xAF))
add_text(s1, "  michaelwegter.com/demos/edtech-school-connect-web-portal/",
         Inches(0.7), Inches(4.4), Inches(9.5), Inches(0.55),
         font_name="Calibri", font_size=14,
         color=C_WHITE)

# Byline
add_text(s1, "Proposal by Michael Wegter  |  June 2026",
         Inches(0.7), Inches(5.2), Inches(8), Inches(0.4),
         font_name="Calibri", font_size=13,
         color=RGBColor(0x93, 0xC5, 0xFD))

# Role badges
add_rect(s1, Inches(0.7), Inches(5.8), Inches(2.2), Inches(0.45),
         RGBColor(0x16, 0x3A, 0x8A))
add_text(s1, "  Parent role (live demo)",
         Inches(0.7), Inches(5.8), Inches(2.2), Inches(0.45),
         font_name="Calibri", font_size=11, color=C_WHITE)

add_rect(s1, Inches(3.1), Inches(5.8), Inches(2.4), Inches(0.45),
         C_ACCENT)
add_text(s1, "  Teacher role (live demo)",
         Inches(3.1), Inches(5.8), Inches(2.4), Inches(0.45),
         font_name="Calibri", font_size=11, color=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, C_BG)
add_rect(s2, 0, 0, SLIDE_W, Inches(0.08), C_PRIMARY)

add_label(s2, "The Problem", Inches(0.7), Inches(0.35))
add_text(s2, "Parents are disconnected from what happens at school",
         Inches(0.7), Inches(0.75), Inches(11), Inches(0.95),
         font_name="Calibri", font_size=36, bold=True, color=C_TEXT)

# Three stat cards
cards_data = [
    ("30 min/day", "Average time teachers spend on parent communication via phone and email", C_PRIMARY),
    ("51M users", "ClassDojo's scale proves that real-time school-parent apps are in massive demand", C_ACCENT),
    ("1 platform", "Schools need one place to send updates, share grades, and message parents directly", C_WARNING),
]
for i, (num, label, col) in enumerate(cards_data):
    cx = Inches(0.65 + i * 4.2)
    cy = Inches(2.0)
    add_rect(s2, cx, cy, Inches(3.9), Inches(3.5), C_CARD)
    # top color bar
    add_rect(s2, cx, cy, Inches(3.9), Inches(0.09), col)
    add_text(s2, num,
             cx + Inches(0.2), cy + Inches(0.3), Inches(3.5), Inches(0.85),
             font_name="Calibri", font_size=38, bold=True, color=col)
    add_text(s2, label,
             cx + Inches(0.2), cy + Inches(1.2), Inches(3.5), Inches(2.0),
             font_name="Calibri", font_size=14, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution (Live Demo)
# ═══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, C_BG)
add_rect(s3, 0, 0, SLIDE_W, Inches(0.08), C_PRIMARY)

add_label(s3, "The Solution", Inches(0.7), Inches(0.35))
add_text(s3, "I built the portal. You can log in right now.",
         Inches(0.7), Inches(0.75), Inches(11), Inches(0.95),
         font_name="Calibri", font_size=36, bold=True, color=C_TEXT)

# URL box
add_rect(s3, Inches(0.7), Inches(1.75), Inches(10.5), Inches(0.65), C_PRIMARY)
add_text(s3, "  Live demo: michaelwegter.com/demos/edtech-school-connect-web-portal/",
         Inches(0.7), Inches(1.75), Inches(10.5), Inches(0.65),
         font_name="Calibri", font_size=15, bold=True, color=C_WHITE)

# Three feature columns
features = [
    ("💬", "Real-Time Messaging",
     "Threaded parent-teacher message bubbles. Backend-persisted, 4-second polling. A message sent on the teacher side appears on the parent side without a refresh.",
     C_PRIMARY, "Hero feature"),
    ("📊", "Student Performance Dashboard",
     "Six subject grade cards with letter grades, numeric scores, and a live Chart.js radar chart. Teacher updates a grade; parent sees it within seconds.",
     C_ACCENT, "Supporting"),
    ("📢", "Activity and Announcements Feed",
     "Teachers post announcements; parents see them immediately at the top of the feed. Category badges: Academic, Event, Attendance, Announcement.",
     C_WARNING, "Supporting"),
]

for i, (icon, title, body, col, badge) in enumerate(features):
    cx = Inches(0.65 + i * 4.2)
    cy = Inches(2.65)
    add_rect(s3, cx, cy, Inches(3.9), Inches(4.2), C_CARD)
    add_rect(s3, cx, cy, Inches(3.9), Inches(0.09), col)
    add_text(s3, icon,
             cx + Inches(0.2), cy + Inches(0.2), Inches(0.6), Inches(0.6),
             font_name="Calibri", font_size=28)
    add_text(s3, badge,
             cx + Inches(0.8), cy + Inches(0.25), Inches(3.0), Inches(0.4),
             font_name="Calibri", font_size=10, bold=True, color=col)
    add_text(s3, title,
             cx + Inches(0.2), cy + Inches(0.8), Inches(3.5), Inches(0.7),
             font_name="Calibri", font_size=16, bold=True, color=C_TEXT)
    add_text(s3, body,
             cx + Inches(0.2), cy + Inches(1.55), Inches(3.5), Inches(2.5),
             font_name="Calibri", font_size=13, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Two Roles
# ═══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, SLIDE_W, SLIDE_H, C_BG)
add_rect(s4, 0, 0, SLIDE_W, Inches(0.08), C_PRIMARY)

add_label(s4, "Demo Walkthrough", Inches(0.7), Inches(0.35))
add_text(s4, "Two logins. Two complete experiences.",
         Inches(0.7), Inches(0.75), Inches(11), Inches(0.9),
         font_name="Calibri", font_size=36, bold=True, color=C_TEXT)

# Teacher card (left)
add_rect(s4, Inches(0.65), Inches(1.85), Inches(5.75), Inches(5.0), C_CARD)
add_rect(s4, Inches(0.65), Inches(1.85), Inches(5.75), Inches(0.55), C_PRIMARY)
add_text(s4, "  Teacher Login",
         Inches(0.65), Inches(1.85), Inches(5.75), Inches(0.55),
         font_name="Calibri", font_size=14, bold=True, color=C_WHITE)
add_text(s4, "teacher@demo.edu / TeachDemo1",
         Inches(0.85), Inches(2.55), Inches(5.35), Inches(0.4),
         font_name="Calibri", font_size=12, color=C_MUTED)
teacher_steps = [
    "See Alex Johnson's class summary on the teacher dashboard",
    "Open Grades, click any subject, update the score, save",
    "Go to Activities, post a new announcement",
    "Open Messages, reply to a parent message",
    "Log out and switch to the parent view to see changes live",
]
for j, step in enumerate(teacher_steps):
    add_text(s4, f"{j+1}.  {step}",
             Inches(0.85), Inches(3.1 + j * 0.52), Inches(5.35), Inches(0.5),
             font_name="Calibri", font_size=13, color=C_TEXT)

# Parent card (right)
add_rect(s4, Inches(6.9), Inches(1.85), Inches(5.75), Inches(5.0), C_CARD)
add_rect(s4, Inches(6.9), Inches(1.85), Inches(5.75), Inches(0.55), C_ACCENT)
add_text(s4, "  Parent Login",
         Inches(6.9), Inches(1.85), Inches(5.75), Inches(0.55),
         font_name="Calibri", font_size=14, bold=True, color=C_WHITE)
add_text(s4, "parent@demo.edu / ParentDemo1",
         Inches(7.1), Inches(2.55), Inches(5.35), Inches(0.4),
         font_name="Calibri", font_size=12, color=C_MUTED)
parent_steps = [
    "Dashboard: GPA card, unread message badge, recent activity feed",
    "Grades view: see updated science grade from the teacher (live!)",
    "Activity feed: see the announcement posted by the teacher",
    "Messages: send a reply to Mrs. Rivera, see the thread grow",
    "Open a second tab as teacher to watch messages arrive in real time",
]
for j, step in enumerate(parent_steps):
    add_text(s4, f"{j+1}.  {step}",
             Inches(7.1), Inches(3.1 + j * 0.52), Inches(5.35), Inches(0.5),
             font_name="Calibri", font_size=13, color=C_TEXT)

add_text(s4, "Student: Alex Johnson, Grade 7  |  Teacher: Mrs. Rivera",
         Inches(0.7), Inches(7.1), Inches(10), Inches(0.35),
         font_name="Calibri", font_size=11, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Requirements Traceability
# ═══════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, SLIDE_W, SLIDE_H, C_BG)
add_rect(s5, 0, 0, SLIDE_W, Inches(0.08), C_PRIMARY)

add_label(s5, "Requirements Coverage", Inches(0.7), Inches(0.35))
add_text(s5, "Every item from your posting is covered",
         Inches(0.7), Inches(0.75), Inches(11), Inches(0.9),
         font_name="Calibri", font_size=36, bold=True, color=C_TEXT)

reqs = [
    ("Two-way school-parent communication",
     "Live messaging module. Both parent and teacher compose, send, and receive. Real database. 4-second polling."),
    ("Real-time student performance updates",
     "Grade cards update when a teacher saves. Radar chart re-renders. Parent sees changes within one polling cycle."),
    ("Real-time activity updates",
     "Teacher posts to the activity feed; the item appears at the top of the parent view on next poll."),
    ("Functional MVP prototype",
     "The demo is live right now at the link above. Two seeded logins, three working features, a real backend API."),
    ("PHP stack",
     "Demo uses Flask + JS for static hosting. Every layer maps 1:1 to PHP/Laravel. Migration: 1 to 2 weeks."),
    ("Fast delivery",
     "This prototype is the production spec. Clean modular architecture. MVP timeline: 4 to 6 weeks."),
]

row_h = Inches(0.73)
for i, (req, detail) in enumerate(reqs):
    cy = Inches(1.9) + i * row_h
    bg = C_CARD if i % 2 == 0 else RGBColor(0xF3, 0xF8, 0xFF)
    add_rect(s5, Inches(0.65), cy, Inches(11.85), row_h, bg)
    # check
    add_text(s5, "✓", Inches(0.75), cy + Inches(0.18), Inches(0.4), Inches(0.4),
             font_name="Calibri", font_size=18, bold=True, color=C_ACCENT)
    # requirement name
    add_text(s5, req,
             Inches(1.2), cy + Inches(0.08), Inches(4.0), Inches(0.35),
             font_name="Calibri", font_size=12, bold=True, color=C_TEXT)
    # detail
    add_text(s5, detail,
             Inches(5.4), cy + Inches(0.08), Inches(7.0), Inches(0.58),
             font_name="Calibri", font_size=12, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PHP Production Mapping
# ═══════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, SLIDE_W, SLIDE_H, C_BG)
add_rect(s6, 0, 0, SLIDE_W, Inches(0.08), C_PRIMARY)

add_label(s6, "PHP Production Path", Inches(0.7), Inches(0.35))
add_text(s6, "The demo is a PHP-ready specification",
         Inches(0.7), Inches(0.75), Inches(11), Inches(0.9),
         font_name="Calibri", font_size=36, bold=True, color=C_TEXT)

add_text(s6, "Your posting requires PHP. The demo runs on Flask + vanilla JS because it lives on a static host. Every layer translates directly to PHP/Laravel. The user flows, data models, and UI components are identical. Only the runtime changes.",
         Inches(0.7), Inches(1.75), Inches(11.5), Inches(0.75),
         font_name="Calibri", font_size=14, color=C_MUTED)

# Header row
add_rect(s6, Inches(0.65), Inches(2.65), Inches(5.85), Inches(0.5), C_PRIMARY)
add_rect(s6, Inches(6.6),  Inches(2.65), Inches(5.9),  Inches(0.5), C_PRIMARY_DK)
add_text(s6, "  Demo layer (prototype)",
         Inches(0.65), Inches(2.65), Inches(5.85), Inches(0.5),
         font_name="Calibri", font_size=13, bold=True, color=C_WHITE)
add_text(s6, "  PHP/Laravel production equivalent",
         Inches(6.6), Inches(2.65), Inches(5.9), Inches(0.5),
         font_name="Calibri", font_size=13, bold=True, color=C_WHITE)

php_rows = [
    ("Flask blueprint + routes", "Laravel Controller + Routes (web.php / api.php)"),
    ("SQLite tables (edtech_*)", "MySQL / PostgreSQL with Eloquent ORM"),
    ("PyJWT HS256 tokens",       "Laravel Sanctum or Passport"),
    ("Vanilla JS frontend",      "Vue.js, Livewire, or keep vanilla JS"),
    ("4-second client polling",  "Laravel Echo + Pusher WebSockets"),
    ("Static GitHub Pages host", "Laravel on shared hosting or a VPS"),
]

for i, (left, right) in enumerate(php_rows):
    cy = Inches(3.25) + i * Inches(0.6)
    bg_l = C_CARD if i % 2 == 0 else RGBColor(0xF3, 0xF8, 0xFF)
    bg_r = C_CARD if i % 2 == 0 else RGBColor(0xF0, 0xFD, 0xF4)
    add_rect(s6, Inches(0.65), cy, Inches(5.85), Inches(0.6), bg_l)
    add_rect(s6, Inches(6.6),  cy, Inches(5.9),  Inches(0.6), bg_r)
    add_text(s6, f"  {left}",
             Inches(0.65), cy + Inches(0.1), Inches(5.85), Inches(0.45),
             font_name="Calibri", font_size=13, color=C_MUTED)
    add_text(s6, f"  {right}",
             Inches(6.6), cy + Inches(0.1), Inches(5.9), Inches(0.45),
             font_name="Calibri", font_size=13, bold=True, color=C_TEXT)

add_text(s6, "Estimated migration effort from this prototype to PHP/Laravel: 1 to 2 weeks for a mid-level developer.",
         Inches(0.7), Inches(7.0), Inches(11), Inches(0.38),
         font_name="Calibri", font_size=12, color=C_PRIMARY, bold=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — About Michael
# ═══════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, SLIDE_W, SLIDE_H, C_BG)
add_rect(s7, 0, 0, SLIDE_W, Inches(0.08), C_PRIMARY)

add_label(s7, "About Michael Wegter", Inches(0.7), Inches(0.35))
add_text(s7, "Full-stack developer. Fast. Reliable. I own what I build.",
         Inches(0.7), Inches(0.75), Inches(11), Inches(0.9),
         font_name="Calibri", font_size=34, bold=True, color=C_TEXT)

proof = [
    ("2.5 years",
     "Sole developer and SME on a React + Python internal platform at U.S. Bank. Approximately 600 users per month, roughly 60,000 lines of code owned and maintained solo. If something broke, I could usually fix it within ten minutes."),
    ("6-month Azure migration",
     "Led the migration of my platform to Azure Cloud as the project's main representative, one of the first apps the bank moved. End-to-end ownership across architecture, implementation, and delivery."),
    ("150+ story points/sprint",
     "Currently at Optum (Angular + .NET + PostgreSQL), onion architecture, large team in strict Agile. My team's go-to for AI-assisted development. Contract won April 2026."),
    ("Full stack across the board",
     "React, Angular, Python (Flask), .NET/C#, Java, SQL, PostgreSQL, Docker, Kubernetes, CI/CD. Personal app portfolio live at michaelwegter.com with Spotify/Apple Music tools, gallery planners, and more."),
]

for i, (head, body) in enumerate(proof):
    cy = Inches(1.85) + i * Inches(1.3)
    add_rect(s7, Inches(0.65), cy, Inches(11.85), Inches(1.2), C_CARD)
    add_rect(s7, Inches(0.65), cy, Inches(0.12), Inches(1.2), C_PRIMARY)
    add_text(s7, head,
             Inches(0.9), cy + Inches(0.1), Inches(11.2), Inches(0.38),
             font_name="Calibri", font_size=14, bold=True, color=C_PRIMARY)
    add_text(s7, body,
             Inches(0.9), cy + Inches(0.5), Inches(11.2), Inches(0.65),
             font_name="Calibri", font_size=13, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CTA / Close
# ═══════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, SLIDE_W, SLIDE_H, C_PRIMARY)

# Accent bottom bar
add_rect(s8, 0, Inches(6.8), SLIDE_W, Inches(0.7), C_ACCENT)

add_text(s8, "The portal is live right now.",
         Inches(1.0), Inches(1.5), Inches(11), Inches(1.3),
         font_name="Calibri", font_size=52, bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(s8, "Log in as parent or teacher and click through the whole platform.\nNo install, no setup. Everything you see is real.",
         Inches(1.0), Inches(2.9), Inches(11), Inches(1.2),
         font_name="Calibri", font_size=20,
         color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

# URL box
add_rect(s8, Inches(2.0), Inches(4.3), Inches(9.0), Inches(0.75),
         RGBColor(0xFF, 0xFF, 0xFF))
add_text(s8, "michaelwegter.com/demos/edtech-school-connect-web-portal/",
         Inches(2.0), Inches(4.3), Inches(9.0), Inches(0.75),
         font_name="Calibri", font_size=17, bold=True,
         color=C_PRIMARY, align=PP_ALIGN.CENTER)

add_text(s8, "Proposal by Michael Wegter  |  June 2026  |  michaelwegter.com",
         Inches(0.7), Inches(6.8), Inches(11.5), Inches(0.5),
         font_name="Calibri", font_size=12,
         color=C_WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "proposal", "deck.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
