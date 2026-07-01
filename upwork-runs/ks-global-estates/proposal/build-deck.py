"""
Build deck.pptx for KS Global Estates proposal.
Run from: upwork-runs/ks-global-estates/
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, sys

# Colors
INK   = RGBColor(0x12, 0x10, 0x0E)
CANVAS= RGBColor(0xF5, 0xF1, 0xEB)
GOLD  = RGBColor(0xB8, 0x97, 0x5A)
MID   = RGBColor(0x6B, 0x65, 0x60)
DARK  = RGBColor(0x1C, 0x19, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

MEDIA = os.path.join(os.path.dirname(__file__), "media")

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, left, top, width, height, fill_rgb=None, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape

def txbox(slide, text, left, top, width, height,
          font_name="DM Sans", size=18, bold=False, italic=False,
          color=INK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def kicker(slide, text, left, top, width):
    tb = txbox(slide, text, left, top, width, Inches(0.35),
               font_name="DM Sans", size=9, bold=True, color=GOLD)
    tb.text_frame.paragraphs[0].runs[0].font.name = "DM Sans"

def add_img(slide, path, left, top, width, height=None):
    if not os.path.exists(path):
        return
    if height:
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        slide.shapes.add_picture(path, left, top, width)

def gold_bar(slide, top=Inches(7.2)):
    rect(slide, 0, top, W, Inches(0.05), GOLD)

# ── SLIDE 1: TITLE ──────────────────────────────────────────────────────────
s = add_slide()
hero = os.path.join(MEDIA, "hero.png")
rect(s, 0, 0, W, H, DARK)
if os.path.exists(hero):
    s.shapes.add_picture(hero, 0, 0, W, H)
overlay = s.shapes.add_shape(1, 0, 0, W, H)
overlay.fill.solid()
overlay.fill.fore_color.rgb = DARK
overlay.line.fill.background()

kicker(s, "WEB DESIGN & DEVELOPMENT PROPOSAL", Inches(0.8), Inches(1.8), Inches(8))
txbox(s, "KS Global Estates", Inches(0.8), Inches(2.15), Inches(9),
      Inches(1.6), font_name="Cormorant Garamond", size=54, bold=False,
      italic=False, color=CANVAS)
txbox(s, "Property Discovery Platform", Inches(0.8), Inches(3.5), Inches(8),
      Inches(0.8), font_name="Cormorant Garamond", size=30, italic=True, color=GOLD)
txbox(s, "Proposal by Michael Wegter", Inches(0.8), Inches(5.0), Inches(6),
      Inches(0.5), font_name="DM Sans", size=13, color=MID)
txbox(s, "michaelwegter.com  |  github.com/mwegter95", Inches(0.8), Inches(5.45),
      Inches(7), Inches(0.4), font_name="DM Sans", size=11, color=MID)
gold_bar(s)

# ── SLIDE 2: THE PROBLEM ────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, CANVAS)
rect(s, 0, 0, Inches(0.08), H, GOLD)
kicker(s, "THE OPPORTUNITY", Inches(0.6), Inches(0.6), Inches(5))
txbox(s, "Global real estate demands a site that earns trust in three seconds",
      Inches(0.6), Inches(1.0), Inches(8.5), Inches(1.5),
      font_name="Cormorant Garamond", size=38, color=INK)
txbox(s,
      "KS Global Estates operates across multiple markets and property types. "
      "A generic template site, or one that reads as AI-generated, undermines credibility "
      "with the high-net-worth buyers and investors who are the actual audience.\n\n"
      "What is needed is a purposefully designed, production-quality platform with real "
      "property browsing functionality: interactive map, live filters, rich listings, and "
      "a design that signals genuine craft.",
      Inches(0.6), Inches(2.6), Inches(8.2), Inches(3.0),
      font_name="DM Sans", size=14, color=INK)

bullets = [
    "No live website currently (ksglobalestates.com returns 403)",
    "High-net-worth buyers judge credibility from first impression",
    "AI-generated design patterns erode trust in the luxury segment",
    "Map-based property search is now a baseline user expectation",
]
for i, b in enumerate(bullets):
    txbox(s, f"  {b}",
          Inches(8.8), Inches(1.6 + i * 1.1), Inches(4.1), Inches(0.9),
          font_name="DM Sans", size=12, color=MID)
    rect(s, Inches(8.75), Inches(1.78 + i * 1.1), Inches(0.04), Inches(0.35), GOLD)

gold_bar(s)

# ── SLIDE 3: THE APPROACH ───────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, CANVAS)
rect(s, 0, 0, Inches(0.08), H, GOLD)
kicker(s, "THE APPROACH", Inches(0.6), Inches(0.5), Inches(5))
txbox(s, "How the real site gets built", Inches(0.6), Inches(0.85), Inches(9),
      Inches(0.9), font_name="Cormorant Garamond", size=36, color=INK)

phases = [
    ("01", "Design System", "Palette, typography, and spacing tokens locked first. Bespoke identity, no templates."),
    ("02", "Frontend App",  "React with Leaflet map, faceted filters, synchronized list/map, galleries, favorites."),
    ("03", "Backend + Data","Node or Python API, PostgreSQL for properties, CMS for agent management of listings."),
    ("04", "QA + Deploy",   "Cross-browser and mobile testing, performance audit, CI/CD pipeline to production."),
]
col_w = Inches(2.9)
for i, (num, title, desc) in enumerate(phases):
    x = Inches(0.6) + i * (col_w + Inches(0.22))
    rect(s, x, Inches(1.9), col_w, Inches(4.5), WHITE)
    rect(s, x, Inches(1.9), Inches(0.04), Inches(4.5), GOLD)
    txbox(s, num, x + Inches(0.18), Inches(2.0), col_w - Inches(0.2), Inches(0.9),
          font_name="Cormorant Garamond", size=34, color=GOLD)
    txbox(s, title, x + Inches(0.18), Inches(2.85), col_w - Inches(0.25), Inches(0.55),
          font_name="DM Sans", size=12, bold=True, color=INK)
    txbox(s, desc, x + Inches(0.18), Inches(3.4), col_w - Inches(0.25), Inches(2.5),
          font_name="DM Sans", size=11, color=MID)

gold_bar(s)

# ── SLIDE 4: THE LIVE DEMO ──────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, DARK)
kicker(s, "THE LIVE DEMO", Inches(0.7), Inches(0.45), Inches(6))
txbox(s, "Working today. Try it yourself.",
      Inches(0.7), Inches(0.8), Inches(8), Inches(0.9),
      font_name="Cormorant Garamond", size=36, color=CANVAS)
txbox(s, "michaelwegter.com/work-samples/ks-global-estates",
      Inches(0.7), Inches(1.55), Inches(9), Inches(0.45),
      font_name="DM Sans", size=13, color=GOLD)

if os.path.exists(hero):
    add_img(s, hero, Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.7))

rect(s, 0, Inches(6.8), W, Inches(0.4), DARK)
txbox(s, "25 properties across 14 global cities  |  Leaflet map with price pins  |  "
         "Faceted filters  |  Detail gallery  |  Favorites",
      Inches(0.7), Inches(6.82), Inches(11), Inches(0.38),
      font_name="DM Sans", size=10, color=MID)
gold_bar(s)

# ── SLIDE 5: DEMO FLOW ──────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, DARK)
kicker(s, "DEMO WALKTHROUGH", Inches(0.6), Inches(0.45), Inches(6))
txbox(s, "From filter to detail in three steps",
      Inches(0.6), Inches(0.8), Inches(9), Inches(0.7),
      font_name="Cormorant Garamond", size=30, color=CANVAS)

stills = ["step-1.png", "step-2.png", "step-3.png"]
captions = [
    "Filter by type: Villa selected, map and list sync instantly",
    "Detail panel opens with specs, price, and inquiry form",
    "Gallery carousel with thumbnail navigation",
]
still_w = Inches(3.9)
still_h = Inches(4.0)
for i, (fn, cap) in enumerate(zip(stills, captions)):
    fp = os.path.join(MEDIA, fn)
    x = Inches(0.55) + i * (still_w + Inches(0.2))
    if os.path.exists(fp):
        add_img(s, fp, x, Inches(1.65), still_w, still_h)
    txbox(s, cap, x, Inches(5.75), still_w, Inches(0.65),
          font_name="DM Sans", size=10, color=MID)

gold_bar(s)

# ── SLIDE 6: REQUIREMENTS MATRIX ────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, CANVAS)
rect(s, 0, 0, Inches(0.08), H, GOLD)
kicker(s, "REQUIREMENTS COVERAGE", Inches(0.6), Inches(0.45), Inches(7))
txbox(s, "Every requirement addressed",
      Inches(0.6), Inches(0.82), Inches(8), Inches(0.75),
      font_name="Cormorant Garamond", size=34, color=INK)

rows = [
    ("R1", "Professional frontend + backend",       "Full React frontend built; Node/Python+PG backend spec for prod"),
    ("R2", "Visually appealing, brand design",      "Cormorant Garamond + DM Sans, warm parchment/gold palette, editorial layout"),
    ("R3", "User-friendly UX",                      "Faceted filters, synchronized map/list, flyTo animation, mobile-responsive"),
    ("R4", "High-quality portfolio-grade work",     "10/10 feature richness, live deployed demo, no placeholder content"),
    ("R5", "Reasonable timeline and budget",        "Discuss scope and phasing once you have reviewed the demo"),
    ("R6", "Property listing and browsing",         "25 global listings, filterable, sortable, with map sync"),
    ("R7", "Mock property database",                "14 cities, 6 regions, 5 property types, full spec data per listing"),
    ("R8", "List view",                             "Horizontal card list, real-time filter updates, scroll-to-selected"),
    ("R9", "Map view with Leaflet",                 "Leaflet + CartoDB Voyager, custom gold price pins, flyTo on select"),
    ("R10","Property images only",                  "All images: exteriors, interiors, waterfront villas, penthouses"),
    ("R11","No AI-generated appearance",            "28 AI-tell patterns explicitly excluded; deliberate editorial design"),
]
col_x = [Inches(0.55), Inches(1.25), Inches(4.05)]
row_h = Inches(0.47)
y0 = Inches(1.7)

rect(s, Inches(0.55), y0 - Inches(0.08), Inches(12.5), Inches(0.38), DARK)
for j, hdr in enumerate(["ID", "Requirement", "How it is addressed"]):
    txbox(s, hdr, col_x[j], y0 - Inches(0.06), Inches(2.5), Inches(0.3),
          font_name="DM Sans", size=9, bold=True, color=CANVAS)

for i, (rid, req, how) in enumerate(rows):
    y = y0 + Inches(0.35) + i * row_h
    if i % 2 == 0:
        rect(s, Inches(0.55), y - Inches(0.06), Inches(12.5), row_h,
             RGBColor(0xEE,0xEA,0xE3))
    txbox(s, rid,  col_x[0], y, Inches(0.6),  row_h, font_name="DM Sans", size=9, bold=True, color=GOLD)
    txbox(s, "Yes", Inches(0.9), y, Inches(0.3), row_h, font_name="DM Sans", size=9, bold=True, color=GOLD)
    txbox(s, req,  col_x[1], y, Inches(2.7),  row_h, font_name="DM Sans", size=9, color=INK)
    txbox(s, how,  col_x[2], y, Inches(9.0),  row_h, font_name="DM Sans", size=9, color=MID)

gold_bar(s)

# ── SLIDE 7: WHY MICHAEL ────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, CANVAS)
rect(s, 0, 0, Inches(0.08), H, GOLD)
kicker(s, "WHY MICHAEL WEGTER", Inches(0.6), Inches(0.45), Inches(6))
txbox(s, "Three proof points that matter here",
      Inches(0.6), Inches(0.82), Inches(10), Inches(0.75),
      font_name="Cormorant Garamond", size=34, color=INK)

points = [
    ("Sole developer, 600 users/month",
     "At U.S. Bank (via Turnberry) I owned a React + Python platform of ~60,000 LOC. "
     "When something broke I fixed it, often within 10 minutes. I also led its full Azure Cloud migration "
     "as the project's main representative."),
    ("Full-stack across the real stack",
     "React, Angular, Python, .NET/C#, Node, PostgreSQL, SQL, Docker, CI/CD. Currently 150+ story points "
     "per sprint at Optum on a large Angular + .NET project. I write the backend, own the frontend, and ship."),
    ("I build demos, not decks",
     "Every proposal I write comes with a working app. This demo was built to production standard, "
     "deployed, and live before you read this. That is how I operate on real projects too."),
]
card_w = Inches(3.8)
for i, (title, body) in enumerate(points):
    x = Inches(0.6) + i * (card_w + Inches(0.25))
    rect(s, x, Inches(1.8), card_w, Inches(4.5), WHITE)
    rect(s, x, Inches(1.8), Inches(0.04), Inches(4.5), GOLD)
    txbox(s, title, x + Inches(0.2), Inches(2.0), card_w - Inches(0.3), Inches(0.7),
          font_name="Cormorant Garamond", size=20, color=INK)
    txbox(s, body, x + Inches(0.2), Inches(2.75), card_w - Inches(0.3), Inches(3.2),
          font_name="DM Sans", size=12, color=MID)

gold_bar(s)

# ── SLIDE 8: SCOPE, TIMELINE, NEXT STEP ─────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, DARK)
kicker(s, "SCOPE + NEXT STEP", Inches(0.7), Inches(0.45), Inches(6))
txbox(s, "Ready to build the real thing",
      Inches(0.7), Inches(0.82), Inches(9), Inches(0.85),
      font_name="Cormorant Garamond", size=38, color=CANVAS)

scope_items = [
    ("Included in phase 1",
     ["Bespoke design system (brand identity)",
      "React frontend (map, filters, listings, detail, gallery)",
      "Node/Python backend + PostgreSQL",
      "CMS or admin UI for managing listings",
      "Mobile-responsive across all views",
      "CI/CD pipeline + staged review"]),
    ("Phase 2 options",
     ["Agent / broker profiles",
      "Mortgage calculator",
      "Saved search and email alerts",
      "Multi-language support",
      "Virtual tour embeds",
      "Real authentication for agents"]),
]
for j, (heading, items) in enumerate(scope_items):
    x = Inches(0.7) + j * Inches(5.8)
    txbox(s, heading, x, Inches(2.0), Inches(5.5), Inches(0.5),
          font_name="DM Sans", size=11, bold=True, color=GOLD)
    for k, item in enumerate(items):
        txbox(s, f"  {item}", x, Inches(2.55) + k * Inches(0.55), Inches(5.5), Inches(0.5),
              font_name="DM Sans", size=12, color=RGBColor(0xC5,0xBF,0xB8))
        rect(s, x - Inches(0.08), Inches(2.73) + k * Inches(0.55),
             Inches(0.03), Inches(0.28), GOLD)

rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.9), RGBColor(0x26,0x22,0x1F))
txbox(s, "Explore the demo, then send me a message on Upwork. "
         "I am happy to discuss scope, timeline, and budget.",
      Inches(0.85), Inches(6.08), Inches(8), Inches(0.65),
      font_name="DM Sans", size=12, color=CANVAS)
txbox(s, "michaelwegter.com/work-samples/ks-global-estates",
      Inches(9.2), Inches(6.12), Inches(3.7), Inches(0.5),
      font_name="DM Sans", size=10, color=GOLD)

gold_bar(s)

# SAVE
out_path = os.path.join(os.path.dirname(__file__), "deck.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
