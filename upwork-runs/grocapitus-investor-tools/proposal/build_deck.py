#!/usr/bin/env python3
"""Build the Grocapitus proposal deck. On-brand: dark surfaces, mustard accent,
Space Grotesk / Inter / JetBrains Mono. No em dashes or en dashes anywhere."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

MEDIA = os.path.join(os.path.dirname(os.path.abspath(__file__)), "media")
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "deck.pptx")

# --- palette ---
MUSTARD   = RGBColor(0xE8, 0xB8, 0x20)
CYAN      = RGBColor(0x12, 0xB4, 0xC8)
GREEN     = RGBColor(0x6E, 0xD4, 0x6A)
PINK      = RGBColor(0xF0, 0x18, 0x6E)
BG_ROOT   = RGBColor(0x12, 0x11, 0x18)
BG_SURF   = RGBColor(0x19, 0x17, 0x20)
BG_CARD   = RGBColor(0x1E, 0x1C, 0x26)
BORDER    = RGBColor(0x2C, 0x2A, 0x38)
T_PRIMARY = RGBColor(0xF2, 0xED, 0xE4)
T_SECOND  = RGBColor(0x8A, 0x88, 0x98)
DARK_INK  = RGBColor(0x15, 0x13, 0x0A)

F_DISPLAY = "Space Grotesk"
F_BODY    = "Inter"
F_MONO    = "JetBrains Mono"

# 16:9 widescreen
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

DEMO_URL = "michaelwegter.com/work-samples/grocapitus-investor-tools"


def slide(bg=BG_ROOT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size, font=F_BODY, color=T_PRIMARY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, line_sp=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_sp
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.name = font; r.font.bold = bold
        r.font.color.rgb = color
        if spacing is not None:
            from pptx.oxml.ns import qn
            rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(spacing))
    return tb


def eyebrow(s, x, y, text, color=MUSTARD):
    txt(s, x, y, Inches(11), Inches(0.35), text.upper(), 12, F_MONO, color,
        bold=True, spacing=220)


def rect(s, x, y, w, h, fill, line=None, line_w=None):
    r = s.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is not None:
        r.line.color.rgb = line; r.line.width = line_w or Pt(1)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def accent_bar(s, x, y, h=Inches(0.06), w=Inches(0.9), color=MUSTARD):
    rect(s, x, y, w, h, color)


def picture_fit(s, path, x, y, w, h):
    """Place image centered/cropped into box w x h (cover)."""
    iw, ih = Image.open(path).size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # image wider: crop sides -> scale to height
        new_w = int(h * img_ratio)
        pic = s.shapes.add_picture(path, x - (new_w - w)//2, y, height=h)
        crop = (new_w - w) / new_w / 2
        pic.crop_left = crop; pic.crop_right = crop
        pic.left = x; pic.width = w
    else:
        new_h = int(w / img_ratio)
        pic = s.shapes.add_picture(path, x, y - (new_h - h)//2, width=w)
        crop = (new_h - h) / new_h / 2
        pic.crop_top = crop; pic.crop_bottom = crop
        pic.top = y; pic.height = h
    pic.line.color.rgb = BORDER; pic.line.width = Pt(1)
    return pic

M = Inches(0.7)  # left margin

# ============================================================ 1. TITLE
s = slide(BG_ROOT)
accent_bar(s, M, Inches(1.5))
eyebrow(s, M, Inches(1.75), "Proposal for Grocapitus Investments")
txt(s, M, Inches(2.35), Inches(11.5), Inches(2.0),
    "AI-Powered\nReal Estate Investor Tools", 54, F_DISPLAY, T_PRIMARY,
    bold=True, line_sp=1.0)
txt(s, M, Inches(4.7), Inches(11.0), Inches(1.2),
    "Brief prompt in. Working, tasteful investor tool out. Built by directing AI, "
    "shipped fast, documented so your team can learn it.", 18, F_BODY, T_SECOND,
    line_sp=1.2)
rect(s, M, Inches(6.15), Inches(5.2), Inches(0.55), BG_CARD, BORDER, Pt(1))
txt(s, M, Inches(6.15), Inches(5.2), Inches(0.55), "  LIVE DEMO   " + DEMO_URL,
    12, F_MONO, MUSTARD, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 2. THE PROBLEM
s = slide(BG_ROOT)
eyebrow(s, M, Inches(0.7), "The opportunity")
txt(s, M, Inches(1.15), Inches(11.8), Inches(1.4),
    "Your investors learn by doing. Give them tools that make the math obvious.",
    32, F_DISPLAY, T_PRIMARY, bold=True, line_sp=1.05)
quote = ("Grocapitus runs on investor education, and you want lightweight, "
         "AI-built apps that put real decisions in front of your community fast. "
         "The job is directing AI tools, shaping the output, and shipping useful "
         "apps on a tight weekly budget. Speed and taste both matter. A half-built "
         "MVP will not teach anyone, and a slow build will not keep up with your team.")
txt(s, M, Inches(2.7), Inches(11.5), Inches(1.4), quote, 18, F_BODY, T_SECOND,
    line_sp=1.3)
# three stakes cards
cards = [
    ("SPEED", "Prompt to working prototype in hours, not days."),
    ("TASTE", "Real usability, not toy demos. Numbers that read at a glance."),
    ("CLARITY", "Tools that help an investor actually decide."),
]
cw = Inches(3.7); gap = Inches(0.35); cx = M; cy = Inches(4.6)
for tag, body in cards:
    rect(s, cx, cy, cw, Inches(1.8), BG_CARD, BORDER, Pt(1))
    txt(s, cx + Inches(0.3), cy + Inches(0.3), cw - Inches(0.6), Inches(0.3),
        tag, 11, F_MONO, CYAN, bold=True, spacing=180)
    txt(s, cx + Inches(0.3), cy + Inches(0.75), cw - Inches(0.6), Inches(0.9),
        body, 15, F_BODY, T_PRIMARY, line_sp=1.2)
    cx += cw + gap

# ============================================================ 3. APPROACH
s = slide(BG_ROOT)
eyebrow(s, M, Inches(0.7), "How I would build the real thing")
txt(s, M, Inches(1.15), Inches(11.8), Inches(0.9),
    "Prompt to prototype, on repeat", 32, F_DISPLAY, T_PRIMARY, bold=True)
steps = [
    ("01", "Start from your brief", "Turn your prompt into a sharp spec: the one decision the tool should make easier for an investor."),
    ("02", "Direct the AI", "Vibe-code a clickable prototype with AI tools. No manual hand-coding, all output shaped through prompts."),
    ("03", "Ship early, iterate", "Get it in front of you fast, then tighten with your feedback instead of polishing in a vacuum."),
    ("04", "Document the workflow", "Capture prompts, iterations, and final logic in a short doc your team can pick up and extend."),
]
cw = Inches(2.83); gap = Inches(0.25); cx = M; cy = Inches(2.5)
for num, head, body in steps:
    rect(s, cx, cy, cw, Inches(3.4), BG_CARD, BORDER, Pt(1))
    accent_bar(s, cx + Inches(0.28), cy + Inches(0.3), w=Inches(0.5))
    txt(s, cx + Inches(0.28), cy + Inches(0.5), cw - Inches(0.5), Inches(0.4),
        num, 14, F_MONO, MUSTARD, bold=True)
    txt(s, cx + Inches(0.28), cy + Inches(1.0), cw - Inches(0.5), Inches(0.7),
        head, 17, F_DISPLAY, T_PRIMARY, bold=True, line_sp=1.05)
    txt(s, cx + Inches(0.28), cy + Inches(1.85), cw - Inches(0.5), Inches(1.4),
        body, 13, F_BODY, T_SECOND, line_sp=1.22)
    cx += cw + gap

# ============================================================ 4. THE LIVE DEMO
s = slide(BG_ROOT)
eyebrow(s, M, Inches(0.55), "The live demo, click it yourself")
txt(s, M, Inches(1.0), Inches(7.0), Inches(0.9),
    "Rental Cash Flow & Cap Rate Calculator", 26, F_DISPLAY, T_PRIMARY, bold=True,
    line_sp=1.0)
# big screenshot right
picture_fit(s, os.path.join(MEDIA, "hero.png"), Inches(6.5), Inches(1.05),
            Inches(6.3), Inches(3.94))
# left descriptor
bullets = [
    ("Live metrics engine", "NOI, cap rate, cash-on-cash, monthly cash flow, and total cash to close recompute on every keystroke."),
    ("Deal-verdict chip", "Color-coded verdict plus a revenue-versus-expense bar. Taste, not a toy."),
    ("Rent sensitivity slider", "Sweep the rent and watch a deal flip from positive to negative live."),
]
by = Inches(2.05)
for head, body in bullets:
    rect(s, M, by + Inches(0.05), Inches(0.06), Inches(0.85), MUSTARD)
    txt(s, M + Inches(0.25), by, Inches(5.6), Inches(0.35), head, 15, F_DISPLAY,
        T_PRIMARY, bold=True)
    txt(s, M + Inches(0.25), by + Inches(0.38), Inches(5.6), Inches(0.9), body,
        12.5, F_BODY, T_SECOND, line_sp=1.2)
    by += Inches(1.35)
# demo url band full width bottom
rect(s, M, Inches(6.45), Inches(11.93), Inches(0.6), BG_CARD, MUSTARD, Pt(1.25))
txt(s, M, Inches(6.45), Inches(11.93), Inches(0.6),
    "LIVE AND VIBE-CODED, NO MANUAL HAND-CODING    " + DEMO_URL, 13, F_MONO,
    MUSTARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 5. REQUIREMENT COVERAGE
s = slide(BG_ROOT)
eyebrow(s, M, Inches(0.55), "Requirement coverage")
txt(s, M, Inches(1.0), Inches(11.8), Inches(0.7),
    "Everything you asked for, mapped", 28, F_DISPLAY, T_PRIMARY, bold=True)
rows = [
    ("Vibe-coding tool experience", "DEMO + PORTFOLIO", "Four live AI-built apps plus this demo.", CYAN),
    ("Concept to working prototype", "DEMO", "The calculator is the proof.", CYAN),
    ("Lightweight apps under CEO direction", "DEMO", "Intentionally lightweight, built to spec.", CYAN),
    ("All code written by AI tools", "DEMO", "This and the portfolio were vibe-coded.", CYAN),
    ("Iterate quickly, ship early", "PORTFOLIO", "Multiple shipped apps are the track record.", GREEN),
    ("Document process so team can learn", "PROCESS", "Short prompt-and-iteration doc per app.", MUSTARD),
    ("Log hours via WebWork", "STATED", "I track all hours accurately in WebWork.", MUSTARD),
    ("Work independently, async comms", "STATED", "Self-directed, clear async updates.", MUSTARD),
    ("Currently enrolled (preferred)", "HONEST", "Not enrolled. A self-directed continuous learner.", PINK),
    ("Curiosity about real estate investing", "DEMO", "Real investor metrics, built first by choice.", CYAN),
    ("Reliable internet and workspace", "STATED", "Reliable connection, dedicated home workspace.", MUSTARD),
]
ty = Inches(1.85)
rh = Inches(0.44)
for i, (req, tag, cov, col) in enumerate(rows):
    if i % 2 == 0:
        rect(s, M, ty, Inches(11.93), rh, BG_SURF)
    txt(s, M + Inches(0.2), ty, Inches(4.0), rh, req, 12.5, F_BODY, T_PRIMARY,
        bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, M + Inches(4.3), ty, Inches(2.4), rh, tag, 10, F_MONO, col, bold=True,
        anchor=MSO_ANCHOR.MIDDLE, spacing=80)
    txt(s, M + Inches(6.8), ty, Inches(5.0), rh, cov, 11.5, F_BODY, T_SECOND,
        anchor=MSO_ANCHOR.MIDDLE)
    ty += rh

# ============================================================ 6. PORTFOLIO
s = slide(BG_ROOT)
eyebrow(s, M, Inches(0.55), "Proof of range")
txt(s, M, Inches(1.0), Inches(11.8), Inches(0.7),
    "Four AI-built apps, already live", 28, F_DISPLAY, T_PRIMARY, bold=True)
apps = [
    ("stage-hero.png", "Stage", "Gallery wall planner with drag-and-drop layout.", False),
    ("ssut-main.png", "Spotify Super User Tools", "Playlist extraction and building tools.", True),
    ("growyard-hero.png", "Growyard", "A garden planner that plots a plantable yard.", True),
    ("life-dashboard-hero.png", "Life Dashboard", "Personal life-tracking in one view.", True),
]
cw = Inches(2.83); gap = Inches(0.25); cx = M; cy = Inches(1.95)
for fname, name, desc, in_port in apps:
    rect(s, cx, cy, cw, Inches(4.4), BG_CARD, BORDER, Pt(1))
    path = os.path.join(MEDIA, "portfolio", fname)
    picture_fit(s, path, cx + Inches(0.12), cy + Inches(0.12), cw - Inches(0.24),
                Inches(1.85))
    txt(s, cx + Inches(0.25), cy + Inches(2.15), cw - Inches(0.5), Inches(0.3),
        "AI-BUILT", 9.5, F_MONO, MUSTARD, bold=True, spacing=160)
    txt(s, cx + Inches(0.25), cy + Inches(2.5), cw - Inches(0.5), Inches(0.7),
        name, 16, F_DISPLAY, T_PRIMARY, bold=True, line_sp=1.0)
    txt(s, cx + Inches(0.25), cy + Inches(3.4), cw - Inches(0.5), Inches(0.9),
        desc, 12, F_BODY, T_SECOND, line_sp=1.2)
    cx += cw + gap

# ============================================================ 7. WHY MICHAEL
s = slide(BG_ROOT)
eyebrow(s, M, Inches(0.7), "Why Michael")
txt(s, M, Inches(1.15), Inches(11.8), Inches(0.8),
    "The short version", 30, F_DISPLAY, T_PRIMARY, bold=True)
whys = [
    ("I ship with AI for real", "By day I am a full-stack developer and my team's go-to for AI-assisted delivery. Directing these tools well is my daily habit, not an experiment."),
    ("Fast and reliable", "I take ownership and move quickly. When something breaks I find it and fix it, often in minutes. You get prototypes early and tight feedback loops."),
    ("I like this problem", "Cap rate, cash-on-cash, and does this actually cash flow are the first questions every new investor asks. That is why I built the calculator first."),
]
cw = Inches(3.7); gap = Inches(0.35); cx = M; cy = Inches(2.6)
for head, body in whys:
    rect(s, cx, cy, cw, Inches(3.0), BG_CARD, BORDER, Pt(1))
    accent_bar(s, cx + Inches(0.3), cy + Inches(0.35), w=Inches(0.7))
    txt(s, cx + Inches(0.3), cy + Inches(0.6), cw - Inches(0.6), Inches(0.8),
        head, 18, F_DISPLAY, MUSTARD, bold=True, line_sp=1.05)
    txt(s, cx + Inches(0.3), cy + Inches(1.45), cw - Inches(0.6), Inches(1.4),
        body, 13, F_BODY, T_SECOND, line_sp=1.25)
    cx += cw + gap

# ============================================================ 8. SCOPE / CLOSE
s = slide(BG_ROOT)
eyebrow(s, M, Inches(0.85), "Scope, timeline, next step")
txt(s, M, Inches(1.35), Inches(11.5), Inches(1.0),
    "Hand me a real prompt", 38, F_DISPLAY, T_PRIMARY, bold=True)
scope = [
    ("ENGAGEMENT", "10 to 20 hours per week, fully remote, flexible schedule. Hours logged in WebWork."),
    ("RATE", "$30 per hour USD, as posted."),
    ("CADENCE", "A clickable prototype early in each cycle, then iterate from your feedback."),
    ("FIRST STEP", "Give me one real tool idea and I will turn it around so you can judge the speed and taste."),
]
cy = Inches(2.7)
for tag, body in scope:
    txt(s, M, cy, Inches(2.6), Inches(0.5), tag, 12, F_MONO, MUSTARD, bold=True,
        spacing=140, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, M + Inches(2.9), cy, Inches(8.8), Inches(0.7), body, 16, F_BODY,
        T_PRIMARY, line_sp=1.15, anchor=MSO_ANCHOR.MIDDLE)
    cy += Inches(0.82)
rect(s, M, Inches(6.4), Inches(11.93), Inches(0.62), BG_CARD, MUSTARD, Pt(1.25))
txt(s, M, Inches(6.4), Inches(11.93), Inches(0.62),
    "TRY THE LIVE DEMO    " + DEMO_URL, 14, F_MONO, MUSTARD,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("Saved", OUT, "slides:", len(prs.slides._sldIdLst))
