#!/usr/bin/env python3
"""Build the proposal deck for the ABA Services Website Upwork run."""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
MEDIA = os.path.join(HERE, "media")
OUT = os.path.join(HERE, "deck.pptx")

# Brand tokens
BG = RGBColor(0x12, 0x11, 0x18)          # --bg-root
CARD = RGBColor(0x1e, 0x1c, 0x26)        # --bg-card
SURFACE = RGBColor(0x19, 0x17, 0x20)     # --bg-surface
BORDER = RGBColor(0x2c, 0x2a, 0x38)
MUSTARD = RGBColor(0xe8, 0xb8, 0x20)
CYAN = RGBColor(0x12, 0xb4, 0xc8)
GREEN = RGBColor(0x6e, 0xd4, 0x6a)
TEXT = RGBColor(0xf2, 0xed, 0xe4)
TEXT_DIM = RGBColor(0x8a, 0x88, 0x98)
INK = RGBColor(0x1a, 0x18, 0x22)

FONT_TITLE = "Space Grotesk"
FONT_BODY = "Inter"
FONT_MONO = "JetBrains Mono"

DEMO_URL = "michaelwegter.com/work-samples/aba-services-website"
DEMO_URL_FULL = "https://michaelwegter.com/work-samples/aba-services-website"

# Hard rule check helper
def assert_no_dashes(s: str):
    if re.search(r"[–—]", s):
        raise AssertionError(f"Em/en dash found in: {s!r}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def add_bg(slide, color=BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    return rect

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=18, color=TEXT,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, mono=False,
             letter_spacing=None):
    assert_no_dashes(text)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT_MONO if mono else font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, fill=CARD, line=None, line_w=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_w is not None:
            rect.line.width = line_w
    rect.shadow.inherit = False
    return rect

def add_eyebrow(slide, x, y, w, text, color=MUSTARD):
    return add_text(slide, x, y, w, Inches(0.3), text.upper(),
                    mono=True, size=11, color=color, bold=True)

def add_footer(slide, page_num, total=8):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "MICHAEL WEGTER  /  PROPOSAL  /  ABA SERVICES WEBSITE",
             mono=True, size=9, color=TEXT_DIM)
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1), Inches(0.3),
             f"{page_num:02d} / {total:02d}",
             mono=True, size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def add_brand_mark(slide, x, y):
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.14), Inches(0.14))
    sq.fill.solid()
    sq.fill.fore_color.rgb = MUSTARD
    sq.line.fill.background()
    sq.shadow.inherit = False
    add_text(slide, x + Inches(0.22), y - Inches(0.06), Inches(3), Inches(0.3),
             "Michael Wegter", font=FONT_TITLE, size=14, color=TEXT, bold=True)

def add_accent_line(slide, x, y, w=Inches(0.8), color=MUSTARD):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(38100))  # ~3pt
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False

# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_brand_mark(s, Inches(0.5), Inches(0.5))
add_eyebrow(s, Inches(0.5), Inches(2.2), Inches(8), "PROPOSAL  /  UPWORK")
add_text(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(2),
         "ABA Services Website,\nDesign and Development",
         font=FONT_TITLE, size=54, bold=True, color=TEXT)
add_accent_line(s, Inches(0.5), Inches(4.7))
add_text(s, Inches(0.5), Inches(4.85), Inches(12), Inches(0.5),
         "A warm, trustworthy clinic site that families can scan in ten seconds.",
         font=FONT_TITLE, size=22, color=TEXT_DIM)
add_eyebrow(s, Inches(0.5), Inches(6.0), Inches(6), "LIVE DEMO", color=CYAN)
add_text(s, Inches(0.5), Inches(6.25), Inches(10), Inches(0.4),
         DEMO_URL_FULL, font=FONT_TITLE, size=16, color=MUSTARD, bold=True)
add_footer(s, 1)

# ---------- Slide 2: Demo at a glance ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_eyebrow(s, Inches(0.5), Inches(0.45), Inches(8), "THE DEMO AT A GLANCE")
add_text(s, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
         "One link. Real interactions. Built for this posting.",
         font=FONT_TITLE, size=32, bold=True, color=TEXT)

# Hero image
hero_path = os.path.join(MEDIA, "hero.png")
img_left = Inches(0.5)
img_top = Inches(1.85)
img_width = Inches(9.5)
add_rect(s, img_left - Emu(38100), img_top - Emu(38100),
         img_width + Emu(76200), Inches(4.7) + Emu(76200),
         fill=BORDER)
s.shapes.add_picture(hero_path, img_left, img_top, width=img_width, height=Inches(4.7))

# Right column callouts
col_x = Inches(10.4)
col_w = Inches(2.5)
add_eyebrow(s, col_x, Inches(1.95), col_w, "CLICK FIRST")
add_text(s, col_x, Inches(2.25), col_w, Inches(0.6),
         "Open a service card",
         font=FONT_TITLE, size=15, bold=True, color=TEXT)
add_text(s, col_x, Inches(2.7), col_w, Inches(0.8),
         "Each card expands to show age range and what to expect.",
         size=11, color=TEXT_DIM)

add_eyebrow(s, col_x, Inches(3.7), col_w, "THEN TRY")
add_text(s, col_x, Inches(4.0), col_w, Inches(0.6),
         "The intake form",
         font=FONT_TITLE, size=15, bold=True, color=TEXT)
add_text(s, col_x, Inches(4.45), col_w, Inches(0.8),
         "Inline validation, accessible labels, privacy note above submit.",
         size=11, color=TEXT_DIM)

add_eyebrow(s, col_x, Inches(5.5), col_w, "LIVE DEMO", color=CYAN)
add_text(s, col_x, Inches(5.8), col_w, Inches(0.4),
         DEMO_URL, mono=True, size=10, color=MUSTARD, bold=True)
add_footer(s, 2)

# ---------- Slide 3: Requirements I addressed ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_eyebrow(s, Inches(0.5), Inches(0.45), Inches(8), "REQUIREMENTS I ADDRESSED")
add_text(s, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
         "Every line in your posting maps to a screen.",
         font=FONT_TITLE, size=32, bold=True, color=TEXT)

rows = [
    ("R1", "Professional ABA website", "Entire single-page demo"),
    ("R2", "Visually appealing, communicates values", "Hero, services strip, warm palette"),
    ("R3", "User-friendly", "Anchor nav, accessible form, FAQ accordion"),
    ("R4", "Communicates services clearly", "Services explorer cards with age ranges"),
    ("R5", "Design and dev integrated", "One artifact, one developer, one timeline"),
    ("R6", "Graphic design elements", "Custom SVG logo, hero illustration, icons"),
]
table_top = Inches(2.0)
row_h = Inches(0.65)
left = Inches(0.5)
total_w = Inches(12.3)
col1_w = Inches(0.9)
col2_w = Inches(5.6)
col3_w = total_w - col1_w - col2_w

# Header
add_text(s, left, table_top, col1_w, Inches(0.35), "REQ",
         mono=True, size=10, color=MUSTARD, bold=True)
add_text(s, left + col1_w, table_top, col2_w, Inches(0.35),
         "WHAT YOU ASKED FOR", mono=True, size=10, color=MUSTARD, bold=True)
add_text(s, left + col1_w + col2_w, table_top, col3_w, Inches(0.35),
         "WHERE IT LIVES IN THE DEMO", mono=True, size=10, color=MUSTARD, bold=True)
add_rect(s, left, table_top + Inches(0.45), total_w, Emu(12700), fill=BORDER)

y = table_top + Inches(0.6)
for rid, ask, where in rows:
    add_text(s, left, y + Inches(0.15), col1_w, Inches(0.4),
             rid, mono=True, size=14, color=MUSTARD, bold=True)
    add_text(s, left + col1_w, y + Inches(0.15), col2_w, Inches(0.4),
             ask, size=14, color=TEXT)
    add_text(s, left + col1_w + col2_w, y + Inches(0.15), col3_w, Inches(0.4),
             where, size=14, color=TEXT_DIM)
    y += row_h
    add_rect(s, left, y, total_w, Emu(6350), fill=BORDER)

add_footer(s, 3)

# ---------- Slide 4: HIPAA-aware form ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_eyebrow(s, Inches(0.5), Inches(0.45), Inches(8), "HOW THE FORM RESPECTS HIPAA")
add_text(s, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
         "The marketing form is not a clinical intake.",
         font=FONT_TITLE, size=32, bold=True, color=TEXT)
add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.6),
         "Practical guidance, not legal advice. The pattern below is what most ABA clinics ship.",
         size=14, color=TEXT_DIM)

points = [
    ("Collect only routing info.", "Name, email, phone, and a service interest dropdown. No free-text medical history field."),
    ("Tell parents not to share PHI.", "A visible note above the submit button steers them to a secure channel for diagnosis details."),
    ("Route real intake through a HIPAA-eligible channel.", "EHR patient portal, or a form provider that signs a Business Associate Agreement."),
    ("Treat the marketing site as marketing.", "Fast, accessible, and low scope. The clinical work lives where it is regulated to live."),
]
y = Inches(2.5)
for head, body in points:
    add_text(s, Inches(0.6), y, Inches(0.45), Inches(0.5),
             "+", font=FONT_TITLE, size=24, color=MUSTARD, bold=True)
    add_text(s, Inches(1.05), y, Inches(11.5), Inches(0.4),
             head, font=FONT_TITLE, size=17, color=TEXT, bold=True)
    add_text(s, Inches(1.05), y + Inches(0.42), Inches(11.5), Inches(0.5),
             body, size=13, color=TEXT_DIM)
    y += Inches(1.05)
add_footer(s, 4)

# ---------- Slide 5: What I would build next ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_eyebrow(s, Inches(0.5), Inches(0.45), Inches(8), "POST UPWORK BUILD")
add_text(s, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
         "What I would ship for the real site.",
         font=FONT_TITLE, size=32, bold=True, color=TEXT)

# Embed page.png as a still on the left
still_path = os.path.join(MEDIA, "page.png")
img_left = Inches(0.5)
img_top = Inches(1.75)
img_w = Inches(5.6)
img_h = Inches(4.9)
add_rect(s, img_left - Emu(38100), img_top - Emu(38100),
         img_w + Emu(76200), img_h + Emu(76200), fill=BORDER)
s.shapes.add_picture(still_path, img_left, img_top, width=img_w, height=img_h)

# Right column: 4 bullets
bullets = [
    ("Five to seven page site", "Home, About, Services, Team, Insurance, Contact. Static build for speed."),
    ("Mobile-first, WCAG AA", "44 by 44 touch targets, keyboard nav, prefers-reduced-motion honored."),
    ("HIPAA-aware intake wiring", "Dropdown-based form posted to your EHR portal or a BAA-covered provider."),
    ("Optional headless CMS", "Sanity or Contentful if you want to edit copy yourselves after launch."),
]
col_x = Inches(6.5)
col_w = Inches(6.4)
y = Inches(1.85)
for head, body in bullets:
    add_text(s, col_x, y, Inches(0.4), Inches(0.4),
             "+", font=FONT_TITLE, size=22, color=MUSTARD, bold=True)
    add_text(s, col_x + Inches(0.45), y, col_w - Inches(0.5), Inches(0.4),
             head, font=FONT_TITLE, size=18, color=TEXT, bold=True)
    add_text(s, col_x + Inches(0.45), y + Inches(0.45), col_w - Inches(0.5), Inches(0.7),
             body, size=13, color=TEXT_DIM)
    y += Inches(1.2)
add_footer(s, 5)

# ---------- Slide 6: Frameworks + similar work ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_eyebrow(s, Inches(0.5), Inches(0.45), Inches(8), "FRAMEWORKS AND RECENT WORK")
add_text(s, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
         "Senior full-stack, ten years shipping real software.",
         font=FONT_TITLE, size=30, bold=True, color=TEXT)

# Two cards
card_w = Inches(6.0)
card_h = Inches(4.6)
gap = Inches(0.3)
left_x = Inches(0.5)
right_x = left_x + card_w + gap
card_y = Inches(2.0)

# Left card: frameworks
add_rect(s, left_x, card_y, card_w, card_h, fill=CARD)
add_eyebrow(s, left_x + Inches(0.4), card_y + Inches(0.4), card_w - Inches(0.8),
            "FRAMEWORKS I WORK IN")
add_text(s, left_x + Inches(0.4), card_y + Inches(0.7), card_w - Inches(0.8), Inches(0.6),
         "Daily delivery",
         font=FONT_TITLE, size=22, bold=True, color=TEXT)
fw_lines = [
    ("Front end", "React, Vite, Angular, vanilla JS, TypeScript"),
    ("Back end", "Flask, Python, .NET / C#, Java"),
    ("Data", "PostgreSQL, SQL Server, SQLite"),
    ("Ship it", "Docker, GitHub Actions, Azure, Playwright"),
]
y = card_y + Inches(1.5)
for label, value in fw_lines:
    add_text(s, left_x + Inches(0.4), y, Inches(1.6), Inches(0.4),
             label.upper(), mono=True, size=10, color=MUSTARD, bold=True)
    add_text(s, left_x + Inches(2.05), y, card_w - Inches(2.5), Inches(0.4),
             value, size=14, color=TEXT)
    y += Inches(0.62)
add_text(s, left_x + Inches(0.4), y + Inches(0.1), card_w - Inches(0.8), Inches(0.8),
         "This demo: plain HTML, CSS, and JS, shipped same-origin from my portfolio.",
         size=12, color=TEXT_DIM)

# Right card: recent work
add_rect(s, right_x, card_y, card_w, card_h, fill=CARD)
add_eyebrow(s, right_x + Inches(0.4), card_y + Inches(0.4), card_w - Inches(0.8),
            "RECENT WORK")
add_text(s, right_x + Inches(0.4), card_y + Inches(0.7), card_w - Inches(0.8), Inches(0.6),
         "Real systems, in production",
         font=FONT_TITLE, size=22, bold=True, color=TEXT)
work = [
    ("Optum RHRP 4", "Current contract. Angular plus .NET. 150+ story points delivered. Team's go-to for AI-assisted development."),
    ("U.S. Bank TDAAS", "2.5 years sole developer and SME. React plus Python plus SQL, 60k LOC, 600 users a month. Led six month Azure migration."),
    ("michaelwegter.com", "Personal apps and work samples gallery. Bright Path ABA mock was built for this posting."),
]
y = card_y + Inches(1.5)
for head, body in work:
    add_text(s, right_x + Inches(0.4), y, card_w - Inches(0.8), Inches(0.35),
             head, font=FONT_TITLE, size=15, bold=True, color=TEXT)
    add_text(s, right_x + Inches(0.4), y + Inches(0.35), card_w - Inches(0.8), Inches(0.7),
             body, size=12, color=TEXT_DIM)
    y += Inches(1.05)
add_footer(s, 6)

# ---------- Slide 7: Kickoff questions ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_eyebrow(s, Inches(0.5), Inches(0.45), Inches(8), "KICKOFF QUESTIONS")
add_text(s, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
         "Three things I would ask on a call.",
         font=FONT_TITLE, size=32, bold=True, color=TEXT)

qs = [
    ("01", "Page count and structure",
     "Single scrolling page or a 5 to 7 page site? My recommendation is the latter; the demo shows the single-page version for speed."),
    ("02", "CMS need",
     "Do you want to edit copy yourselves after launch? If yes, I would add a small headless CMS. If not, a static build is cheaper to maintain."),
    ("03", "Brand assets",
     "Do you have a logo, color direction, or photography? If not, I can propose a small visual direction before I build."),
]
y = Inches(1.95)
for num, head, body in qs:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.4), fill=CARD)
    add_text(s, Inches(0.85), y + Inches(0.25), Inches(1.1), Inches(0.9),
             num, font=FONT_TITLE, size=42, bold=True, color=MUSTARD)
    add_text(s, Inches(2.0), y + Inches(0.25), Inches(10.5), Inches(0.45),
             head, font=FONT_TITLE, size=20, bold=True, color=TEXT)
    add_text(s, Inches(2.0), y + Inches(0.75), Inches(10.5), Inches(0.55),
             body, size=13, color=TEXT_DIM)
    y += Inches(1.6)
add_footer(s, 7)

# ---------- Slide 8: Contact + demo ----------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_brand_mark(s, Inches(0.5), Inches(0.5))
add_eyebrow(s, Inches(0.5), Inches(2.4), Inches(8), "NEXT STEP")
add_text(s, Inches(0.5), Inches(2.75), Inches(12.3), Inches(1.6),
         "Open the demo,\nthen let's talk.",
         font=FONT_TITLE, size=58, bold=True, color=TEXT)
add_accent_line(s, Inches(0.5), Inches(5.05))

add_eyebrow(s, Inches(0.5), Inches(5.3), Inches(6), "LIVE DEMO", color=CYAN)
add_text(s, Inches(0.5), Inches(5.6), Inches(12), Inches(0.5),
         DEMO_URL_FULL, font=FONT_TITLE, size=20, color=MUSTARD, bold=True)

add_eyebrow(s, Inches(0.5), Inches(6.25), Inches(6), "PORTFOLIO", color=CYAN)
add_text(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.4),
         "michaelwegter.com  /  Michael Wegter",
         font=FONT_TITLE, size=16, color=TEXT)
add_footer(s, 8)

prs.save(OUT)
print(f"wrote {OUT}")
print(f"slides: {len(prs.slides)}")
