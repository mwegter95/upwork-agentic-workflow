#!/usr/bin/env python3
"""Build the TeleHealthy Voice Ops Console proposal deck."""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE   = os.path.dirname(os.path.abspath(__file__))
MEDIA  = os.path.join(HERE, "..", "media")
OUT    = os.path.join(HERE, "deck.pptx")

# ---- Design system (clinical teal, dark version for impact) ----
BG       = RGBColor(0x0B, 0x2A, 0x33)   # deep clinical teal
SURFACE  = RGBColor(0x0F, 0x3A, 0x46)
CARD     = RGBColor(0x13, 0x44, 0x53)
BORDER   = RGBColor(0x1C, 0x58, 0x6B)
PRIMARY  = RGBColor(0x0D, 0x94, 0x88)   # teal-600
CLINICAL = RGBColor(0x05, 0x96, 0x69)   # emerald-600
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)   # live call
RED      = RGBColor(0xEF, 0x44, 0x44)   # escalation
TEXT     = RGBColor(0xF0, 0xFD, 0xFA)   # teal-50
MUTED    = RGBColor(0x8E, 0xCD, 0xC8)
INK      = RGBColor(0x07, 0x1C, 0x22)

FONT_HEAD = "Inter"
FONT_BODY = "Inter"
FONT_MONO = "DM Mono"

DEMO_SHORT = "michaelwegter.com/work-samples/telehealthy"
DEMO_FULL  = "https://michaelwegter.com/work-samples/telehealthy"

def assert_no_dashes(s):
    if re.search(r"[–—]", s):
        raise AssertionError(f"Em/en dash found in: {s!r}")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width
H = prs.slide_height

def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r

def rect(slide, x, y, w, h, fill=CARD, line=None, lw=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        if lw: r.line.width = lw
    r.shadow.inherit = False
    return r

def txt(slide, x, y, w, h, text, *, font=FONT_BODY, size=16, color=TEXT,
        bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, mono=False):
    assert_no_dashes(text)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT_MONO if mono else font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def eyebrow(slide, x, y, w, label, color=PRIMARY):
    return txt(slide, x, y, w, Inches(0.3), label.upper(),
               mono=True, size=10, color=color, bold=True)

def accent_bar(slide, x, y, w=Inches(0.7), color=PRIMARY):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(38100))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False

def mark(slide, x, y):
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.13), Inches(0.13))
    sq.fill.solid(); sq.fill.fore_color.rgb = PRIMARY
    sq.line.fill.background(); sq.shadow.inherit = False
    txt(slide, x + Inches(0.21), y - Inches(0.05), Inches(3), Inches(0.3),
        "Michael Wegter", font=FONT_HEAD, size=14, bold=True)

def footer(slide, n, total=8):
    txt(slide, Inches(0.5), Inches(7.07), Inches(9), Inches(0.28),
        "MICHAEL WEGTER  /  PROPOSAL  /  TELEHEALTHY VOICE OPS CONSOLE",
        mono=True, size=8, color=MUTED)
    txt(slide, Inches(11.8), Inches(7.07), Inches(1.1), Inches(0.28),
        f"{n:02d} / {total:02d}", mono=True, size=8, color=MUTED, align=PP_ALIGN.RIGHT)

def meter_bar(slide, x, y, w, pct, fill_color):
    """Draw a single horizontal meter bar."""
    rect(slide, x, y, w, Emu(76200), fill=BORDER)   # track
    if pct > 0:
        rect(slide, x, y, int(w * pct), Emu(76200), fill=fill_color)

def img(slide, path, x, y, w, h=None):
    if not os.path.isfile(path):
        print(f"  [WARN] image not found: {path}")
        rect(slide, x, y, w, h or Inches(3), fill=BORDER)
        return
    if h:
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        slide.shapes.add_picture(path, x, y, width=w)

TOTAL = 8

# ==================== Slide 1: Title ====================
s = prs.slides.add_slide(BLANK)
bg(s)
mark(s, Inches(0.5), Inches(0.5))

eyebrow(s, Inches(0.5), Inches(2.0), Inches(9), "Proposal / Upwork / AI Voice Engineer")
txt(s, Inches(0.5), Inches(2.35), Inches(12.3), Inches(2.2),
    "TeleHealthy\nVoice Ops Console",
    font=FONT_HEAD, size=56, bold=True)
accent_bar(s, Inches(0.5), Inches(4.75))
txt(s, Inches(0.5), Inches(4.92), Inches(12), Inches(0.5),
    "AI receptionist + outbound dialer + analytics, running live in your browser.",
    font=FONT_HEAD, size=20, color=MUTED)

eyebrow(s, Inches(0.5), Inches(5.85), Inches(5), "Live Demo", color=AMBER)
txt(s, Inches(0.5), Inches(6.1), Inches(12), Inches(0.45),
    DEMO_FULL, font=FONT_HEAD, size=16, color=AMBER, bold=True)
footer(s, 1, TOTAL)

# ==================== Slide 2: Demo at a glance ====================
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.5), Inches(0.4), Inches(9), "The Demo at a Glance")
txt(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6),
    "One link. Every requirement. Click it right now.",
    font=FONT_HEAD, size=30, bold=True)

# Hero screenshot left
hero = os.path.join(MEDIA, "01-hero-landing.png")
rect(s, Inches(0.4), Inches(1.6), Inches(8.8) + Emu(76200), Inches(5.1) + Emu(76200), fill=BORDER)
img(s, hero, Inches(0.5), Inches(1.7), Inches(8.6), Inches(4.9))

# Right callouts
cx = Inches(9.6); cw = Inches(3.4)
eyebrow(s, cx, Inches(1.7), cw, "Click first", color=PRIMARY)
txt(s, cx, Inches(2.0), cw, Inches(0.45),
    "Simulate Inbound Call", font=FONT_HEAD, size=15, bold=True)
txt(s, cx, Inches(2.5), cw, Inches(0.75),
    "Pick an intent: Book, FAQ, or escalate. Watch the AI carry the conversation and fire an SMS on booking.", size=12, color=MUTED)

eyebrow(s, cx, Inches(3.6), cw, "Then try", color=PRIMARY)
txt(s, cx, Inches(3.9), cw, Inches(0.45),
    "Outbound + Analytics", font=FONT_HEAD, size=15, bold=True)
txt(s, cx, Inches(4.4), cw, Inches(0.75),
    "Pick a CRM contact for a dynamic call script. Analytics tab shows 155 calls, intent breakdown, and transcript log.", size=12, color=MUTED)

eyebrow(s, cx, Inches(5.55), cw, "Live Demo", color=AMBER)
txt(s, cx, Inches(5.8), cw, Inches(0.35),
    DEMO_SHORT, mono=True, size=10, color=AMBER, bold=True)
footer(s, 2, TOTAL)

# ==================== Slide 3: Requirements addressed ====================
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.5), Inches(0.4), Inches(9), "Requirements Addressed")
txt(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6),
    "Every line in the posting maps to a working screen.",
    font=FONT_HEAD, size=30, bold=True)

rows = [
    ("req-1", "AI receptionist: greeting, intent, booking, escalate", "Inbound tab: 6 intent paths, warm transfer at <0.45 confidence"),
    ("req-2", "Appointment booking with details captured",              "Multi-turn booking flow: new/existing patient, date, provider"),
    ("req-3", "FAQ handling: insurance, hours, telehealth",             "3 FAQ intents, EN/ES responses, natural phrasing"),
    ("req-4", "SMS confirmation on booking",                            "Real-time SMS log fires with patient name + appt details"),
    ("req-5", "Outbound dialer, dynamic CRM scripts, 2+ scenarios",    "Outbound tab: 8 contacts, reminder + follow-up scripts"),
    ("req-6", "CRM data, analytics, transcript logging",               "Analytics: call volume chart, intents, confidence, transcripts"),
]

table_top = Inches(1.65)
row_h = Inches(0.73)
left = Inches(0.5)
total_w = Inches(12.3)
c1 = Inches(1.1); c2 = Inches(4.8); c3 = total_w - c1 - c2

# Header row
rect(s, left, table_top, total_w, Inches(0.38), fill=SURFACE)
txt(s, left + Inches(0.1), table_top + Inches(0.06), c1, Inches(0.3),
    "REQ", mono=True, size=9, color=PRIMARY, bold=True)
txt(s, left + c1 + Inches(0.1), table_top + Inches(0.06), c2, Inches(0.3),
    "WHAT YOU ASKED FOR", mono=True, size=9, color=PRIMARY, bold=True)
txt(s, left + c1 + c2 + Inches(0.1), table_top + Inches(0.06), c3, Inches(0.3),
    "WHERE IT LIVES IN THE DEMO", mono=True, size=9, color=PRIMARY, bold=True)

y = table_top + Inches(0.38)
for i, (rid, ask, where) in enumerate(rows):
    fill_c = CARD if i % 2 == 0 else SURFACE
    rect(s, left, y, total_w, row_h, fill=fill_c)
    txt(s, left + Inches(0.1), y + Inches(0.18), c1, Inches(0.4),
        rid, mono=True, size=12, color=AMBER, bold=True)
    txt(s, left + c1 + Inches(0.1), y + Inches(0.18), c2 - Inches(0.2), Inches(0.45),
        ask, size=12, color=TEXT)
    txt(s, left + c1 + c2 + Inches(0.1), y + Inches(0.18), c3 - Inches(0.2), Inches(0.45),
        where, size=12, color=MUTED)
    y += row_h

footer(s, 3, TOTAL)

# ==================== Slide 4: AI Confidence + Warm Transfer ====================
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.5), Inches(0.4), Inches(9), "Hero Feature: AI Confidence + Warm Transfer")
txt(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6),
    "The system hands off before it makes a mistake.",
    font=FONT_HEAD, size=30, bold=True)

# Left: escalation screenshot
esc = os.path.join(MEDIA, "05-warm-transfer-escalation.png")
rect(s, Inches(0.4), Inches(1.6), Inches(5.5) + Emu(76200), Inches(4.9) + Emu(76200), fill=BORDER)
img(s, esc, Inches(0.5), Inches(1.7), Inches(5.4), Inches(4.7))

# Right: confidence meter viz
cx = Inches(6.5); cw = Inches(6.4)
eyebrow(s, cx, Inches(1.7), cw, "Live confidence scoring")
txt(s, cx, Inches(2.0), cw, Inches(0.5),
    "Each intent gets a real-time confidence score.", font=FONT_HEAD, size=16, bold=True)

meters = [
    ("Booking intent",   0.88, CLINICAL),
    ("Insurance FAQ",    0.82, CLINICAL),
    ("Telehealth FAQ",   0.61, AMBER),
    ("Urgent / escalate",0.34, RED),
]
my = Inches(2.7)
for label, pct, fill in meters:
    txt(s, cx, my, Inches(2.0), Inches(0.28), label, mono=True, size=10, color=MUTED)
    meter_bar(s, cx + Inches(2.15), my + Inches(0.06), Inches(3.5), pct, fill)
    txt(s, cx + Inches(5.75), my, Inches(0.6), Inches(0.28),
        f"{pct:.2f}", mono=True, size=10, color=TEXT)
    my += Inches(0.48)

txt(s, cx, my + Inches(0.1), cw, Inches(0.5),
    "Below 0.45: escalation badge fires, warm transfer message plays, human agent alert logged.",
    size=12, color=MUTED)

rect(s, cx, my + Inches(0.8), cw, Inches(0.7), fill=RED)
txt(s, cx + Inches(0.2), my + Inches(0.93), cw - Inches(0.4), Inches(0.45),
    "ESCALATE  /  confidence 0.34  /  transferring to care coordinator",
    mono=True, size=11, color=TEXT, bold=True)

footer(s, 4, TOTAL)

# ==================== Slide 5: Analytics + CRM ====================
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.5), Inches(0.4), Inches(9), "Analytics Dashboard + CRM")
txt(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6),
    "Every call logged. Every metric tracked.",
    font=FONT_HEAD, size=30, bold=True)

# Left: analytics screenshot
analytics = os.path.join(MEDIA, "06-analytics-dashboard.png")
rect(s, Inches(0.4), Inches(1.6), Inches(5.5) + Emu(76200), Inches(4.9) + Emu(76200), fill=BORDER)
img(s, analytics, Inches(0.5), Inches(1.7), Inches(5.4), Inches(4.7))

# Right: CRM screenshot + bullets
crm = os.path.join(MEDIA, "07-outbound-crm-contacts.png")
rect(s, Inches(6.4), Inches(1.6), Inches(6.4) + Emu(76200), Inches(2.2) + Emu(76200), fill=BORDER)
img(s, crm, Inches(6.5), Inches(1.7), Inches(6.3), Inches(2.0))

kpis = [
    ("155 calls", "logged in session analytics"),
    ("67%", "booking rate"),
    ("8%", "escalation rate"),
    ("2m 14s", "avg handle time"),
]
ky = Inches(4.1)
for val, label in kpis:
    txt(s, Inches(6.5), ky, Inches(1.5), Inches(0.35), val,
        font=FONT_HEAD, size=18, bold=True, color=PRIMARY)
    txt(s, Inches(8.1), ky + Inches(0.04), Inches(4.6), Inches(0.3),
        label, size=12, color=MUTED)
    ky += Inches(0.52)

footer(s, 5, TOTAL)

# ==================== Slide 6: Automation pipeline + real stack ====================
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.5), Inches(0.4), Inches(9), "Demo Pipeline + Real-World Architecture")
txt(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6),
    "The demo maps directly to a production stack.",
    font=FONT_HEAD, size=30, bold=True)

# Demo pipeline nodes
nodes = [
    ("Greeting", PRIMARY),
    ("Intent", PRIMARY),
    ("Booking", CLINICAL),
    ("Webhook", AMBER),
    ("SMS", AMBER),
    ("Escalate", RED),
]
nx = Inches(0.5)
ny = Inches(1.7)
node_w = Inches(1.8)
for i, (label, color) in enumerate(nodes):
    rect(s, nx, ny, node_w, Inches(0.65), fill=SURFACE, line=color, lw=Pt(1.5))
    txt(s, nx + Inches(0.1), ny + Inches(0.15), node_w - Inches(0.2), Inches(0.4),
        label, mono=True, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        txt(s, nx + node_w, ny + Inches(0.18), Inches(0.3), Inches(0.35),
            ">", size=18, color=MUTED, align=PP_ALIGN.CENTER)
        nx += node_w + Inches(0.3)
    nx += Inches(0)

txt(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.35),
    "Pipeline lights up step by step during a live call in the demo Inbound tab. Each node is a real webhook event in production.",
    size=12, color=MUTED)

# Real stack table
stacks = [
    ("Telephony",     "Vapi or Retell",          "Inbound call handling, confidence scoring, HIPAA recording, GoHighLevel native integration"),
    ("Outbound",      "Retell + Twilio",          "Automated dialer, dynamic scripts from CRM contact fields pulled at call time"),
    ("Orchestration", "n8n / Make.com",           "Webhook-driven: intent classifier, slot-fill, CRM update, SMS dispatch, escalation routing"),
    ("CRM",           "GoHighLevel API",          "Contact lookup, appointment sync, call activity log, pipeline stage updates"),
    ("Analytics",     "PostgreSQL + React",        "Transcript + confidence log, daily aggregation, admin dashboard same as demo Analytics tab"),
]

sy = Inches(3.2)
for layer, tool, desc in stacks:
    rect(s, Inches(0.5), sy, Inches(12.3), Inches(0.62), fill=CARD)
    txt(s, Inches(0.7), sy + Inches(0.13), Inches(2.0), Inches(0.4),
        layer.upper(), mono=True, size=10, color=PRIMARY, bold=True)
    txt(s, Inches(2.8), sy + Inches(0.13), Inches(2.4), Inches(0.4),
        tool, size=12, bold=True, color=AMBER)
    txt(s, Inches(5.4), sy + Inches(0.13), Inches(7.2), Inches(0.4),
        desc, size=11, color=MUTED)
    sy += Inches(0.72)

footer(s, 6, TOTAL)

# ==================== Slide 7: Why Me ====================
s = prs.slides.add_slide(BLANK)
bg(s)
eyebrow(s, Inches(0.5), Inches(0.4), Inches(9), "Why I Am the Right Fit")
txt(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6),
    "Shipped real software at scale. Builds real AI, not wrappers.",
    font=FONT_HEAD, size=28, bold=True)

proofs = [
    ("01", "U.S. Bank TDAAS (2.5 yrs, sole SME)",
     "React + Python + SQL. 600 users/month. ~60k LOC. Became SME within months. Fixed critical bugs within ten minutes of discovery. Led the six-month Azure Cloud migration as the team's primary representative."),
    ("02", "Optum RHRP 4 (current, Angular + .NET/C#)",
     "150+ story points delivered in strict Agile. Team's go-to for AI-assisted development. Contract started April 2026. Large-team, onion-architecture, production system."),
    ("03", "Real AI integration, not wrappers",
     "The WebGPU in-browser LLM in this demo runs actual Llama-3 inference. Azure OpenAI integrated in production at U.S. Bank. Vapi and Retell docs read. GoHighLevel API documented and understood."),
    ("04", "Full stack the project needs",
     "React, TypeScript, Python/Flask, .NET, PostgreSQL, Docker, GitHub Actions, Playwright. Healthcare domain context (HIPAA, clinical terminology, psychiatry practice workflows) present throughout this demo."),
]
py = Inches(1.65)
for num, head, body in proofs:
    rect(s, Inches(0.5), py, Inches(12.3), Inches(1.3), fill=CARD)
    txt(s, Inches(0.65), py + Inches(0.3), Inches(0.8), Inches(0.7),
        num, font=FONT_HEAD, size=36, bold=True, color=PRIMARY)
    txt(s, Inches(1.55), py + Inches(0.12), Inches(10.9), Inches(0.4),
        head, font=FONT_HEAD, size=15, bold=True)
    txt(s, Inches(1.55), py + Inches(0.55), Inches(10.9), Inches(0.65),
        body, size=12, color=MUTED)
    py += Inches(1.42)

footer(s, 7, TOTAL)

# ==================== Slide 8: Close ====================
s = prs.slides.add_slide(BLANK)
bg(s)
mark(s, Inches(0.5), Inches(0.5))
eyebrow(s, Inches(0.5), Inches(2.3), Inches(9), "Next Step")
txt(s, Inches(0.5), Inches(2.65), Inches(12.3), Inches(1.7),
    "Open the demo,\nthen let's talk.",
    font=FONT_HEAD, size=56, bold=True)
accent_bar(s, Inches(0.5), Inches(4.9))

eyebrow(s, Inches(0.5), Inches(5.15), Inches(6), "Live Demo", color=AMBER)
txt(s, Inches(0.5), Inches(5.42), Inches(12.3), Inches(0.5),
    DEMO_FULL, font=FONT_HEAD, size=18, color=AMBER, bold=True)

eyebrow(s, Inches(0.5), Inches(6.15), Inches(6), "Portfolio + Contact", color=PRIMARY)
txt(s, Inches(0.5), Inches(6.42), Inches(12.3), Inches(0.4),
    "michaelwegter.com  /  zweetztuph@gmail.com",
    font=FONT_HEAD, size=15, color=MUTED)

footer(s, 8, TOTAL)

# ---- Save ----
prs.save(OUT)
print(f"wrote {OUT}")
print(f"slides: {len(prs.slides)}")
