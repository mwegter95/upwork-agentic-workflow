"""Build proposal deck for Orschell Excavating e-commerce proposal."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Colors
BG       = RGBColor(0x1C, 0x1F, 0x24)
SURFACE  = RGBColor(0x25, 0x29, 0x2F)
ORANGE   = RGBColor(0xF7, 0x6B, 0x10)
YELLOW   = RGBColor(0xFF, 0xD6, 0x00)
WHITE    = RGBColor(0xE8, 0xEC, 0xF0)
MUTED    = RGBColor(0x9B, 0xA3, 0xAD)
BORDER   = RGBColor(0x35, 0x3A, 0x42)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank

def add_slide():
    s = prs.slides.add_slide(blank_layout)
    # dark background
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return s

def box(slide, left, top, width, height, fill_color=None, border_color=None, border_pt=0):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt2(border_pt)
    else:
        shape.line.fill.background()
    return shape

def label(slide, text, left, top, width, height,
          font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
          font_name='Calibri'):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or WHITE
    run.font.name = font_name
    return txb

def accent_bar(slide, top=Inches(0), height=Pt(6)):
    """Hazard stripe bar at given vertical position."""
    stripe = box(slide, 0, top, SLIDE_W, height, fill_color=ORANGE)

def divider(slide, top):
    shape = box(slide, Inches(0.6), top, SLIDE_W - Inches(1.2), Pt(1), fill_color=BORDER)

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
s1 = add_slide()
accent_bar(s1, 0, Pt(8))
accent_bar(s1, SLIDE_H - Pt(8), Pt(8))

# Left panel
box(s1, 0, 0, Inches(6.5), SLIDE_H, fill_color=SURFACE)

label(s1, 'PROPOSAL FOR', Inches(0.55), Inches(1.6), Inches(5.5), Inches(0.5),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s1, 'Full-Stack E-Commerce Backend', Inches(0.55), Inches(2.1), Inches(5.5), Inches(1.8),
      font_size=36, bold=True, color=WHITE, font_name='Calibri')
label(s1, 'Node.js + TypeScript + PostgreSQL', Inches(0.55), Inches(3.7), Inches(5.5), Inches(0.5),
      font_size=16, bold=False, color=ORANGE, font_name='Calibri')
label(s1, 'Auth / Catalog / Inventory / Orders / REST API / Admin CMS',
      Inches(0.55), Inches(4.15), Inches(5.5), Inches(0.5),
      font_size=11, bold=False, color=MUTED, font_name='Calibri')

divider(s1, Inches(4.8))
label(s1, 'michaelwegter.com', Inches(0.55), Inches(5.0), Inches(5.5), Inches(0.4),
      font_size=12, bold=False, color=MUTED, font_name='Calibri')
label(s1, 'Michael Wegter', Inches(0.55), Inches(5.4), Inches(5.5), Inches(0.5),
      font_size=14, bold=True, color=WHITE, font_name='Calibri')
label(s1, '$1,700  |  Fixed Price  |  2.5 to 3 weeks', Inches(0.55), Inches(5.85), Inches(5.5), Inches(0.4),
      font_size=11, bold=False, color=MUTED, font_name='Calibri')

# Right panel: demo callout
label(s1, 'LIVE DEMO', Inches(7.2), Inches(2.5), Inches(5.5), Inches(0.4),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s1, 'Click and try it now', Inches(7.2), Inches(2.9), Inches(5.5), Inches(0.4),
      font_size=18, bold=True, color=WHITE, font_name='Calibri')
label(s1, 'michaelwegter.com/demos/\norschell-excavating-e-commerce-full/',
      Inches(7.2), Inches(3.35), Inches(5.5), Inches(0.8),
      font_size=11, bold=False, color=ORANGE, font_name='Calibri')
label(s1, 'Admin: admin@orschellsupply.com / Admin1234!\nCustomer: demo@customer.com / Demo1234!',
      Inches(7.2), Inches(4.2), Inches(5.5), Inches(0.7),
      font_size=10, bold=False, color=MUTED, font_name='Calibri')

# ── SLIDE 2: THE PROBLEM ──────────────────────────────────────────────────────
s2 = add_slide()
accent_bar(s2, 0, Pt(6))

label(s2, 'YOUR PROJECT', Inches(0.6), Inches(0.55), Inches(6), Inches(0.4),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s2, 'Backend features you need built', Inches(0.6), Inches(0.9), Inches(12), Inches(0.55),
      font_size=28, bold=True, color=WHITE, font_name='Calibri')

reqs = [
    ('🔐', 'User authentication + accounts',   'JWT roles, register/login, order history, bcrypt passwords'),
    ('📦', 'Product catalog management',        'Admin CMS: CRUD, pricing, SKU, categories, featured flag'),
    ('📊', 'Inventory tracking',                'Per-SKU stock, LOW/OUT-OF-STOCK alerts, atomic decrements on order'),
    ('🛒', 'Order processing workflows',        'Cart persistence, checkout, 4-step status pipeline, customer history'),
    ('💳', 'Payment gateway integration',       'Stripe-ready checkout form (mocked + labeled in demo)'),
    ('🔌', 'REST API development',              '22 typed Express endpoints, input validation, error handling'),
    ('⚡', 'Performance optimization',          'Indexed queries, WAL mode, connection pooling (pg-pool in prod)'),
    ('🗄️', 'Database improvements',            'Normalized schema, FKs, CHECK constraints, Postgres-compatible'),
]

col_w = Inches(5.8)
for i, (icon, title, desc) in enumerate(reqs):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (col_w + Inches(0.4))
    top  = Inches(1.65) + row * Inches(1.2)
    box(s2, left, top, col_w, Inches(1.1), fill_color=SURFACE)
    label(s2, icon + '  ' + title, left + Inches(0.18), top + Inches(0.1),
          col_w - Inches(0.3), Inches(0.38), font_size=12, bold=True, color=WHITE, font_name='Calibri')
    label(s2, desc, left + Inches(0.18), top + Inches(0.48),
          col_w - Inches(0.3), Inches(0.5), font_size=10, bold=False, color=MUTED, font_name='Calibri')

label(s2, 'All 8 requirements traceable to working features in the live demo.',
      Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
      font_size=10, bold=False, color=MUTED, font_name='Calibri')

# ── SLIDE 3: THE DEMO ─────────────────────────────────────────────────────────
s3 = add_slide()
accent_bar(s3, 0, Pt(6))

label(s3, 'LIVE DEMO', Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s3, 'Built before submitting this proposal', Inches(0.6), Inches(0.9), Inches(12), Inches(0.55),
      font_size=28, bold=True, color=WHITE, font_name='Calibri')

# 4 screenshot placeholders (boxes with captions)
screens = [
    ('01-hero.png',           'Home / Hero'),
    ('02-shop.png',           'Shop / Product Catalog'),
    ('06-admin-dashboard.png','Admin Dashboard'),
    ('07-admin-products.png', 'Admin Product CMS'),
]
img_dir = os.path.dirname(os.path.abspath(__file__))
sw = Inches(5.9)
sh = Inches(3.5)
for i, (fname, caption) in enumerate(screens):
    col = i % 2
    row = i // 2
    left = Inches(0.4) + col * (sw + Inches(0.3))
    top  = Inches(1.65) + row * (sh + Inches(0.65))
    fpath = os.path.join(img_dir, fname)
    if os.path.exists(fpath):
        try:
            s3.shapes.add_picture(fpath, left, top, sw, sh)
        except Exception:
            box(s3, left, top, sw, sh, fill_color=SURFACE)
    else:
        box(s3, left, top, sw, sh, fill_color=SURFACE)
    label(s3, caption, left, top + sh + Inches(0.06), sw, Inches(0.35),
          font_size=10, bold=True, color=MUTED, font_name='Calibri')

label(s3, 'michaelwegter.com/demos/orschell-excavating-e-commerce-full/',
      Inches(0.6), Inches(7.1), Inches(12), Inches(0.35),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')

# ── SLIDE 4: TECH STACK ───────────────────────────────────────────────────────
s4 = add_slide()
accent_bar(s4, 0, Pt(6))

label(s4, 'TECH STACK', Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s4, 'Your exact stack. Zero ramp-up.', Inches(0.6), Inches(0.9), Inches(12), Inches(0.55),
      font_size=28, bold=True, color=WHITE, font_name='Calibri')

categories = [
    ('Backend (as specified)', [
        'Node.js 22  (LTS)',
        'TypeScript (strict)',
        'Express 4  (REST router)',
        'PostgreSQL schema (SQLite in demo; pg + pg-pool in prod)',
        'JWT authentication  (jsonwebtoken)',
        'bcrypt password hashing',
    ]),
    ('Frontend', [
        'Vite 5 + React 18 + TypeScript',
        'React Router 6  (SPA routing)',
        'Recharts  (admin analytics)',
        'Axios  (typed API client)',
    ]),
    ('Database design', [
        'Normalized: users, products, categories,',
        'inventory, cart_items, orders, order_items',
        'FKs, CHECK constraints, covering indexes',
        'WAL journal mode; Postgres-identical SQL',
    ]),
    ('Production path', [
        'Swap node:sqlite -> pg + pg-pool',
        'Schema unchanged (no SQLite-isms)',
        'Add Stripe SDK for payment gateway',
        'Rate limiting, helmet, CORS hardening',
    ]),
]

cw = Inches(5.8)
for i, (cat, items) in enumerate(categories):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (cw + Inches(0.4))
    top  = Inches(1.7) + row * Inches(2.6)
    box(s4, left, top, cw, Inches(2.45), fill_color=SURFACE)
    # orange left bar
    box(s4, left, top, Inches(0.06), Inches(2.45), fill_color=ORANGE)
    label(s4, cat, left + Inches(0.18), top + Inches(0.12),
          cw - Inches(0.3), Inches(0.38), font_size=12, bold=True, color=ORANGE, font_name='Calibri')
    txt = '\n'.join(items)
    label(s4, txt, left + Inches(0.18), top + Inches(0.55),
          cw - Inches(0.3), Inches(1.75), font_size=10, bold=False, color=MUTED, font_name='Calibri')

# ── SLIDE 5: ADMIN CMS ────────────────────────────────────────────────────────
s5 = add_slide()
accent_bar(s5, 0, Pt(6))

label(s5, 'ADMIN CMS', Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s5, 'Full catalog + inventory + order management', Inches(0.6), Inches(0.9), Inches(12), Inches(0.55),
      font_size=28, bold=True, color=WHITE, font_name='Calibri')

screens2 = [
    ('06-admin-dashboard.png','Dashboard: Revenue KPIs + 30-day chart'),
    ('07-admin-products.png', 'Products: CRUD  (create / edit / deactivate)'),
    ('08-admin-orders.png',   'Orders: Status pipeline (Pending -> Delivered)'),
    ('09-admin-inventory.png','Inventory: Stock levels + inline edit'),
]
sw2 = Inches(5.9)
sh2 = Inches(3.4)
for i, (fname, caption) in enumerate(screens2):
    col = i % 2
    row = i // 2
    left = Inches(0.4) + col * (sw2 + Inches(0.3))
    top  = Inches(1.65) + row * (sh2 + Inches(0.65))
    fpath = os.path.join(img_dir, fname)
    if os.path.exists(fpath):
        try:
            s5.shapes.add_picture(fpath, left, top, sw2, sh2)
        except Exception:
            box(s5, left, top, sw2, sh2, fill_color=SURFACE)
    else:
        box(s5, left, top, sw2, sh2, fill_color=SURFACE)
    label(s5, caption, left, top + sh2 + Inches(0.06), sw2, Inches(0.38),
          font_size=10, bold=False, color=MUTED, font_name='Calibri')

# ── SLIDE 6: TIMELINE + PRICE ─────────────────────────────────────────────────
s6 = add_slide()
accent_bar(s6, 0, Pt(6))

label(s6, 'DELIVERY PLAN', Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s6, '2.5 to 3 weeks  |  $1,700 fixed price', Inches(0.6), Inches(0.9), Inches(12), Inches(0.55),
      font_size=28, bold=True, color=WHITE, font_name='Calibri')

weeks = [
    ('Week 1', 'Audit + Foundation',
     'Review existing codebase. Set up Postgres, migrate schema, wire JWT auth, stub remaining endpoints. TypeScript strict mode throughout.'),
    ('Week 2', 'Core features',
     'Product catalog API, inventory decrement, cart persistence, full order pipeline. Admin CMS endpoints. Stripe integration if approved.'),
    ('Week 3', 'Performance + QA',
     'Query analysis, indexes, connection pooling, rate limiting, input validation. Integration tests. Frontend wiring. Documented handoff.'),
]

for i, (week, title, desc) in enumerate(weeks):
    top = Inches(1.8) + i * Inches(1.5)
    box(s6, Inches(0.5), top, Inches(1.1), Inches(1.2), fill_color=ORANGE)
    label(s6, week, Inches(0.55), top + Inches(0.35), Inches(1.0), Inches(0.5),
          font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Calibri')
    label(s6, title, Inches(1.85), top + Inches(0.08), Inches(10.8), Inches(0.4),
          font_size=14, bold=True, color=WHITE, font_name='Calibri')
    label(s6, desc, Inches(1.85), top + Inches(0.52), Inches(10.8), Inches(0.65),
          font_size=11, bold=False, color=MUTED, font_name='Calibri')

# Price highlight
box(s6, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.85), fill_color=SURFACE)
box(s6, Inches(0.5), Inches(6.35), Inches(0.06), Inches(0.85), fill_color=ORANGE)
label(s6, '$1,700 fixed  |  All 8 features  |  Node.js + TypeScript + PostgreSQL  |  Documented handoff',
      Inches(0.75), Inches(6.5), Inches(12), Inches(0.5),
      font_size=13, bold=True, color=WHITE, font_name='Calibri')

# ── SLIDE 7: WHY ME ───────────────────────────────────────────────────────────
s7 = add_slide()
accent_bar(s7, 0, Pt(6))

label(s7, 'WHY ME', Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s7, 'Shipped it before. Can own it fast.', Inches(0.6), Inches(0.9), Inches(12), Inches(0.55),
      font_size=28, bold=True, color=WHITE, font_name='Calibri')

proof = [
    ('U.S. Bank / TDAAS  (2.5 years)',
     'Became sole developer and SME on a React + Python + SQL internal platform within months. 60k+ LOC. Led the Azure Cloud migration. If something broke, I fixed it, usually within 10 minutes.'),
    ('Optum RHRP 4  (current)',
     'Angular + .NET/C# + PostgreSQL enterprise project. Onion architecture, strict Agile. 150+ story points delivered. Go-to for AI-assisted development. Contract started April 2026.'),
    ('TypeScript + Node.js daily',
     'TypeScript is my primary backend language. The 22-endpoint API here was written without scaffolding, strict mode throughout, fully typed request/response interfaces.'),
    ('Full-stack portfolio',
     'Multiple shipped apps at michaelwegter.com: gallery planner, SEO analyzer, Spotify/Apple Music tools. Code at github.com/mwegter95. Same quality as what you just saw in the demo.'),
]

pw = Inches(5.8)
for i, (title, desc) in enumerate(proof):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (pw + Inches(0.4))
    top  = Inches(1.75) + row * Inches(2.3)
    box(s7, left, top, pw, Inches(2.1), fill_color=SURFACE)
    box(s7, left, top, Inches(0.06), Inches(2.1), fill_color=ORANGE)
    label(s7, title, left + Inches(0.18), top + Inches(0.15),
          pw - Inches(0.3), Inches(0.4), font_size=13, bold=True, color=WHITE, font_name='Calibri')
    label(s7, desc, left + Inches(0.18), top + Inches(0.6),
          pw - Inches(0.3), Inches(1.35), font_size=10, bold=False, color=MUTED, font_name='Calibri')

# ── SLIDE 8: CLOSE ────────────────────────────────────────────────────────────
s8 = add_slide()
accent_bar(s8, 0, Pt(8))
accent_bar(s8, SLIDE_H - Pt(8), Pt(8))

label(s8, 'NEXT STEP', Inches(0.6), Inches(1.9), Inches(12), Inches(0.45),
      font_size=11, bold=True, color=ORANGE, font_name='Calibri')
label(s8, "If the demo looks close to what you're building, I'd love a short call to walk through your existing codebase and nail down a delivery plan.",
      Inches(0.6), Inches(2.35), Inches(12.1), Inches(1.2),
      font_size=22, bold=False, color=WHITE, font_name='Calibri')

label(s8, 'Michael Wegter', Inches(0.6), Inches(4.0), Inches(8), Inches(0.5),
      font_size=24, bold=True, color=WHITE, font_name='Calibri')
label(s8, 'michaelwegter.com  |  github.com/mwegter95', Inches(0.6), Inches(4.55), Inches(8), Inches(0.4),
      font_size=13, bold=False, color=MUTED, font_name='Calibri')

label(s8, 'Live demo:',
      Inches(0.6), Inches(5.3), Inches(3), Inches(0.4),
      font_size=13, bold=True, color=ORANGE, font_name='Calibri')
label(s8, 'michaelwegter.com/demos/orschell-excavating-e-commerce-full/',
      Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
      font_size=13, bold=False, color=ORANGE, font_name='Calibri')
label(s8, 'Admin: admin@orschellsupply.com / Admin1234!    Customer: demo@customer.com / Demo1234!',
      Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
      font_size=10, bold=False, color=MUTED, font_name='Calibri')

# Save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'deck.pptx')
prs.save(out)
print(f'Saved: {out}')
