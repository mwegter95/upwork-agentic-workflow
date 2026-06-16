"""
Generate the proposal deck for CivicQ Landing Page.
Design: bespoke civic-tech dark palette (near-black, teal, muted gold),
Newsreader (headlines) + Inter (body) per plan.md.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Color palette (from plan.md / demo design system)
C_BG        = RGBColor(0x0B, 0x0E, 0x11)  # near-black base
C_SURFACE   = RGBColor(0x13, 0x17, 0x1C)  # card surface
C_SURF_HI   = RGBColor(0x1A, 0x1F, 0x26)  # elevated / hover
C_BORDER    = RGBColor(0x26, 0x2C, 0x34)  # hairline
C_TEXT      = RGBColor(0xE7, 0xEA, 0xEE)  # primary text
C_MUTED     = RGBColor(0x9A, 0xA4, 0xB1)  # secondary text
C_TEAL      = RGBColor(0x3D, 0xBF, 0xA8)  # accent, "record" verified-data color
C_TEAL_BG   = RGBColor(0x10, 0x22, 0x1D)  # LIVE pill background
C_GOLD      = RGBColor(0xC9, 0xA2, 0x27)  # secondary accent, IN DEVELOPMENT
C_GOLD_BG   = RGBColor(0x22, 0x1C, 0x0E)  # IN DEVELOPMENT pill background
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_outline_rect(slide, left, top, width, height, line_color, fill_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.color.rgb = line_color
    shape.line.width = Pt(0.75)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def add_label(slide, text, left, top, width=Inches(6), color=None):
    txb = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = color or C_TEAL
    return txb


def pill(slide, text, left, top, fill, text_color):
    w = Inches(1.7)
    h = Inches(0.32)
    r = add_rect(slide, left, top, w, h, fill)
    r.line.fill.background()
    txt = slide.shapes.add_textbox(left, top, w, h)
    tf = txt.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = text_color
    return r, txt


def base_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_BG)
    return s


def section_header(s, label, title):
    add_rect(s, 0, 0, SLIDE_W, Inches(0.07), C_TEAL)
    add_label(s, label, Inches(0.7), Inches(0.35))
    add_text(s, title, Inches(0.7), Inches(0.75), Inches(11.5), Inches(0.95),
             font_name="Georgia", font_size=32, bold=True, color=C_TEXT)


# ============================================================
# SLIDE 1 - Title / Hero
# ============================================================
s1 = base_slide()
add_rect(s1, 0, Inches(6.85), SLIDE_W, Inches(0.65), C_TEAL)

add_text(s1, "CivicQ Landing Page",
         Inches(0.7), Inches(1.5), Inches(11.5), Inches(1.4),
         font_name="Georgia", font_size=48, bold=True, color=C_TEXT)

add_text(s1, "All 10 sections, the methodology page, and all three required data\nvisualizations, built to spec in a faithful civic-tech dark design system.",
         Inches(0.7), Inches(2.9), Inches(10.5), Inches(1.2),
         font_name="Calibri", font_size=18, color=C_MUTED)

add_rect(s1, Inches(0.7), Inches(4.35), Inches(9.0), Inches(0.55), C_SURFACE)
add_text(s1, "  michaelwegter.com/work-samples/civicq",
         Inches(0.7), Inches(4.35), Inches(9.0), Inches(0.55),
         font_name="Calibri", font_size=14, bold=True, color=C_TEAL)

add_text(s1, "Proposal by Michael Wegter  |  June 2026",
         Inches(0.7), Inches(5.1), Inches(8), Inches(0.4),
         font_name="Calibri", font_size=12, color=C_MUTED)

pill(s1, "LIVE DEMO", Inches(0.7), Inches(5.7), C_TEAL_BG, C_TEAL)
pill(s1, "3 CODED VISUALIZATIONS", Inches(2.6), Inches(5.7), C_GOLD_BG, C_GOLD)
pill(s1, "REAL BACKEND", Inches(5.5), Inches(5.7), C_TEAL_BG, C_TEAL)


# ============================================================
# SLIDE 2 - The 10 sections, in order
# ============================================================
s2 = base_slide()
section_header(s2, "What is built", "All 10 sections, in your exact order")

sections = [
    "Hero, video thumbnail + primary/secondary CTAs",
    "Problem statement, single-column text",
    "Feature card grid, 2x2, LIVE / IN DEV pills",
    "Scoring methodology, six-dimension chart",
    "Honest status disclosure, two-column",
    "Organization Hub, wireframe + email capture",
    "Fundraising module, $75K breakdown + CTA",
    "Partners / Press / Orgs, three-column",
    "Feedback banner + footer",
    "Sticky nav, anchor links + search/auth",
]
col_w = Inches(5.7)
row_h = Inches(0.62)
for i, label in enumerate(sections):
    col = i // 5
    row = i % 5
    cx = Inches(0.7) + col * Inches(6.0)
    cy = Inches(1.95) + row * row_h
    add_outline_rect(s2, cx, cy, col_w, Inches(0.55), C_BORDER, C_SURFACE)
    add_text(s2, f"{i+1:02d}", cx + Inches(0.15), cy + Inches(0.08), Inches(0.6), Inches(0.4),
              font_name="Georgia", font_size=16, bold=True, color=C_TEAL)
    add_text(s2, label, cx + Inches(0.75), cy + Inches(0.1), Inches(4.85), Inches(0.4),
              font_name="Calibri", font_size=12, color=C_TEXT)

add_text(s2, "Plus a dedicated /methodology page, same design system, reachable from the nav.",
         Inches(0.7), Inches(5.35), Inches(11.5), Inches(0.4),
         font_name="Calibri", font_size=13, color=C_MUTED)


# ============================================================
# SLIDE 3 - Three coded data visualizations
# ============================================================
s3 = base_slide()
section_header(s3, "The thing you asked about directly", "Three custom graphics, coded by hand")

viz = [
    ("Six-Dimension Scoring Chart",
     "Weighted horizontal bars, each dimension labeled with name and weight, animated draw-in, hidden data table fallback for screen readers."),
    ("$75,000 Funds Breakdown",
     "Segmented bar summing exactly to $75K. Category and dollar figure labeled directly on each segment, readable without relying on color."),
    ("Organization Hub Mockup",
     "Coded browser-chrome frame around a stylized dashboard, built to look like an intentional product shot, not a gray wireframe box."),
]
for i, (title, body) in enumerate(viz):
    cx = Inches(0.65 + i * 4.2)
    cy = Inches(2.0)
    add_outline_rect(s3, cx, cy, Inches(3.9), Inches(4.3), C_BORDER, C_SURFACE)
    add_rect(s3, cx, cy, Inches(3.9), Inches(0.07), C_TEAL)
    add_text(s3, title, cx + Inches(0.25), cy + Inches(0.3), Inches(3.4), Inches(0.7),
              font_name="Calibri", font_size=16, bold=True, color=C_TEXT)
    add_text(s3, body, cx + Inches(0.25), cy + Inches(1.1), Inches(3.4), Inches(2.9),
              font_name="Calibri", font_size=13, color=C_MUTED)
    add_text(s3, "Coded by Michael, no designer assets",
              cx + Inches(0.25), cy + Inches(3.75), Inches(3.4), Inches(0.4),
              font_name="Calibri", font_size=10, bold=True, color=C_GOLD)


# ============================================================
# SLIDE 4 - Email capture + backend
# ============================================================
s4 = base_slide()
section_header(s4, "Real backend, not a mock", "Email capture that actually persists")

add_outline_rect(s4, Inches(0.65), Inches(2.0), Inches(11.9), Inches(4.4), C_BORDER, C_SURFACE)
add_text(s4, "POST /api/civicq/subscribe",
         Inches(1.0), Inches(2.35), Inches(8), Inches(0.5),
         font_name="Consolas", font_size=18, bold=True, color=C_TEAL)
add_text(s4, "Validates the email, writes it to a real SQLite-backed table with a timestamp, and returns a clear success, duplicate, or invalid response.",
         Inches(1.0), Inches(2.9), Inches(10.5), Inches(0.7),
         font_name="Calibri", font_size=14, color=C_MUTED)

add_text(s4, "GET /api/civicq/csv-export",
         Inches(1.0), Inches(3.85), Inches(8), Inches(0.5),
         font_name="Consolas", font_size=18, bold=True, color=C_TEAL)
add_text(s4, "Streams every captured email as a downloadable CSV file, proving the \"Mailchimp/ConvertKit with CSV backup\" requirement with real persistence.",
         Inches(1.0), Inches(4.4), Inches(10.5), Inches(0.7),
         font_name="Calibri", font_size=14, color=C_MUTED)

add_text(s4, "Tested end to end on the live deployment: browser form submit, API 201 response, row visible in the CSV export.",
         Inches(1.0), Inches(5.4), Inches(10.5), Inches(0.7),
         font_name="Calibri", font_size=13, bold=True, color=C_GOLD)


# ============================================================
# SLIDE 5 - Requirements traceability
# ============================================================
s5 = base_slide()
section_header(s5, "Requirements covered", "Every item in your posting, addressed")

reqs = [
    ("React, dark-mode system matched precisely", "Faithful institutional dark stand-in, ready to swap for your real tokens"),
    ("Responsive 390px and 1440px", "Verified at both breakpoints, grids collapse cleanly on mobile"),
    ("Lighthouse 85+ on mobile", "Lean build, no chart library weight, self-hosted fonts"),
    ("Three custom data visualizations", "All three coded as SVG, live on the demo, see prior slide"),
    ("Email capture with CSV backup", "Real backend, persists and exports, see prior slide"),
    ("Analytics (Plausible/Fathom)", "Privacy-respecting script hook and event tracking wired"),
    ("WCAG 2.1 AA accessibility", "Semantic landmarks, focus states, contrast-checked palette, ARIA labels"),
    ("/methodology page", "Separate route, same system, full scoring explanation"),
]
row_h2 = Inches(0.56)
for i, (req, detail) in enumerate(reqs):
    cy = Inches(1.95) + i * row_h2
    bg = C_SURFACE if i % 2 == 0 else C_SURF_HI
    add_rect(s5, Inches(0.65), cy, Inches(11.9), row_h2, bg)
    add_text(s5, "✓", Inches(0.8), cy + Inches(0.1), Inches(0.4), Inches(0.4),
             font_name="Calibri", font_size=16, bold=True, color=C_TEAL)
    add_text(s5, req, Inches(1.25), cy + Inches(0.07), Inches(4.3), Inches(0.45),
             font_name="Calibri", font_size=12, bold=True, color=C_TEXT)
    add_text(s5, detail, Inches(5.7), cy + Inches(0.07), Inches(6.7), Inches(0.45),
             font_name="Calibri", font_size=12, color=C_MUTED)


# ============================================================
# SLIDE 6 - Timeline / Milestones
# ============================================================
s6 = base_slide()
section_header(s6, "Timeline", "Three milestones inside your 2 to 4 week window")

milestones = [
    ("Week 1", "40%", "Structure", "Design-system study against your real staging tokens, all 10 sections and sticky nav built to spec, responsive scaffold at 390px and 1440px."),
    ("Week 2", "40%", "Viz + Forms", "All three coded data visualizations, email capture wired to your real Mailchimp or ConvertKit setup with CSV backup, analytics hooks installed."),
    ("Weeks 3-4", "20%", "Polish", "Full WCAG 2.1 AA audit, Lighthouse tuning to 85+ mobile, cross-device QA, methodology page finalization, handoff."),
]
for i, (wk, pct, title, body) in enumerate(milestones):
    cx = Inches(0.65 + i * 4.2)
    cy = Inches(2.0)
    add_outline_rect(s6, cx, cy, Inches(3.9), Inches(4.3), C_BORDER, C_SURFACE)
    add_rect(s6, cx, cy, Inches(3.9), Inches(0.07), C_GOLD)
    add_text(s6, wk, cx + Inches(0.25), cy + Inches(0.25), Inches(2.5), Inches(0.4),
              font_name="Calibri", font_size=13, bold=True, color=C_MUTED)
    add_text(s6, pct, cx + Inches(2.9), cy + Inches(0.2), Inches(0.9), Inches(0.5),
              font_name="Georgia", font_size=20, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)
    add_text(s6, title, cx + Inches(0.25), cy + Inches(0.75), Inches(3.4), Inches(0.6),
              font_name="Georgia", font_size=20, bold=True, color=C_TEXT)
    add_text(s6, body, cx + Inches(0.25), cy + Inches(1.5), Inches(3.4), Inches(2.6),
              font_name="Calibri", font_size=13, color=C_MUTED)

add_text(s6, "The GoFundMe launch dependency is a known external blocker. The funds module and CTA will be built launch-ready on my end regardless of campaign-side timing.",
         Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.6),
         font_name="Calibri", font_size=12, color=C_MUTED)


# ============================================================
# SLIDE 7 - Pricing
# ============================================================
s7 = base_slide()
section_header(s7, "Fixed price bid", "$4,200, milestone payments welcome")

add_outline_rect(s7, Inches(0.65), Inches(2.1), Inches(11.9), Inches(2.0), C_BORDER, C_SURFACE)
add_text(s7, "$4,200", Inches(1.0), Inches(2.45), Inches(4), Inches(1.3),
         font_name="Georgia", font_size=56, bold=True, color=C_TEAL)
add_text(s7, "Within your $3,500 to $5,000 range. Milestone split: 40% on structure,\n40% on visualizations and forms, 20% on polish and handoff.",
         Inches(5.3), Inches(2.7), Inches(7.0), Inches(1.2),
         font_name="Calibri", font_size=15, color=C_MUTED)

add_text(s7, "Approach to the three graphics",
         Inches(0.7), Inches(4.6), Inches(8), Inches(0.5),
         font_name="Calibri", font_size=16, bold=True, color=C_TEXT)
add_text(s7, "I build all three myself in code. No designer assets needed. The live demo already shows all three working.",
         Inches(0.7), Inches(5.15), Inches(11), Inches(0.7),
         font_name="Calibri", font_size=14, color=C_MUTED)


# ============================================================
# SLIDE 8 - About Michael
# ============================================================
s8 = base_slide()
section_header(s8, "About Michael Wegter", "Spec fidelity, ownership, speed")

proof = [
    ("2.5 years", "Sole developer and SME on a React and Python platform at U.S. Bank, about 600 users a month. If something broke, I could usually fix it within 10 minutes."),
    ("6-month Azure migration", "Led that platform's move to Azure Cloud as the project's main representative, one of the first apps the bank moved."),
    ("150+ story points/sprint", "Currently at Optum (Angular + .NET + PostgreSQL), large team in strict Agile, team's go-to for AI-assisted development."),
    ("US-based, available now", "Full stack across React, Python (Flask), .NET, SQL, and PostgreSQL. The CivicQ demo runs on the same backend infrastructure as my own site."),
]
for i, (head, body) in enumerate(proof):
    cy = Inches(1.95) + i * Inches(1.18)
    add_rect(s8, Inches(0.65), cy, Inches(11.9), Inches(1.05), C_SURFACE)
    add_rect(s8, Inches(0.65), cy, Inches(0.1), Inches(1.05), C_TEAL)
    add_text(s8, head, Inches(0.9), cy + Inches(0.08), Inches(11.3), Inches(0.35),
              font_name="Calibri", font_size=14, bold=True, color=C_TEAL)
    add_text(s8, body, Inches(0.9), cy + Inches(0.45), Inches(11.3), Inches(0.55),
              font_name="Calibri", font_size=12, color=C_MUTED)


# ============================================================
# SLIDE 9 - CTA / Close
# ============================================================
s9 = base_slide()
add_rect(s9, 0, Inches(6.8), SLIDE_W, Inches(0.7), C_TEAL)

add_text(s9, "Your home page is already live.",
         Inches(1.0), Inches(1.6), Inches(11), Inches(1.3),
         font_name="Georgia", font_size=44, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

add_text(s9, "Open it, click through all 10 sections, try the email form.\nNothing here is a mockup.",
         Inches(1.0), Inches(2.95), Inches(11), Inches(1.1),
         font_name="Calibri", font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)

add_rect(s9, Inches(2.2), Inches(4.3), Inches(8.9), Inches(0.7), C_SURFACE)
add_text(s9, "michaelwegter.com/work-samples/civicq",
         Inches(2.2), Inches(4.3), Inches(8.9), Inches(0.7),
         font_name="Calibri", font_size=16, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

add_text(s9, "Proposal by Michael Wegter  |  June 2026  |  michaelwegter.com",
         Inches(0.7), Inches(6.95), Inches(11.5), Inches(0.45),
         font_name="Calibri", font_size=11, color=C_BG, align=PP_ALIGN.CENTER)

add_text(s9, "Not affiliated with any political party, campaign, or candidate.",
         Inches(0.7), Inches(5.4), Inches(11.5), Inches(0.4),
         font_name="Calibri", font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER)


out_path = os.path.join(os.path.dirname(__file__), "deck.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
