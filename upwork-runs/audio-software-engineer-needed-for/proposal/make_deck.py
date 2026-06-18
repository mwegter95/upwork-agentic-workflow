"""
make_deck.py -- Generates deck.pptx for the Worship Platform Audit proposal.
Run: python3 make_deck.py  (from this directory)
Requires: pip install python-pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette (worship stage-console dark) ──────────────────────────────
BG      = RGBColor(0x0C, 0x0C, 0x0F)
SURFACE = RGBColor(0x18, 0x18, 0x22)
BORDER  = RGBColor(0x2A, 0x2A, 0x36)
WHITE   = RGBColor(0xF0, 0xEF, 0xE9)
MUTED   = RGBColor(0x7C, 0x7A, 0x86)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
BLACK   = RGBColor(0x00, 0x00, 0x00)

# ── Slide dimensions (16:9 widescreen) ─────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

MEDIA = os.path.join(os.path.dirname(__file__), "../media")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank layout


# ── Helpers ─────────────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None):
    from pptx.util import Inches
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Inter"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    run.font.name      = font
    return box

def multiline(slide, lines, l, t, w, h, font="Inter"):
    """lines = list of (text, size, bold, color, align)"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, align) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if text is None:
            continue
        run = para.add_run()
        run.text = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color
        run.font.name      = font
    return box

def amber_stripe(slide, height=0.055):
    rect(slide, 0, 0, 13.33, height, AMBER)

def pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


# ── SLIDE 1: Title ───────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
set_bg(s1)
amber_stripe(s1)

# Left block
txt(s1, "WORSHIP PLATFORM AUDIT", 0.7, 0.9, 7, 0.4,
    10, bold=True, color=AMBER, font="JetBrains Mono")
multiline(s1, [
    ("Sanctuary", 52, True, WHITE, PP_ALIGN.LEFT),
], 0.7, 1.35, 8, 1.2)
txt(s1, "Worship Playback Engine", 0.72, 2.5, 8, 0.5,
    22, bold=False, color=MUTED)
rect(s1, 0.7, 3.15, 2.5, 0.045, AMBER)
multiline(s1, [
    ("Michael Wegter", 18, True, WHITE, PP_ALIGN.LEFT),
    ("Audio Software Engineer", 13, False, MUTED, PP_ALIGN.LEFT),
], 0.7, 3.35, 7, 0.9)

txt(s1, "Live demo:", 0.7, 4.5, 2, 0.3, 11, color=MUTED)
txt(s1, "michaelwegter.com/demos/audio-software-engineer-needed-for/",
    0.7, 4.78, 9, 0.35, 11, color=BLUE, font="JetBrains Mono")

# Right: screenshot
pic(s1, os.path.join(MEDIA, "02-playback-active.png"),
    7.4, 0.8, 5.5, 5.8)


# ── SLIDE 2: Nine Pain Points ────────────────────────────────────────────────
PAINS = [
    "Playback transport reliability",
    "Section navigation behavior",
    "Audio synchronization",
    "Click track timing",
    "Guide cue timing and triggering",
    "Stem routing and playback",
    "Mixer channel behavior",
    "Count-in functionality",
    "Audio engine stability",
]

s2 = prs.slides.add_slide(BLANK)
set_bg(s2)
amber_stripe(s2)
txt(s2, "NINE PROBLEM AREAS", 0.55, 0.22, 10, 0.35,
    10, bold=True, color=AMBER, font="JetBrains Mono")
txt(s2, "Every issue you listed, addressed in the live demo", 0.55, 0.65, 10, 0.5,
    24, bold=True, color=WHITE)

# 3x3 grid
cols = 3
cell_w = 3.9
cell_h = 1.1
margin_l = 0.55
margin_t = 1.35
gap_x = 0.18
gap_y = 0.14

for i, pain in enumerate(PAINS):
    col = i % cols
    row = i // cols
    lx = margin_l + col * (cell_w + gap_x)
    ty = margin_t + row * (cell_h + gap_y)
    rect(s2, lx, ty, cell_w, cell_h, SURFACE, BORDER)
    # green check
    txt(s2, "✓", lx + 0.15, ty + 0.2, 0.35, 0.6, 16, bold=True, color=GREEN)
    txt(s2, pain, lx + 0.5, ty + 0.22, cell_w - 0.6, 0.7, 12, color=WHITE)

# Footer note
txt(s2, "Each implemented with sample-accurate Web Audio API scheduling.",
    0.55, 7.1, 9, 0.3, 10, color=MUTED, font="JetBrains Mono")


# ── SLIDE 3: Live Demo ────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
set_bg(s3)
amber_stripe(s3)
txt(s3, "LIVE DEMO", 0.55, 0.22, 8, 0.35,
    10, bold=True, color=AMBER, font="JetBrains Mono")
txt(s3, "Sanctuary Worship Playback Engine", 0.55, 0.65, 9, 0.5,
    24, bold=True, color=WHITE)

pic(s3, os.path.join(MEDIA, "01-initial-state.png"), 0.55, 1.3, 12.2, 4.9)

txt(s3, "Transport  |  Stems  |  Section Nav  |  Click Track  |  Count-In  |  A/B Loop  |  Full Mixer",
    0.55, 6.35, 9.5, 0.4, 10, color=MUTED, font="JetBrains Mono")
txt(s3, "michaelwegter.com/demos/audio-software-engineer-needed-for/",
    0.55, 6.85, 10, 0.35, 12, bold=True, color=BLUE, font="JetBrains Mono")


# ── SLIDE 4: Demo Features (two screenshots) ─────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
set_bg(s4)
amber_stripe(s4)
txt(s4, "DEMO FEATURES", 0.55, 0.22, 8, 0.35,
    10, bold=True, color=AMBER, font="JetBrains Mono")
txt(s4, "What the demo implements", 0.55, 0.65, 9, 0.5,
    24, bold=True, color=WHITE)

pic(s4, os.path.join(MEDIA, "04-section-nav.png"), 0.55, 1.3, 6.0, 4.1)
pic(s4, os.path.join(MEDIA, "05-mixer-stems.png"), 6.85, 1.3, 6.0, 4.1)

txt(s4, "Section navigation with BPM-locked seek",
    0.55, 5.5, 5.8, 0.35, 10, color=MUTED, align=PP_ALIGN.CENTER)
txt(s4, "6-stem mixer: drums, bass, keys, acoustic, vocals, click",
    6.85, 5.5, 5.8, 0.35, 10, color=MUTED, align=PP_ALIGN.CENTER)

txt(s4, "Wilson-pattern lookahead scheduler. OfflineAudioContext stem synthesis. Sample-accurate transport.",
    0.55, 6.85, 12.2, 0.35, 10, color=MUTED, font="JetBrains Mono")


# ── SLIDE 5: My Background (Music) ──────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
set_bg(s5)
amber_stripe(s5)
txt(s5, "MUSIC BACKGROUND", 0.55, 0.22, 8, 0.35,
    10, bold=True, color=AMBER, font="JetBrains Mono")
txt(s5, "Musician, DAW operator, and live audio professional", 0.55, 0.65, 11, 0.5,
    24, bold=True, color=WHITE)

items = [
    ("Bachelor of Music, St. Olaf College",
     "Formal music theory, ear training, composition, and audio production."),
    ("DAWs since the fifth grade",
     "Fifteen-plus years with Ableton Live, Logic Pro, and Pro Tools."),
    ("Professional Wedding DJ",
     "Live events require audio that never fails. I am obsessed with getting it exactly right."),
    ("Built Python MIDI Bridge for Live Gigs",
     "Real-time MIDI signal processing, running in production at every performance. Zero failures."),
]

for i, (title, body) in enumerate(items):
    ty = 1.45 + i * 1.3
    rect(s5, 0.55, ty, 5.9, 1.1, SURFACE, BORDER)
    txt(s5, title, 0.75, ty + 0.1, 5.5, 0.4, 13, bold=True, color=WHITE)
    txt(s5, body,  0.75, ty + 0.55, 5.5, 0.5, 11, color=MUTED)

pic(s5, os.path.join(MEDIA, "03-waveform-progress.png"), 7.0, 1.35, 5.9, 5.2)


# ── SLIDE 6: My Background (Engineering) ────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
set_bg(s6)
amber_stripe(s6)
txt(s6, "ENGINEERING BACKGROUND", 0.55, 0.22, 9, 0.35,
    10, bold=True, color=AMBER, font="JetBrains Mono")
txt(s6, "Full-stack engineer, real-time audio systems", 0.55, 0.65, 10, 0.5,
    24, bold=True, color=WHITE)

eng_items = [
    ("Optum / RHRP 4 (Current)",
     "Angular, .NET/C#, PostgreSQL, onion architecture. 150+ story points. Team AI lead."),
    ("U.S. Bank via Turnberry (2.5 years)",
     "React, Python, SQL, Java. Became sole dev and SME on 60,000-line platform in 2 months."),
    ("Python MIDI Bridge",
     "rtmidi, custom signal processing, real-time mapping. Production-deployed for live gigs."),
    ("Sanctuary Demo (this proposal)",
     "Web Audio API, OfflineAudioContext, GainNode, StereoPannerNode, Wilson lookahead scheduler."),
]

for i, (title, body) in enumerate(eng_items):
    col = i % 2
    row = i // 2
    lx = 0.55 + col * 6.45
    ty = 1.45 + row * 2.3
    rect(s6, lx, ty, 5.9, 2.05, SURFACE, BORDER)
    txt(s6, title, lx + 0.2, ty + 0.18, 5.5, 0.4, 13, bold=True, color=AMBER)
    txt(s6, body,  lx + 0.2, ty + 0.68, 5.5, 1.1, 12, color=WHITE)


# ── SLIDE 7: Audit Approach ──────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
set_bg(s7)
amber_stripe(s7)
txt(s7, "AUDIT APPROACH", 0.55, 0.22, 8, 0.35,
    10, bold=True, color=AMBER, font="JetBrains Mono")
txt(s7, "Diagnose first. Fix what is confirmed broken.", 0.55, 0.65, 11, 0.5,
    24, bold=True, color=WHITE)

steps = [
    ("01", "Review the codebase",
     "Read the full audio architecture: engine, scheduler, transport, and routing. Map the signal flow before touching anything."),
    ("02", "Reproduce and root-cause",
     "Systematically reproduce each of your nine issues. Trace each one to its specific root cause. Document before writing fixes."),
    ("03", "Fix high-priority bugs",
     "Targeted, surgical patches for transport reliability, sync drift, and timing issues. Each fix tested in isolation."),
    ("04", "Deliver recommendations",
     "Written analysis of root causes, fixes applied, and architectural improvements for the platform's next stage."),
]

sw = 2.9
sh = 4.4
gap = 0.22
ml = 0.55
mt = 1.4

for i, (num, title, body) in enumerate(steps):
    lx = ml + i * (sw + gap)
    rect(s7, lx, mt, sw, sh, SURFACE, BORDER)
    txt(s7, num,   lx + 0.2, mt + 0.2,  sw - 0.3, 0.5, 28, bold=True, color=AMBER)
    txt(s7, title, lx + 0.2, mt + 0.9,  sw - 0.3, 0.55, 14, bold=True, color=WHITE)
    txt(s7, body,  lx + 0.2, mt + 1.6,  sw - 0.3, 2.5, 11, color=MUTED)


# ── SLIDE 8: Next Steps / CTA ────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
set_bg(s8)
amber_stripe(s8)

txt(s8, "NEXT STEPS", 0.55, 0.22, 8, 0.35,
    10, bold=True, color=AMBER, font="JetBrains Mono")
txt(s8, "Let's talk about your codebase.", 0.55, 0.68, 11, 0.6,
    30, bold=True, color=WHITE)

txt(s8,
    "30 minutes to walk through the nine issues and look at the relevant parts of the architecture. "
    "I will give you an honest read on what I see before any engagement starts.",
    0.55, 1.5, 8.5, 1.0, 15, color=MUTED)

links = [
    ("Live demo:",    "michaelwegter.com/demos/audio-software-engineer-needed-for/", BLUE),
    ("Portfolio:",    "michaelwegter.com/work-samples",                               BLUE),
]
for i, (label, url, color) in enumerate(links):
    ty = 2.9 + i * 0.55
    txt(s8, label, 0.55, ty, 1.6, 0.45, 12, bold=True, color=MUTED)
    txt(s8, url,   2.2,  ty, 7.0, 0.45, 12, color=color, font="JetBrains Mono")

rect(s8, 0.55, 3.9, 0.04, 2.2, AMBER)

multiline(s8, [
    ("Michael Wegter", 22, True,  WHITE, PP_ALIGN.LEFT),
    ("Audio Software Engineer  |  Musician  |  Full-Stack Engineer", 12, False, MUTED, PP_ALIGN.LEFT),
], 0.75, 4.0, 7, 1.4)

# Right side: final screenshot
pic(s8, os.path.join(MEDIA, "02-playback-active.png"), 7.5, 1.2, 5.5, 5.2)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "deck.pptx")
prs.save(out)
print(f"Saved {out} ({len(prs.slides)} slides)")
