"""
Build the proposal deck for maryland-driveway-restore.
Design system matches the MDR demo:
  charcoal:  #1A1D1A
  green:     #2F5D3A
  green-acc: #4A9060
  white:     #FFFFFF
  off-white: #F5F5F3
  muted:     #7A847A
Fonts: Calibri (system) as fallback for Playfair Display / Inter (PPTX can't embed web fonts).
Run: python build-deck.py  (requires: pip install python-pptx)
Output: maryland-driveway-restore-proposal.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ──────────────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

CHARCOAL   = rgb('#1A1D1A')
CHARCOAL_L = rgb('#252925')
CARD       = rgb('#2A2E2A')
GREEN      = rgb('#2F5D3A')
GREEN_ACC  = rgb('#4A9060')
WHITE      = rgb('#FFFFFF')
OFF_WHITE  = rgb('#F5F5F3')
MUTED      = rgb('#7A847A')
TEXT_DIM   = rgb('#A8B4A8')
BORDER     = rgb('#3A3F3A')

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def fill_bg(slide, color=None):
    from pptx.oxml.ns import qn
    from lxml import etree
    c = color or CHARCOAL
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = c

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or WHITE
    run.font.name = font_name
    return txBox

def slide_num(prs):
    return len(prs.slides)

# ── Build ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]  # blank layout

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)

# green accent strip left
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)

# eyebrow
add_text(s, 'PHASE 1 WEBSITE PROPOSAL', Inches(0.55), Inches(1.6), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)

# headline
add_text(s, 'Maryland Driveway Restore', Inches(0.55), Inches(2.1), Inches(9), Inches(1),
         font_size=40, bold=True, color=WHITE)

add_text(s, 'A premium, mobile-first WordPress site that converts visitors into phone calls and estimate requests. Live in 7 days.',
         Inches(0.55), Inches(3.2), Inches(8.5), Inches(1), font_size=18, color=TEXT_DIM)

# demo box
add_rect(s, Inches(0.55), Inches(4.4), Inches(7.5), Inches(0.75), CARD, GREEN)
add_text(s, 'LIVE DEMO:  https://api.michaelwegter.com/demos/maryland-driveway-restore/',
         Inches(0.7), Inches(4.5), Inches(7.3), Inches(0.5), font_size=13, color=GREEN_ACC)

# by line
add_text(s, 'Proposal by Michael Wegter  |  michaelwegter.com',
         Inches(0.55), Inches(5.5), Inches(7), Inches(0.4), font_size=12, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, 'THE OPPORTUNITY', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, 'Most Maryland driveway contractors look like it is 2008 online.', Inches(0.55), Inches(1.1), Inches(9), Inches(1),
         font_size=30, bold=True, color=WHITE)

problems = [
    ('No mobile optimization', 'Over 60% of local searches happen on phones. Sites that break on mobile lose those leads instantly.'),
    ('No lead capture', 'A phone number in the footer is not a conversion strategy. You need a form, sticky CTAs, and GA4 tracking every action.'),
    ('No local SEO foundation', 'Without schema, optimized titles, and a GSC-connected sitemap, Google cannot rank a site for "Maryland driveway restoration."'),
    ('No trust signals', 'Before/after photography, reviews, and a clear process walk visitors through the decision. Most sites skip all three.'),
]

y = Inches(2.2)
for title, detail in problems:
    add_rect(s, Inches(0.55), y, Inches(0.05), Inches(0.55), GREEN_ACC)
    add_text(s, title, Inches(0.75), y, Inches(5), Inches(0.35), font_size=13, bold=True, color=WHITE)
    add_text(s, detail, Inches(0.75), y + Inches(0.32), Inches(8.5), Inches(0.35), font_size=11, color=TEXT_DIM)
    y += Inches(0.85)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Demo
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, 'LIVE DEMO', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, 'A real WordPress site. Log in and explore the admin.', Inches(0.55), Inches(1.1), Inches(9), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

add_text(s, 'This is not a mockup. It is a fully functional WordPress install with all 9 Phase 1 pages, real schema, real forms, real Rank Math configuration, and real Fluent Forms lead capture. The same setup gets migrated to Rocket.net on Day 1.',
         Inches(0.55), Inches(2.0), Inches(8.5), Inches(0.9), font_size=13, color=TEXT_DIM)

add_rect(s, Inches(0.55), Inches(3.1), Inches(9.5), Inches(1.0), CARD, GREEN)
add_text(s, 'Demo URL', Inches(0.75), Inches(3.15), Inches(3), Inches(0.3), font_size=10, color=MUTED)
add_text(s, 'https://api.michaelwegter.com/demos/maryland-driveway-restore/',
         Inches(0.75), Inches(3.45), Inches(9), Inches(0.4), font_size=13, bold=True, color=GREEN_ACC)

add_rect(s, Inches(0.55), Inches(4.25), Inches(4.5), Inches(0.9), CARD)
add_text(s, 'WP Admin Login', Inches(0.75), Inches(4.3), Inches(4), Inches(0.3), font_size=10, color=MUTED)
add_text(s, 'demo_admin  /  DemoView2026!', Inches(0.75), Inches(4.6), Inches(4), Inches(0.35), font_size=13, color=WHITE)

add_rect(s, Inches(5.25), Inches(4.25), Inches(4.75), Inches(0.9), CARD)
add_text(s, 'WP Admin URL', Inches(5.45), Inches(4.3), Inches(4), Inches(0.3), font_size=10, color=MUTED)
add_text(s, '.../wp-login.php', Inches(5.45), Inches(4.6), Inches(4.5), Inches(0.35), font_size=12, color=GREEN_ACC)

add_text(s, 'Real WordPress. Real dashboard. Real demo.', Inches(0.55), Inches(5.45), Inches(8), Inches(0.4),
         font_size=14, bold=True, color=GREEN_ACC, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Phase 1 Scope
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, 'PHASE 1 SCOPE', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, '9 pages. Every requirement in the brief.', Inches(0.55), Inches(1.0), Inches(9), Inches(0.7),
         font_size=26, bold=True, color=WHITE)

pages = [
    'Homepage (hero, CTAs, benefits, process, pricing, gallery, reviews, FAQ, form)',
    'Driveway Restoration (service detail, process, before/after)',
    'Pricing (tiered table, add-ons, CTA)',
    'Gallery (39 images, lightbox, organized by type)',
    'About (mission, team, environmental commitment)',
    'FAQ (38 questions, accordion, FAQ schema)',
    'Contact / Free Estimate (photo upload, Fluent Forms)',
    'Service Areas (Maryland coverage, Annapolis focus)',
    '/free-estimate (dedicated QR landing page, not the homepage)',
]

y = Inches(1.9)
col1 = pages[:5]
col2 = pages[5:]

for item in col1:
    add_rect(s, Inches(0.55), y + Inches(0.08), Inches(0.12), Inches(0.12), GREEN_ACC)
    add_text(s, item, Inches(0.82), y, Inches(5.5), Inches(0.45), font_size=11, color=TEXT_DIM)
    y += Inches(0.52)

y = Inches(1.9)
for item in col2:
    add_rect(s, Inches(7.0), y + Inches(0.08), Inches(0.12), Inches(0.12), GREEN_ACC)
    add_text(s, item, Inches(7.27), y, Inches(5.5), Inches(0.45), font_size=11, color=TEXT_DIM)
    y += Inches(0.52)

add_text(s, 'ALSO INCLUDED: Rank Math SEO (4 schema types) | Fluent Forms | ShortPixel | Wordfence | UpdraftPlus | GA4 | GSC | Google Business Profile',
         Inches(0.55), Inches(6.8), Inches(12), Inches(0.45), font_size=10, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Mobile-First + SEO
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, 'MOBILE-FIRST + LOCAL SEO', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, 'Built to rank. Built to convert on phones.', Inches(0.55), Inches(1.0), Inches(9), Inches(0.7),
         font_size=26, bold=True, color=WHITE)

mobile_items = [
    ('Sticky "Call Now" + "Text a Photo"', 'Fixed at bottom of screen on mobile, 56px touch targets, visible at every scroll position'),
    ('Responsive mobile-first CSS', 'Designed for phones first, expanded for tablet/desktop. Hamburger nav on mobile.'),
    ('Core Web Vitals optimized', 'ShortPixel + Rocket CDN + lazy loading = Lighthouse 90+ target before launch'),
    ('Fast-loading design', 'Minimal plugin footprint, no bloat, images compressed and named for SEO'),
]

seo_items = [
    ('Rank Math SEO', 'Unique title + meta per page, H1-H3 hierarchy, XML sitemap auto-generated'),
    ('4 Schema Types', 'LocalBusiness, Service, FAQ, and Review structured data embedded at launch'),
    ('Priority Keywords', '"Maryland driveway restoration", "Annapolis driveway restoration", "blacktop restoration" + 4 more'),
    ('Google integrations', 'GA4, Search Console, Business Profile all configured and verified on Day 5'),
]

y = Inches(2.0)
for title, detail in mobile_items:
    add_rect(s, Inches(0.55), y + Inches(0.08), Inches(0.12), Inches(0.12), GREEN_ACC)
    add_text(s, title, Inches(0.82), y, Inches(5.5), Inches(0.25), font_size=12, bold=True, color=WHITE)
    add_text(s, detail, Inches(0.82), y + Inches(0.27), Inches(5.6), Inches(0.3), font_size=10, color=TEXT_DIM)
    y += Inches(0.72)

y = Inches(2.0)
for title, detail in seo_items:
    add_rect(s, Inches(7.0), y + Inches(0.08), Inches(0.12), Inches(0.12), GREEN_ACC)
    add_text(s, title, Inches(7.27), y, Inches(5.5), Inches(0.25), font_size=12, bold=True, color=WHITE)
    add_text(s, detail, Inches(7.27), y + Inches(0.27), Inches(5.6), Inches(0.3), font_size=10, color=TEXT_DIM)
    y += Inches(0.72)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — 7-Day Timeline
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, '7-DAY LAUNCH TIMELINE', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, 'Day-by-day delivery plan.', Inches(0.55), Inches(1.0), Inches(9), Inches(0.6),
         font_size=26, bold=True, color=WHITE)

timeline = [
    ('Day 1', 'Kickoff + Hosting', 'Rocket.net provisioned, WP installed, brand assets received, Bricks Builder licensed'),
    ('Day 2', 'Theme + Core Pages', 'Custom theme built (charcoal/green/white). Homepage, Driveway Restoration, Pricing completed'),
    ('Day 3', 'All Pages + Forms', 'Gallery, About, FAQ, Contact, Service Areas, /free-estimate. Fluent Forms wired with photo upload'),
    ('Day 4', 'SEO Config', 'Rank Math: title/meta/H-structure/schema/sitemap/image ALT text/internal links set per page'),
    ('Day 5', 'Analytics + GBP', 'GA4 + GSC via Site Kit. Google Business Profile connected and optimized. Event tracking live'),
    ('Day 6', 'Performance QA', 'ShortPixel run. Rocket CDN verified. Lighthouse 90+ confirmed. Mobile CTA testing on devices'),
    ('Day 7', 'Launch + Handoff', 'DNS cutover, SSL verified, owner walkthrough. You can edit the site without a developer.'),
]

y = Inches(1.85)
for day, title, detail in timeline:
    add_rect(s, Inches(0.55), y, Inches(0.85), Inches(0.55), GREEN)
    add_text(s, day, Inches(0.55), y + Inches(0.08), Inches(0.85), Inches(0.4),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.55), y, Inches(2.8), Inches(0.3), font_size=11, bold=True, color=WHITE)
    add_text(s, detail, Inches(1.55), y + Inches(0.28), Inches(11), Inches(0.3), font_size=10, color=TEXT_DIM)
    y += Inches(0.72)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Pricing
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, 'INVESTMENT', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, 'Fixed price. No surprises.', Inches(0.55), Inches(1.0), Inches(9), Inches(0.7),
         font_size=28, bold=True, color=WHITE)

# Phase 1 card
add_rect(s, Inches(0.55), Inches(1.9), Inches(5.8), Inches(4.8), CARD, GREEN)
add_text(s, 'PHASE 1 -- FIXED PRICE', Inches(0.75), Inches(2.05), Inches(5.4), Inches(0.4),
         font_size=10, bold=True, color=GREEN_ACC)
add_text(s, '$1,100', Inches(0.75), Inches(2.5), Inches(5.4), Inches(0.8),
         font_size=44, bold=True, color=WHITE)
add_text(s, 'flat  |  Rocket.net hosting, first 2 months included', Inches(0.75), Inches(3.3), Inches(5.4), Inches(0.4),
         font_size=11, color=GREEN_ACC)

p1_items = [
    '9-page WordPress site (all Phase 1 pages)',
    'Bricks Builder visual editor (owner-editable)',
    'Rank Math + all 4 schema types at launch',
    'Fluent Forms lead capture + photo upload',
    'ShortPixel, Wordfence, UpdraftPlus configured',
    'GA4, GSC, Google Business Profile setup',
    '30-day post-launch support window',
    'Owner walkthrough + WP admin training',
]
y = Inches(3.8)
for item in p1_items:
    add_rect(s, Inches(0.75), y + Inches(0.07), Inches(0.1), Inches(0.1), GREEN_ACC)
    add_text(s, item, Inches(1.0), y, Inches(5.1), Inches(0.35), font_size=10, color=TEXT_DIM)
    y += Inches(0.38)

# Phase 2 card
add_rect(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(4.8), CARD)
add_text(s, 'PHASE 2 -- ONGOING PARTNERSHIP', Inches(7.0), Inches(2.05), Inches(5.6), Inches(0.4),
         font_size=10, bold=True, color=GREEN_ACC)
add_text(s, '$3,500 - $6,500', Inches(7.0), Inches(2.5), Inches(5.6), Inches(0.8),
         font_size=32, bold=True, color=WHITE)
add_text(s, 'phased rollout  |  no rebuild required', Inches(7.0), Inches(3.3), Inches(5.6), Inches(0.4),
         font_size=11, color=TEXT_DIM)

p2_items = [
    '30 to 50 Maryland city landing pages',
    'Neighborhood hyperlocal pages',
    'Learning Center (blog + SEO content hub)',
    'Interactive estimate calculator',
    'Video library (time-lapses, testimonials)',
    'Expanded gallery (60+ images)',
    'Case studies with outcomes + photos',
    'Advanced local SEO, geo-targeted keywords',
]
y = Inches(3.8)
for item in p2_items:
    add_rect(s, Inches(7.0), y + Inches(0.07), Inches(0.1), Inches(0.1), GREEN)
    add_text(s, item, Inches(7.25), y, Inches(5.3), Inches(0.35), font_size=10, color=TEXT_DIM)
    y += Inches(0.38)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Phase 2 Roadmap
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, 'PHASE 2 ROADMAP', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, 'Multi-state growth, built on Phase 1.', Inches(0.55), Inches(1.0), Inches(9), Inches(0.7),
         font_size=26, bold=True, color=WHITE)
add_text(s, 'Phase 2 layers on top of Phase 1. Same WordPress install, same theme, same SEO foundation. No rebuild.',
         Inches(0.55), Inches(1.8), Inches(9), Inches(0.5), font_size=13, color=TEXT_DIM)

roadmap = [
    ('Q1', '30-50 City Pages', 'Annapolis, Baltimore, Columbia, Frederick, Towson + more. Each page targets exact-match local keywords for that market.'),
    ('Q1', 'Learning Center', 'Blog + resource hub on restoration tips, maintenance, longevity. Builds topical authority and long-tail search traffic.'),
    ('Q2', 'Neighborhood Pages', 'Hyperlocal targeting below city level. Ranks for "driveway restoration [neighborhood]" searches.'),
    ('Q2', 'Estimate Calculator', 'Interactive tool: enter property size and surface type, get an instant range. Major lead qualifier and trust signal.'),
    ('Q3', 'Video Library', 'Before/after time-lapses, testimonial videos, process walkthroughs. SEO-titled and embedded for engagement.'),
    ('Q3', 'Case Studies + Expanded Gallery', 'Detailed project write-ups with photos, outcomes, and client quotes. Closes larger jobs.'),
]

y = Inches(2.5)
for quarter, title, detail in roadmap:
    add_rect(s, Inches(0.55), y, Inches(0.7), Inches(0.55), GREEN)
    add_text(s, quarter, Inches(0.55), y + Inches(0.1), Inches(0.7), Inches(0.35),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.4), y, Inches(3.2), Inches(0.28), font_size=12, bold=True, color=WHITE)
    add_text(s, detail, Inches(1.4), y + Inches(0.28), Inches(11), Inches(0.3), font_size=10, color=TEXT_DIM)
    y += Inches(0.75)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Proof / Background
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)
add_text(s, 'PROOF OF WORK', Inches(0.55), Inches(0.5), Inches(8), Inches(0.4),
         font_size=11, bold=True, color=GREEN_ACC)
add_text(s, 'Relevant experience, not a resume.', Inches(0.55), Inches(1.0), Inches(9), Inches(0.7),
         font_size=26, bold=True, color=WHITE)

# fortherecordmn
add_rect(s, Inches(0.55), Inches(2.0), Inches(12.3), Inches(1.6), CARD, GREEN)
add_text(s, 'fortherecordmn.com  --  SEO + Google Analytics Implementation',
         Inches(0.75), Inches(2.1), Inches(11.8), Inches(0.4), font_size=14, bold=True, color=WHITE)
add_text(s, 'SEO architecture and GA4 implementation for a local business: structured data setup, keyword-targeted meta across all pages, Google Search Console configuration, sitemap submission, and GA4 event tracking for user interactions. The same set of tasks as Phase 1 of this project, applied to a local business context.',
         Inches(0.75), Inches(2.55), Inches(11.8), Inches(0.85), font_size=11, color=TEXT_DIM)

# background
add_rect(s, Inches(0.55), Inches(3.75), Inches(12.3), Inches(1.6), CARD)
add_text(s, 'Full-Stack Development Background',
         Inches(0.75), Inches(3.85), Inches(11.8), Inches(0.4), font_size=14, bold=True, color=WHITE)
add_text(s, 'U.S.-based full-stack developer. U.S. Bank (via Turnberry): React + Python + SQL, internal platform serving 600 users/month, became sole developer and SME within months. Currently at Optum: Angular + .NET/C# + PostgreSQL, 150+ story points delivered, team go-to for AI-assisted development. Portfolio at michaelwegter.com (GitHub: mwegter95).',
         Inches(0.75), Inches(4.3), Inches(11.8), Inches(0.85), font_size=11, color=TEXT_DIM)

add_text(s, 'I build fast, become the expert on what I take on, and communicate plainly about where things stand.',
         Inches(0.55), Inches(5.55), Inches(12), Inches(0.5), font_size=13, bold=True, color=GREEN_ACC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — CTA / Close
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
fill_bg(s, CHARCOAL)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.25), GREEN)
add_rect(s, Inches(0), SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), GREEN)

add_text(s, 'Ready to move forward?', Inches(1.2), Inches(1.6), Inches(11), Inches(0.9),
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, 'The demo is live and clickable right now. Log into the WordPress admin and see exactly how your site will work. Questions, a call, or a go-ahead -- any of those works.',
         Inches(1.2), Inches(2.7), Inches(11), Inches(0.9), font_size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# demo link box
add_rect(s, Inches(2.5), Inches(3.8), Inches(8.5), Inches(0.9), CARD, GREEN)
add_text(s, 'https://api.michaelwegter.com/demos/maryland-driveway-restore/',
         Inches(2.6), Inches(3.9), Inches(8.3), Inches(0.5), font_size=14, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)
add_text(s, 'Login: demo_admin / DemoView2026!', Inches(2.6), Inches(4.4), Inches(8.3), Inches(0.25),
         font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, 'Michael Wegter', Inches(1.2), Inches(5.3), Inches(11), Inches(0.5),
         font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 'michaelwegter.com  |  github.com/mwegter95  |  U.S.-based',
         Inches(1.2), Inches(5.85), Inches(11), Inches(0.4), font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), 'maryland-driveway-restore-proposal.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Slides: {len(prs.slides)}')
