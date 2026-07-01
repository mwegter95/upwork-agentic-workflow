"""Build the Pathwaize Intelligence Center proposal deck (dark execution-console aesthetic)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette (matches demo design system) ──────────────────────────────────────
BG_BASE      = RGBColor(0x0B, 0x0F, 0x1A)
BG_SURFACE   = RGBColor(0x14, 0x19, 0x29)
BG_ELEVATED  = RGBColor(0x1A, 0x22, 0x36)
BORDER       = RGBColor(0x1E, 0x27, 0x40)
ACCENT_BLUE  = RGBColor(0x0E, 0xA5, 0xE9)
ACCENT_CYAN  = RGBColor(0x06, 0xB6, 0xD4)
STATUS_OK    = RGBColor(0x22, 0xC5, 0x5E)
STATUS_WARN  = RGBColor(0xF5, 0x9E, 0x0B)
STATUS_RUN   = RGBColor(0x8B, 0x5C, 0xF6)
TEXT_PRIMARY = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_SECOND  = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED   = RGBColor(0x47, 0x55, 0x69)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # blank layout


def new_slide():
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_BASE
    return slide


def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_width=Pt(0.75)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_name="Inter", font_size=Pt(14),
             bold=False, color=TEXT_PRIMARY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_label(slide, text, l, t, w=Inches(8)):
    """Small monospace eyebrow label."""
    add_text(slide, text.upper(), l, t, w, Inches(0.3),
             font_name="DM Mono", font_size=Pt(9),
             color=ACCENT_BLUE, align=PP_ALIGN.LEFT)


def slide_footer(slide, text="Michael Wegter  |  michaelwegter.com  |  Pathwaize Intelligence Center Proposal"):
    add_text(slide, text,
             Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
             font_name="DM Mono", font_size=Pt(7.5),
             color=TEXT_MUTED, align=PP_ALIGN.LEFT)


# ─── SLIDE 1: Cover ───────────────────────────────────────────────────────────
s1 = new_slide()

# Accent rectangle left bar
add_rect(s1, Inches(0), Inches(0), Inches(0.18), H, ACCENT_BLUE)

# Subtle surface panel right
add_rect(s1, Inches(7.8), Inches(0.8), Inches(5.33), Inches(5.9), BG_SURFACE, BORDER)

add_label(s1, "Proposal  /  Full-Stack SaaS Developer", Inches(0.5), Inches(1.4))

add_text(s1, "Pathwaize Intelligence Center", Inches(0.5), Inches(1.85), Inches(7.2), Inches(1.6),
         font_name="Inter", font_size=Pt(42), bold=True, color=TEXT_PRIMARY)

add_text(s1, "MVP", Inches(0.5), Inches(3.3), Inches(3), Inches(0.7),
         font_name="Inter", font_size=Pt(42), bold=True, color=ACCENT_BLUE)

add_text(s1, "A working execution console for real estate investors.", Inches(0.5), Inches(4.1), Inches(7), Inches(0.45),
         font_name="Inter", font_size=Pt(16), color=TEXT_SECOND)

add_text(s1, "michaelwegter.com/demos/real-estate-intelligence-center/",
         Inches(0.5), Inches(4.7), Inches(9), Inches(0.35),
         font_name="DM Mono", font_size=Pt(10), color=ACCENT_CYAN)

add_text(s1, "Michael Wegter  |  Full-Stack Developer  |  July 2026",
         Inches(0.5), Inches(5.4), Inches(7), Inches(0.3),
         font_name="DM Mono", font_size=Pt(9), color=TEXT_MUTED)

# Right panel content
add_text(s1, "DEMO LIVE NOW", Inches(8.1), Inches(1.3), Inches(4.8), Inches(0.3),
         font_name="DM Mono", font_size=Pt(8), color=ACCENT_BLUE)
add_text(s1, "Run an engine.", Inches(8.1), Inches(1.7), Inches(4.8), Inches(0.35),
         font_name="Inter", font_size=Pt(13), bold=True, color=TEXT_PRIMARY)
add_text(s1, "Review the output.", Inches(8.1), Inches(2.1), Inches(4.8), Inches(0.35),
         font_name="Inter", font_size=Pt(13), bold=True, color=TEXT_PRIMARY)
add_text(s1, "Approve it.", Inches(8.1), Inches(2.5), Inches(4.8), Inches(0.35),
         font_name="Inter", font_size=Pt(13), bold=True, color=TEXT_PRIMARY)
add_text(s1, "Push it.", Inches(8.1), Inches(2.9), Inches(4.8), Inches(0.35),
         font_name="Inter", font_size=Pt(13), bold=True, color=STATUS_OK)
add_text(s1, "That is the workflow. It is running at the link above.",
         Inches(8.1), Inches(3.5), Inches(4.8), Inches(0.6),
         font_name="Inter", font_size=Pt(10.5), color=TEXT_SECOND)

slide_footer(s1)


# ─── SLIDE 2: The Core Workflow ───────────────────────────────────────────────
s2 = new_slide()
add_rect(s2, Inches(0), Inches(0), Inches(0.18), H, ACCENT_BLUE)
add_label(s2, "Core Workflow  /  Run > Generate > Review > Revise > Approve > Save/Push", Inches(0.5), Inches(0.35))
add_text(s2, "The workflow, exactly as you scoped it", Inches(0.5), Inches(0.65), Inches(12), Inches(0.55),
         font_name="Inter", font_size=Pt(28), bold=True, color=TEXT_PRIMARY)
add_text(s2, "Approval is a hard gate. Outputs cannot push to any destination until the user explicitly clicks Approve.",
         Inches(0.5), Inches(1.2), Inches(11.5), Inches(0.35),
         font_name="Inter", font_size=Pt(12), color=TEXT_SECOND)

steps = [
    ("01", "Run", "User fills the brief and fires an engine"),
    ("02", "Generate", "LLM router routes request and returns structured output"),
    ("03", "Review", "Output lands in Pending Review queue (status: needs-review)"),
    ("04", "Revise", "User requests edits; engine re-runs with revision notes"),
    ("05", "Approve", "User clicks Approve; status flips to approved"),
    ("06", "Save / Push", "Push button unlocks; output saved to destination"),
]

box_w = Inches(2.0)
box_h = Inches(3.2)
gap   = Inches(0.08)
start_l = Inches(0.4)
top   = Inches(1.75)

colors_step = [BG_ELEVATED, BG_SURFACE, BG_ELEVATED, BG_SURFACE, BG_ELEVATED, BG_SURFACE]
accent_cols = [ACCENT_BLUE, TEXT_SECOND, STATUS_WARN, TEXT_SECOND, STATUS_OK, ACCENT_CYAN]

for i, (num, label, desc) in enumerate(steps):
    l = start_l + i * (box_w + gap)
    add_rect(s2, l, top, box_w, box_h, colors_step[i], BORDER)
    add_text(s2, num, l + Inches(0.15), top + Inches(0.15), box_w - Inches(0.3), Inches(0.3),
             font_name="DM Mono", font_size=Pt(9), color=ACCENT_BLUE)
    add_text(s2, label, l + Inches(0.15), top + Inches(0.5), box_w - Inches(0.3), Inches(0.4),
             font_name="Inter", font_size=Pt(14), bold=True, color=accent_cols[i])
    add_text(s2, desc, l + Inches(0.15), top + Inches(1.0), box_w - Inches(0.3), Inches(2.0),
             font_name="Inter", font_size=Pt(10.5), color=TEXT_SECOND)
    # Arrow between steps
    if i < len(steps) - 1:
        add_text(s2, ">>", l + box_w - Inches(0.02), top + Inches(1.5), Inches(0.25), Inches(0.35),
                 font_name="DM Mono", font_size=Pt(10), color=TEXT_MUTED)

slide_footer(s2)


# ─── SLIDE 3: Three Engines ───────────────────────────────────────────────────
s3 = new_slide()
add_rect(s3, Inches(0), Inches(0), Inches(0.18), H, ACCENT_BLUE)
add_label(s3, "V1 Engines  /  Three AI Engines Wired End-to-End", Inches(0.5), Inches(0.35))
add_text(s3, "Three engines, all running in the demo today", Inches(0.5), Inches(0.65), Inches(12), Inches(0.55),
         font_name="Inter", font_size=Pt(28), bold=True, color=TEXT_PRIMARY)

engines = [
    ("Priority 1", "Authority Engine",
     "Generates founder-authority content: LinkedIn posts, X threads, thought-leadership pieces.",
     "Input: topic, tone, platform\nOutput: formatted post + revision option\nPush: LinkedIn or X (approval-gated)"),
    ("Priority 2", "Newsletter Engine",
     "Generates newsletter body copy, three subject-line options, and preview text. Subject lines are selectable cards.",
     "Input: topic, audience, brand voice\nOutput: body + 3 subjects + preview text\nPush: save as draft or schedule (approval-gated)"),
    ("Priority 3", "Knowledge Engine Lite",
     "Generates business KB entries. Destination picker: Google Drive or GitHub. Save receipt with file path and asset link.",
     "Input: topic, business context\nOutput: KB entry (formatted markdown)\nPush: Drive or GitHub (approval-gated)"),
]

card_w = Inches(4.05)
card_h = Inches(5.0)
top_e  = Inches(1.4)
gap_e  = Inches(0.25)

for i, (badge, title, desc, io) in enumerate(engines):
    l = Inches(0.4) + i * (card_w + gap_e)
    add_rect(s3, l, top_e, card_w, card_h, BG_SURFACE, BORDER)
    # Badge
    add_rect(s3, l + Inches(0.18), top_e + Inches(0.18), Inches(1.0), Inches(0.28),
             RGBColor(0x0E, 0x3A, 0x58), BORDER, Pt(0.5))
    add_text(s3, badge, l + Inches(0.2), top_e + Inches(0.17), Inches(1.2), Inches(0.3),
             font_name="DM Mono", font_size=Pt(7.5), color=ACCENT_BLUE)
    add_text(s3, title, l + Inches(0.18), top_e + Inches(0.58), card_w - Inches(0.36), Inches(0.45),
             font_name="Inter", font_size=Pt(15), bold=True, color=TEXT_PRIMARY)
    add_text(s3, desc, l + Inches(0.18), top_e + Inches(1.1), card_w - Inches(0.36), Inches(1.2),
             font_name="Inter", font_size=Pt(10.5), color=TEXT_SECOND)
    # Divider line via rect
    add_rect(s3, l + Inches(0.18), top_e + Inches(2.45), card_w - Inches(0.36), Inches(0.01), BORDER)
    add_text(s3, io, l + Inches(0.18), top_e + Inches(2.6), card_w - Inches(0.36), Inches(2.2),
             font_name="DM Mono", font_size=Pt(9), color=TEXT_SECOND)

slide_footer(s3)


# ─── SLIDE 4: LLM Routing Architecture ───────────────────────────────────────
s4 = new_slide()
add_rect(s4, Inches(0), Inches(0), Inches(0.18), H, ACCENT_BLUE)
add_label(s4, "Backend Architecture  /  Provider-Agnostic LLM Routing", Inches(0.5), Inches(0.35))
add_text(s4, "Frontend never calls an LLM provider directly", Inches(0.5), Inches(0.65), Inches(12), Inches(0.55),
         font_name="Inter", font_size=Pt(28), bold=True, color=TEXT_PRIMARY)
add_text(s4, "One provider-agnostic contract. Adding Claude when OpenAI is live = one config entry, not a rebuild.",
         Inches(0.5), Inches(1.2), Inches(11.5), Inches(0.35),
         font_name="Inter", font_size=Pt(12), color=TEXT_SECOND)

nodes = [
    ("Next.js / React", "GHL iframe embed\nno LLM calls"),
    ("POST /route", "{ engine, payload,\nworkspace_id }"),
    ("LLM Router", "selects model from\nconfig registry"),
    ("Provider API", "OpenAI / Claude /\nlocal model"),
    ("Supabase Log", "run_id, model, tokens\ncost, status"),
]
node_w = Inches(2.35)
node_h = Inches(1.6)
gap_n  = Inches(0.18)
top_n  = Inches(2.0)
for i, (title, sub) in enumerate(nodes):
    l = Inches(0.35) + i * (node_w + gap_n)
    col = ACCENT_BLUE if i == 1 else BG_ELEVATED
    tcol = WHITE if i == 1 else TEXT_PRIMARY
    add_rect(s4, l, top_n, node_w, node_h, col, BORDER)
    add_text(s4, title, l + Inches(0.12), top_n + Inches(0.15), node_w - Inches(0.24), Inches(0.4),
             font_name="Inter", font_size=Pt(11), bold=True, color=tcol)
    add_text(s4, sub, l + Inches(0.12), top_n + Inches(0.65), node_w - Inches(0.24), Inches(0.8),
             font_name="DM Mono", font_size=Pt(8.5), color=TEXT_SECOND if i != 1 else RGBColor(0xBA, 0xE4, 0xF7))
    if i < len(nodes) - 1:
        add_text(s4, ">>", l + node_w + Inches(0.04), top_n + Inches(0.65), Inches(0.18), Inches(0.35),
                 font_name="DM Mono", font_size=Pt(10), color=TEXT_MUTED)

# Stack table below
stack_items = [
    ("Frontend", "Next.js / React in GHL iframe. Calls only /route, never any provider."),
    ("Router (AWS Lambda)", "Provider-agnostic. Config-driven model selection. Logs every run."),
    ("Data (Supabase)", "Runs, outputs, logs, workspace RLS. Your existing infrastructure."),
    ("Integrations", "Drive + GitHub OAuth server-side. LinkedIn + X via approved APIs."),
]
row_h = Inches(0.6)
top_st = Inches(4.0)
for i, (k, v) in enumerate(stack_items):
    bg = BG_SURFACE if i % 2 == 0 else BG_ELEVATED
    add_rect(s4, Inches(0.4), top_st + i * row_h, Inches(12.5), row_h, bg, BORDER, Pt(0.5))
    add_text(s4, k, Inches(0.55), top_st + i * row_h + Inches(0.14), Inches(2.5), Inches(0.35),
             font_name="DM Mono", font_size=Pt(9), bold=True, color=ACCENT_CYAN)
    add_text(s4, v, Inches(3.1), top_st + i * row_h + Inches(0.14), Inches(9.6), Inches(0.35),
             font_name="Inter", font_size=Pt(10), color=TEXT_SECOND)

slide_footer(s4)


# ─── SLIDE 5: Approval Gating ────────────────────────────────────────────────
s5 = new_slide()
add_rect(s5, Inches(0), Inches(0), Inches(0.18), H, STATUS_OK)
add_label(s5, "Approval Gating  /  Two Independent Layers", Inches(0.5), Inches(0.35))
add_text(s5, "Outputs cannot publish without explicit user approval", Inches(0.5), Inches(0.65), Inches(12), Inches(0.55),
         font_name="Inter", font_size=Pt(28), bold=True, color=TEXT_PRIMARY)
add_text(s5, "Frontend status machine + backend server-side guard enforce the same rule independently.",
         Inches(0.5), Inches(1.2), Inches(11.5), Inches(0.35),
         font_name="Inter", font_size=Pt(12), color=TEXT_SECOND)

gates = [
    ("Frontend Status Machine",
     "Push/Save button rendered disabled + cursor: not-allowed unless status = 'approved'. No code path advances to 'pushed' without explicit Approve click. Verify in the demo: try clicking Push before Approve.",
     ACCENT_BLUE),
    ("Backend Server-Side Guard",
     "Push endpoints check output status in Supabase before executing. Returns 403 if not approved, even if the frontend gate is somehow bypassed. Frontend and backend enforce the same rule independently.",
     ACCENT_CYAN),
    ("Output Status Lifecycle",
     "generating > needs-review > (revised >) approved > pushed. Each transition is user-triggered or explicitly confirmed. No automatic advancement.",
     STATUS_OK),
]

gate_w = Inches(3.9)
gate_h = Inches(3.8)
gap_g  = Inches(0.3)
top_g  = Inches(1.85)
for i, (title, body, color) in enumerate(gates):
    l = Inches(0.4) + i * (gate_w + gap_g)
    add_rect(s5, l, top_g, gate_w, gate_h, BG_SURFACE, BORDER)
    # Color top bar
    add_rect(s5, l, top_g, gate_w, Inches(0.06), color)
    add_text(s5, title, l + Inches(0.18), top_g + Inches(0.2), gate_w - Inches(0.36), Inches(0.45),
             font_name="Inter", font_size=Pt(13), bold=True, color=color)
    add_text(s5, body, l + Inches(0.18), top_g + Inches(0.85), gate_w - Inches(0.36), Inches(2.7),
             font_name="Inter", font_size=Pt(10.5), color=TEXT_SECOND)

slide_footer(s5)


# ─── SLIDE 6: Milestone Plan ──────────────────────────────────────────────────
s6 = new_slide()
add_rect(s6, Inches(0), Inches(0), Inches(0.18), H, ACCENT_BLUE)
add_label(s6, "Build Plan  /  Milestone-Based, Fixed Price", Inches(0.5), Inches(0.35))
add_text(s6, "6 to 7 weeks. $9,000 fixed price.", Inches(0.5), Inches(0.65), Inches(12), Inches(0.55),
         font_name="Inter", font_size=Pt(28), bold=True, color=TEXT_PRIMARY)
add_text(s6, "Each milestone is deployable. No billing until you have seen the working increment.",
         Inches(0.5), Inches(1.2), Inches(11.5), Inches(0.35),
         font_name="Inter", font_size=Pt(12), color=TEXT_SECOND)

milestones = [
    ("M1", "UI Shell + Authority Engine end-to-end",
     "3-col console layout, GHL iframe, run/review/approve/push on Authority Engine, run records, LLM routing layer",
     "1.5 wks", "$2,500"),
    ("M2", "Newsletter Engine end-to-end",
     "Full body + 3 subject lines + preview text, selectable subjects, revision panel, approval-gated save/schedule",
     "1 wk", "$1,500"),
    ("M3", "Knowledge Engine Lite + Drive/GitHub save",
     "KB generation, destination picker, server-side OAuth for Drive + GitHub, save receipt with file path",
     "1.5 wks", "$2,000"),
    ("M4", "Outputs queue + approve/revise/re-run flow",
     "Pending Review queue, click-to-review, approve/request-edits/re-run actions, run history panel",
     "1 wk", "$1,500"),
    ("M5", "Intel Brief stub + polish + GHL QA + deploy",
     "6-section Intel Brief, workspace context, GHL postMessage integration, final QA, production deploy",
     "1 wk", "$1,500"),
]

row_h  = Inches(0.86)
top_m  = Inches(1.75)
for i, (badge, title, desc, dur, price) in enumerate(milestones):
    bg = BG_SURFACE if i % 2 == 0 else BG_ELEVATED
    add_rect(s6, Inches(0.4), top_m + i * row_h, Inches(12.5), row_h, bg, BORDER, Pt(0.5))
    # Badge
    add_text(s6, badge, Inches(0.55), top_m + i * row_h + Inches(0.2), Inches(0.55), Inches(0.35),
             font_name="DM Mono", font_size=Pt(9), bold=True, color=ACCENT_BLUE)
    add_text(s6, title, Inches(1.2), top_m + i * row_h + Inches(0.1), Inches(7.5), Inches(0.35),
             font_name="Inter", font_size=Pt(11), bold=True, color=TEXT_PRIMARY)
    add_text(s6, desc, Inches(1.2), top_m + i * row_h + Inches(0.46), Inches(7.5), Inches(0.3),
             font_name="Inter", font_size=Pt(9.5), color=TEXT_SECOND)
    add_text(s6, dur, Inches(8.85), top_m + i * row_h + Inches(0.24), Inches(1.3), Inches(0.35),
             font_name="DM Mono", font_size=Pt(10), color=TEXT_SECOND, align=PP_ALIGN.CENTER)
    add_text(s6, price, Inches(10.3), top_m + i * row_h + Inches(0.24), Inches(2.2), Inches(0.35),
             font_name="DM Mono", font_size=Pt(11), bold=True, color=STATUS_OK, align=PP_ALIGN.RIGHT)

# Total row
total_y = top_m + len(milestones) * row_h + Inches(0.06)
add_rect(s6, Inches(0.4), total_y, Inches(12.5), Inches(0.42), BG_BASE, BORDER, Pt(0.5))
add_text(s6, "TOTAL", Inches(1.2), total_y + Inches(0.06), Inches(4), Inches(0.3),
         font_name="DM Mono", font_size=Pt(9), bold=True, color=TEXT_MUTED)
add_text(s6, "6-7 weeks", Inches(8.85), total_y + Inches(0.06), Inches(1.3), Inches(0.3),
         font_name="DM Mono", font_size=Pt(9.5), bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
add_text(s6, "$9,000", Inches(10.3), total_y + Inches(0.04), Inches(2.2), Inches(0.35),
         font_name="DM Mono", font_size=Pt(13), bold=True, color=STATUS_OK, align=PP_ALIGN.RIGHT)

slide_footer(s6)


# ─── SLIDE 7: Why Michael ─────────────────────────────────────────────────────
s7 = new_slide()
add_rect(s7, Inches(0), Inches(0), Inches(0.18), H, ACCENT_BLUE)
add_label(s7, "Why Michael Wegter  /  Proof Points", Inches(0.5), Inches(0.35))
add_text(s7, "Full-stack ownership. LLMs in production. No overengineering.",
         Inches(0.5), Inches(0.65), Inches(12), Inches(0.55),
         font_name="Inter", font_size=Pt(28), bold=True, color=TEXT_PRIMARY)

proofs = [
    ("Sole developer, 600 users/month",
     "U.S. Bank (via Turnberry): sole dev and SME on TDAAS, ~60k LOC, React + Python + SQL + Java. Became go-to within months. Led the 6-month Azure Cloud migration as the project's main rep."),
    ("LLM workflows in production",
     "At Optum (current, April 2026): team's go-to for AI-assisted development using Azure OpenAI. The demo implements prompt chaining, JSON-schema output enforcement, sliding context windows, and 1-retry reflection."),
    ("The demo is the application",
     "I built a working Intelligence Center before applying. All three engines, approval gating, run records, cost logging, and the LLM routing layer are functional at the demo link. Working software, not slides about potential."),
    ("No overengineering on principle",
     "The demo is 12 files. The production design is one routing service + 4 Supabase tables. V1 should validate the core workflow quickly. I read the posting: skip advanced analytics, skip future engines, ship the loop."),
]

card_w2 = Inches(5.9)
card_h2 = Inches(2.5)
gap2    = Inches(0.25)
top2    = Inches(1.55)

for i, (title, body) in enumerate(proofs):
    row = i // 2
    col = i % 2
    l = Inches(0.4) + col * (card_w2 + gap2)
    t = top2 + row * (card_h2 + gap2)
    add_rect(s7, l, t, card_w2, card_h2, BG_SURFACE, BORDER)
    add_text(s7, title, l + Inches(0.18), t + Inches(0.18), card_w2 - Inches(0.36), Inches(0.4),
             font_name="Inter", font_size=Pt(12), bold=True, color=ACCENT_BLUE)
    add_text(s7, body, l + Inches(0.18), t + Inches(0.7), card_w2 - Inches(0.36), Inches(1.6),
             font_name="Inter", font_size=Pt(10.5), color=TEXT_SECOND)

slide_footer(s7)


# ─── SLIDE 8: Call to Action ──────────────────────────────────────────────────
s8 = new_slide()
add_rect(s8, Inches(0), Inches(0), Inches(0.18), H, ACCENT_BLUE)

# Center panel
add_rect(s8, Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.1), BG_SURFACE, BORDER)

add_text(s8, "The demo is live.", Inches(2.0), Inches(1.65), Inches(9.3), Inches(0.7),
         font_name="Inter", font_size=Pt(36), bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)

add_text(s8, "Click it. Run an engine. See the approval gate in action.",
         Inches(2.0), Inches(2.45), Inches(9.3), Inches(0.4),
         font_name="Inter", font_size=Pt(14), color=TEXT_SECOND, align=PP_ALIGN.CENTER)

add_text(s8, "michaelwegter.com/demos/real-estate-intelligence-center/",
         Inches(2.0), Inches(3.05), Inches(9.3), Inches(0.4),
         font_name="DM Mono", font_size=Pt(13), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_text(s8, "Michael Wegter  |  michaelwegter.com  |  github.com/mwegter95",
         Inches(2.0), Inches(3.7), Inches(9.3), Inches(0.35),
         font_name="DM Mono", font_size=Pt(9.5), color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_text(s8, "Ready to start on the real build when you are.",
         Inches(2.0), Inches(4.4), Inches(9.3), Inches(0.4),
         font_name="Inter", font_size=Pt(13), color=TEXT_SECOND, align=PP_ALIGN.CENTER)

slide_footer(s8)


# ─── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/michaelwegter/Desktop/Projects/michaelwegter.com/public/demos/real-estate-intelligence-center/proposal/deck.pptx"
prs.save(out)
print(f"Saved: {out}")
